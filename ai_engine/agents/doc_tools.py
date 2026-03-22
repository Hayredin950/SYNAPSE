"""
ai_engine.agents.doc_tools
~~~~~~~~~~~~~~~~~~~~~~~~~~
💎 DIAMOND-LEVEL Document generation tools for SYNAPSE AI.

Generators:
  1. PDF       — ReportLab: gradient cover, embedded Matplotlib charts (bar,
                 pie, line), callout boxes, pull-quotes, infographic stat cards,
                 styled TOC, branded headers/footers, decorative geometry.
  2. PPTX      — python-pptx: gradient title/section slides, chart slides
                 (bar, pie via Matplotlib PNG), icon-bullet content slides,
                 two-column agenda, stats dashboard slide, closing slide.
  3. DOCX      — python-docx: shaded cover block, embedded charts, styled
                 tables, callout boxes with left accent, TOC field, page nums.
  4. Markdown  — Extended MD: YAML front-matter, badges, ASCII-art stat
                 cards, admonition blocks, emoji TOC, summary tables.
  5. HTML      — Standalone HTML5 with embedded CSS (glassmorphism, gradients,
                 dark sidebar nav), Chart.js charts (bar, pie, line),
                 animated stat cards, responsive layout — zero dependencies.

All generators: save to MEDIA_ROOT/documents/<user_id>/<uuid>.<ext>
                return a plain string "Path: ...\nSize: ..." for the agent.
"""

from __future__ import annotations

import io
import logging
import math
import os
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

from langchain_core.tools import StructuredTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# 💎 Brand palette
# ─────────────────────────────────────────────────────────────────────────────
class _Brand:
    INDIGO = "#4F46E5"
    INDIGO_DARK = "#3730A3"
    INDIGO_MID = "#6366F1"
    INDIGO_LIGHT = "#818CF8"
    INDIGO_PALE = "#C7D2FE"
    INDIGO_BG = "#EEF2FF"
    VIOLET = "#7C3AED"
    VIOLET_LIGHT = "#A78BFA"
    CYAN = "#06B6D4"
    EMERALD = "#10B981"
    AMBER = "#F59E0B"
    ROSE = "#F43F5E"
    DARK = "#1E1B4B"
    GRAY_900 = "#111827"
    GRAY_700 = "#374151"
    GRAY_500 = "#6B7280"
    GRAY_300 = "#D1D5DB"
    GRAY_100 = "#F3F4F6"
    WHITE = "#FFFFFF"

    # Accent palette for charts (cycles)
    CHART = [
        "#4F46E5",
        "#7C3AED",
        "#06B6D4",
        "#10B981",
        "#F59E0B",
        "#F43F5E",
        "#818CF8",
        "#A78BFA",
        "#34D399",
        "#FCD34D",
    ]

    @staticmethod
    def rl(h: str):
        from reportlab.lib.colors import HexColor

        return HexColor(h)

    @staticmethod
    def rgb(h: str):
        from pptx.dml.color import RGBColor

        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def mpl(h: str) -> str:
        """Return hex suitable for matplotlib."""
        return h


# ─────────────────────────────────────────────────────────────────────────────
# Storage helpers
# ─────────────────────────────────────────────────────────────────────────────
def _doc_dir(user_id: str = "anonymous") -> Path:
    """Return the per-user document directory, creating it if needed.

    SEC-05: Sanitise user_id to prevent path traversal attacks.
    A malicious user_id like '../../etc' could escape the media root.
    We strip all path separators and non-alphanumeric characters.
    """
    import re as _re

    # Allow only alphanumeric, hyphens, and underscores — reject everything else.
    # This covers path separators (/ \\), dots (..), null bytes, and Unicode tricks.
    safe_user_id = _re.sub(r"[^a-zA-Z0-9_\-]", "_", str(user_id))
    if not safe_user_id or safe_user_id in ("", "_", "__"):
        safe_user_id = "anonymous"

    root = Path(
        os.environ.get("DJANGO_MEDIA_ROOT") or os.environ.get("MEDIA_ROOT") or "media"
    ).resolve()
    d = root / "documents" / safe_user_id

    # Final guard: ensure the resolved path is still inside root
    # (defence-in-depth against symlink attacks)
    try:
        d.resolve().relative_to(root)
    except ValueError:
        raise PermissionError(f"Resolved document path escapes media root: {d}")

    d.mkdir(parents=True, exist_ok=True)
    return d


def _rel(abs_path: Path) -> str:
    root = Path(
        os.environ.get("DJANGO_MEDIA_ROOT") or os.environ.get("MEDIA_ROOT") or "media"
    )
    try:
        return str(abs_path.relative_to(root))
    except:
        return str(abs_path)


# ─────────────────────────────────────────────────────────────────────────────
# 📊 Chart factory  (returns PNG bytes via Matplotlib)
# ─────────────────────────────────────────────────────────────────────────────
def _chart_bytes(
    chart_type: str,  # "bar" | "pie" | "line" | "hbar"
    labels: List[str],
    values: List[float],
    title: str = "",
    width_in: float = 5.5,
    height_in: float = 3.2,
    dark_bg: bool = False,
) -> bytes:
    """Render a Matplotlib chart and return PNG bytes."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.patches as mpatches
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch

    colors = [_Brand.CHART[i % len(_Brand.CHART)] for i in range(len(labels))]

    bg = "#1E1B4B" if dark_bg else "#F8F9FF"
    fg = "#EEF2FF" if dark_bg else "#1E1B4B"
    grid = "#3730A3" if dark_bg else "#E0E4FF"

    fig, ax = plt.subplots(figsize=(width_in, height_in))
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)

    if chart_type == "pie":
        wedge_props = {"linewidth": 2, "edgecolor": bg}
        wedges, texts, autotexts = ax.pie(
            values,
            labels=None,
            colors=colors,
            autopct="%1.1f%%",
            startangle=140,
            wedgeprops=wedge_props,
            pctdistance=0.78,
        )
        for at in autotexts:
            at.set_color(bg)
            at.set_fontsize(8)
            at.set_fontweight("bold")
        ax.legend(
            wedges,
            labels,
            loc="center left",
            bbox_to_anchor=(0.85, 0.5),
            fontsize=7.5,
            frameon=False,
            labelcolor=fg,
        )
    elif chart_type in ("bar", "hbar"):
        if chart_type == "bar":
            bars = ax.bar(
                labels, values, color=colors, width=0.6, edgecolor=bg, linewidth=1.5
            )
            ax.set_xlabel("", color=fg)
            ax.tick_params(axis="x", colors=fg, labelsize=8, rotation=15)
            for bar, val in zip(bars, values):
                ax.text(
                    bar.get_x() + bar.get_width() / 2,
                    bar.get_height() + max(values) * 0.01,
                    f"{val:,.0f}",
                    ha="center",
                    va="bottom",
                    color=fg,
                    fontsize=8,
                    fontweight="bold",
                )
        else:
            bars = ax.barh(labels, values, color=colors, edgecolor=bg, linewidth=1.5)
            ax.tick_params(axis="y", colors=fg, labelsize=8)
            for bar, val in zip(bars, values):
                ax.text(
                    bar.get_width() + max(values) * 0.01,
                    bar.get_y() + bar.get_height() / 2,
                    f"{val:,.0f}",
                    va="center",
                    color=fg,
                    fontsize=8,
                    fontweight="bold",
                )
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(grid)
        ax.spines["bottom"].set_color(grid)
        ax.tick_params(axis="y", colors=fg, labelsize=8)
        ax.yaxis.grid(True, color=grid, linewidth=0.5, linestyle="--")
        ax.set_axisbelow(True)
    elif chart_type == "line":
        x = list(range(len(labels)))
        ax.plot(
            x,
            values,
            color=colors[0],
            linewidth=2.5,
            marker="o",
            markersize=6,
            markerfacecolor=bg,
            markeredgecolor=colors[0],
            markeredgewidth=2,
        )
        ax.fill_between(x, values, alpha=0.15, color=colors[0])
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=8, color=fg, rotation=15)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(grid)
        ax.spines["bottom"].set_color(grid)
        ax.tick_params(axis="y", colors=fg, labelsize=8)
        ax.yaxis.grid(True, color=grid, linewidth=0.5, linestyle="--")
        ax.set_axisbelow(True)

    if title:
        ax.set_title(title, color=fg, fontsize=10, fontweight="bold", pad=8)

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor=bg)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_b64(
    chart_type, labels, values, title="", width_in=5.5, height_in=3.2, dark_bg=False
) -> str:
    """Return base64-encoded PNG for embedding in HTML."""
    import base64

    data = _chart_bytes(chart_type, labels, values, title, width_in, height_in, dark_bg)
    return base64.b64encode(data).decode()


def _generate_cover_image(
    title: str,
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    doc_type: str = "pdf",
    width: int = 1200,
    height: int = 630,
):
    """Generate a premium format-specific cover image using Pillow. Returns PIL Image."""
    from PIL import Image, ImageDraw, ImageFont

    # ── Format-specific colour themes ────────────────────────────────
    THEMES = {
        "pdf": {
            "c1": (55, 48, 163),
            "c2": (124, 58, 237),
            "accent": (129, 140, 248),
            "icon": "PDF",
            "badge": (79, 70, 229),
            "label": "PDF REPORT",
        },
        "ppt": {
            "c1": (180, 83, 9),
            "c2": (217, 119, 6),
            "accent": (251, 191, 36),
            "icon": "PPT",
            "badge": (245, 158, 11),
            "label": "POWERPOINT",
        },
        "word": {
            "c1": (29, 78, 216),
            "c2": (37, 99, 235),
            "accent": (96, 165, 250),
            "icon": "DOC",
            "badge": (59, 130, 246),
            "label": "WORD DOCUMENT",
        },
        "markdown": {
            "c1": (5, 150, 105),
            "c2": (16, 185, 129),
            "accent": (52, 211, 153),
            "icon": "MD",
            "badge": (16, 185, 129),
            "label": "MARKDOWN",
        },
        "html": {
            "c1": (8, 145, 178),
            "c2": (6, 182, 212),
            "accent": (103, 232, 249),
            "icon": "HTML",
            "badge": (6, 182, 212),
            "label": "HTML PAGE",
        },
    }
    theme = THEMES.get(doc_type, THEMES["pdf"])
    c1, c2 = theme["c1"], theme["c2"]
    accent = theme["accent"]
    badge = theme["badge"]
    label = theme["label"]
    icon_txt = theme["icon"]

    img = Image.new("RGBA", (width, height), (0, 0, 0, 255))
    draw = ImageDraw.Draw(img, "RGBA")

    # ── Diagonal gradient background ────────────────────────────────
    for y in range(height):
        t = y / height
        r = int(c1[0] + t * (c2[0] - c1[0]))
        g = int(c1[1] + t * (c2[1] - c1[1]))
        b = int(c1[2] + t * (c2[2] - c1[2]))
        draw.line([(0, y), (width, y)], fill=(r, g, b, 255))

    # ── Decorative geometry: large translucent circles ───────────────
    for cx, cy, cr, alpha in [
        (width - 60, -100, 380, 22),
        (width + 80, 120, 500, 12),
        (-80, height - 80, 300, 18),
        (width // 2 + 200, height + 60, 400, 10),
    ]:
        draw.ellipse([cx - cr, cy - cr, cx + cr, cy + cr], fill=(*accent, alpha))

    # ── Format badge (top-right) ─────────────────────────────────────
    badge_w, badge_h = 160, 56
    bx, by = width - badge_w - 40, 36
    # Rounded rect approximation via ellipse + rectangle
    draw.rectangle([bx + 14, by, bx + badge_w - 14, by + badge_h], fill=(*badge, 230))
    draw.ellipse([bx, by, bx + 28, by + badge_h], fill=(*badge, 230))
    draw.ellipse(
        [bx + badge_w - 28, by, bx + badge_w, by + badge_h], fill=(*badge, 230)
    )

    # ── Diagonal decorative stripe ───────────────────────────────────
    stripe_pts = [
        (0, height * 0.72),
        (width * 0.45, height * 0.72),
        (width * 0.55, height),
        (0, height),
    ]
    draw.polygon(stripe_pts, fill=(255, 255, 255, 12))

    # ── Left accent bar (gradient top-to-bottom) ─────────────────────
    for y in range(height):
        alpha_bar = int(180 + 75 * (y / height))
        draw.line([(0, y), (10, y)], fill=(*accent, alpha_bar))

    # ── Bottom frosted strip ─────────────────────────────────────────
    draw.rectangle([0, height - 80, width, height], fill=(255, 255, 255, 245))

    # ── Large icon text (top-left decorative watermark) ─────────────
    # Semi-transparent large icon in background
    draw.rectangle(
        [44, 36, 44 + badge_w - 60, 36 + badge_h - 4], fill=(255, 255, 255, 25)
    )

    # ── Load fonts ───────────────────────────────────────────────────
    import os as _os

    FONT_PATHS = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "C:/Windows/Fonts/arialbd.ttf",
    ]
    FONT_REG = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]

    def load_font(paths, size):
        for fp in paths:
            if _os.path.exists(fp):
                try:
                    return ImageFont.truetype(fp, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    f_title = load_font(FONT_PATHS, 58)
    f_sub = load_font(FONT_REG, 24)
    f_badge = load_font(FONT_PATHS, 20)
    f_meta = load_font(FONT_REG, 17)
    f_label = load_font(FONT_PATHS, 13)

    # ── Draw badge text ──────────────────────────────────────────────
    draw.text(
        (bx + badge_w // 2 - 30, by + 16),
        icon_txt,
        fill=(255, 255, 255, 255),
        font=f_badge,
    )

    # ── SYNAPSE AI label (top-left) ──────────────────────────────────
    draw.text((52, 46), "SYNAPSE AI", fill=(255, 255, 255, 200), font=f_label)
    draw.rectangle([52, 66, 52 + 60, 68], fill=(*accent, 180))

    # ── Format label pill (below SYNAPSE AI) ─────────────────────────
    draw.text((52, 76), label, fill=(*accent, 255), font=f_label)

    # ── Title (white, wrapping) ──────────────────────────────────────
    title_x, title_y = 52, 155
    words = title.split()
    lines_out, cur = [], ""
    for word in words:
        test = (cur + " " + word).strip()
        if len(test) > 26:
            if cur:
                lines_out.append(cur)
            cur = word
        else:
            cur = test
    if cur:
        lines_out.append(cur)
    lines_out = lines_out[:3]

    for li, line in enumerate(lines_out):
        # Drop shadow
        draw.text(
            (title_x + 2, title_y + li * 72 + 2), line, fill=(0, 0, 0, 60), font=f_title
        )
        draw.text(
            (title_x, title_y + li * 72), line, fill=(255, 255, 255, 255), font=f_title
        )

    # ── Subtitle ─────────────────────────────────────────────────────
    sub_y = title_y + len(lines_out) * 72 + 18
    if subtitle and sub_y < height - 140:
        draw.text((title_x, sub_y), subtitle[:70], fill=(*accent, 220), font=f_sub)

    # ── Bottom strip content ─────────────────────────────────────────
    strip_y = height - 60
    draw.text((52, strip_y + 8), author.upper(), fill=(*badge, 255), font=f_meta)
    right_txt = f"SYNAPSE AI  ·  DIAMOND LEVEL  ·  {label}"
    # right-align: estimate width
    draw.text(
        (width - len(right_txt) * 9, strip_y + 8),
        right_txt,
        fill=(107, 114, 128, 255),
        font=f_meta,
    )

    # Convert RGBA → RGB for saving as JPEG/PNG without alpha issues
    final = Image.new("RGB", (width, height), (255, 255, 255))
    final.paste(img, mask=img.split()[3])
    return final


def _cover_image_bytes(
    title: str,
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    doc_type: str = "pdf",
) -> bytes:
    """Return PNG bytes of the format-specific cover image."""
    img = _generate_cover_image(title, subtitle, author, doc_type=doc_type)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    buf.seek(0)
    return buf.read()


# =============================================================================
# 💎 1. DIAMOND PDF GENERATOR
# =============================================================================


class GeneratePDFInput(BaseModel):
    title: str = Field(..., description="Document title on cover page")
    sections: List[Dict[str, str]] = Field(
        ..., description="List of {'heading','content'} dicts"
    )
    subtitle: str = Field(default="", description="Optional subtitle on cover")
    author: str = Field(
        default="SYNAPSE AI", description="Author shown in metadata/footer"
    )
    user_id: str = Field(default="anonymous", description="User ID for storage path")


def _generate_pdf(
    title: str,
    sections: List[Dict[str, str]],
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    try:
        from reportlab.lib.colors import Color, HexColor, black, white
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm, mm
        from reportlab.platypus import (
            BaseDocTemplate,
            Frame,
            HRFlowable,
        )
        from reportlab.platypus import Image as RLImage
        from reportlab.platypus import (
            KeepTogether,
            NextPageTemplate,
            PageBreak,
            PageTemplate,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
        )
        from reportlab.platypus.flowables import Flowable

        W, H = A4
        LM, RM, TM, BM = 2.0 * cm, 2.0 * cm, 2.8 * cm, 2.2 * cm
        CW = W - LM - RM  # content width

        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.pdf"
        now_str = datetime.now(timezone.utc).strftime("%B %Y")

        def hc(h):
            return HexColor(h)

        # ── Custom Flowables ──────────────────────────────────────────

        class LinearGradientRect(Flowable):
            """Full-width horizontal gradient band."""

            def __init__(self, w, h, c1, c2, vertical=False):
                super().__init__()
                self.width = w
                self.height = h
                self.c1 = c1
                self.c2 = c2
                self.vertical = vertical

            def draw(self):
                steps = 100
                for i in range(steps):
                    t = i / steps
                    r = self.c1.red + t * (self.c2.red - self.c1.red)
                    g = self.c1.green + t * (self.c2.green - self.c1.green)
                    b = self.c1.blue + t * (self.c2.blue - self.c1.blue)
                    self.canv.setFillColorRGB(r, g, b)
                    if self.vertical:
                        self.canv.rect(
                            0,
                            i * self.height / steps,
                            self.width,
                            self.height / steps + 1,
                            fill=1,
                            stroke=0,
                        )
                    else:
                        self.canv.rect(
                            i * self.width / steps,
                            0,
                            self.width / steps + 1,
                            self.height,
                            fill=1,
                            stroke=0,
                        )

        class StatCard(Flowable):
            """Small infographic stat card: icon + big number + label."""

            def __init__(self, value, label, accent, width=3.5 * cm, height=2.2 * cm):
                super().__init__()
                self.value = value
                self.label = label
                self.accent = accent
                self.width = width
                self.height = height

            def draw(self):
                c = self.canv
                # Card background
                c.setFillColor(hc(_Brand.INDIGO_BG))
                c.roundRect(0, 0, self.width, self.height, 6, fill=1, stroke=0)
                # Top accent strip
                c.setFillColor(self.accent)
                c.roundRect(0, self.height - 4, self.width, 4, 2, fill=1, stroke=0)
                # Value
                c.setFillColor(hc(_Brand.DARK))
                c.setFont("Helvetica-Bold", 16)
                c.drawCentredString(self.width / 2, self.height * 0.42, str(self.value))
                # Label
                c.setFillColor(hc(_Brand.GRAY_500))
                c.setFont("Helvetica", 7)
                c.drawCentredString(
                    self.width / 2, self.height * 0.18, self.label.upper()
                )

        class CalloutBox(Flowable):
            """Shaded left-accent callout box."""

            def __init__(self, text, width, bg, accent, style):
                super().__init__()
                self._text = text
                self.width = width
                self._bg = bg
                self._accent = accent
                self._style = style
                from reportlab.platypus import Paragraph as _P

                p = _P(text, style)
                _, self.height = p.wrap(width - 1.4 * cm, 10000)
                self.height += 0.7 * cm

            def draw(self):
                c = self.canv
                c.setFillColor(self._bg)
                c.roundRect(0, 0, self.width, self.height, 5, fill=1, stroke=0)
                c.setFillColor(self._accent)
                c.rect(0, 0, 5, self.height, fill=1, stroke=0)
                from reportlab.platypus import Paragraph as _P

                p = _P(self._text, self._style)
                p.wrapOn(c, self.width - 1.4 * cm, self.height)
                p.drawOn(c, 0.6 * cm, 0.3 * cm)

        class ChartFlowable(Flowable):
            """Embed a Matplotlib PNG chart."""

            def __init__(self, png_bytes, width, height):
                super().__init__()
                self._bytes = png_bytes
                self.width = width
                self.height = height

            def draw(self):
                buf = io.BytesIO(self._bytes)
                self.canv.drawImage(
                    RLImage(buf, width=self.width, height=self.height)._doc,  # noqa
                    0,
                    0,
                    self.width,
                    self.height,
                    preserveAspectRatio=True,
                    mask="auto",
                )

            def wrap(self, aw, ah):
                return self.width, self.height

        # Use ImageReader directly for charts
        def _chart_image(
            chart_type, labels, values, title="", w_cm=12, h_cm=7, dark=False
        ):
            from reportlab.lib.utils import ImageReader

            png = _chart_bytes(
                chart_type,
                labels,
                values,
                title,
                w_cm / 2.54,
                h_cm / 2.54,
                dark_bg=dark,
            )
            buf = io.BytesIO(png)
            return RLImage(buf, width=w_cm * cm, height=h_cm * cm)

        # ── Styles ────────────────────────────────────────────────────
        S = getSampleStyleSheet()

        def ps(name, parent="Normal", **kw):
            return ParagraphStyle(name, parent=S[parent], **kw)

        sCoverTitle = ps(
            "CT",
            "Normal",
            fontSize=40,
            textColor=hc(_Brand.WHITE),
            fontName="Helvetica-Bold",
            leading=48,
            spaceAfter=8,
        )
        sCoverSub = ps(
            "CS",
            "Normal",
            fontSize=16,
            textColor=hc(_Brand.INDIGO_PALE),
            leading=22,
            spaceAfter=4,
        )
        sCoverMeta = ps(
            "CM", "Normal", fontSize=9, textColor=hc(_Brand.INDIGO_LIGHT), leading=14
        )

        sH1 = ps(
            "H1",
            "Normal",
            fontSize=20,
            textColor=hc(_Brand.DARK),
            fontName="Helvetica-Bold",
            spaceBefore=16,
            spaceAfter=4,
            leading=26,
        )
        sBody = ps(
            "Body",
            "Normal",
            fontSize=10.5,
            leading=17.5,
            textColor=hc(_Brand.GRAY_700),
            spaceAfter=9,
            alignment=TA_JUSTIFY,
        )
        sBodyLead = ps(
            "Lead",
            "Normal",
            fontSize=12,
            leading=20,
            textColor=hc(_Brand.GRAY_900),
            spaceAfter=9,
            alignment=TA_JUSTIFY,
            fontName="Helvetica-Bold",
        )
        sCallout = ps(
            "Call", "Normal", fontSize=10, leading=16, textColor=hc(_Brand.DARK)
        )
        sTocNum = ps(
            "TN",
            "Normal",
            fontSize=10,
            textColor=hc(_Brand.INDIGO_MID),
            fontName="Helvetica-Bold",
            alignment=TA_RIGHT,
        )
        sTocTxt = ps(
            "TT",
            "Normal",
            fontSize=10.5,
            textColor=hc(_Brand.DARK),
            fontName="Helvetica-Bold",
        )
        sTocSub = ps("TS", "Normal", fontSize=9, textColor=hc(_Brand.GRAY_500))
        sCaption = ps(
            "Cap",
            "Normal",
            fontSize=8,
            textColor=hc(_Brand.GRAY_500),
            alignment=TA_CENTER,
            spaceAfter=4,
        )

        # ── Page callbacks ────────────────────────────────────────────
        def _on_cover(canv, doc):
            canv.saveState()
            # Full-bleed diagonal gradient (DARK→INDIGO→VIOLET)
            steps = 120
            c1 = hc(_Brand.DARK)
            c2 = hc(_Brand.VIOLET)
            for i in range(steps):
                t = i / steps
                r = c1.red + t * (c2.red - c1.red)
                g = c1.green + t * (c2.green - c1.green)
                b = c1.blue + t * (c2.blue - c1.blue)
                canv.setFillColorRGB(r, g, b)
                canv.rect(0, i * H / steps, W, H / steps + 1, fill=1, stroke=0)

            # Decorative geometric circles (top right)
            for r_, a_ in [(7 * cm, 0.07), (10 * cm, 0.04), (13 * cm, 0.03)]:
                canv.setFillColor(Color(1, 1, 1, alpha=a_))
                canv.circle(W - 2 * cm, H - 1 * cm, r_, fill=1, stroke=0)

            # Bottom white arch strip
            canv.setFillColor(hc(_Brand.WHITE))
            p = canv.beginPath()
            p.moveTo(0, 0)
            p.lineTo(W, 0)
            p.lineTo(W, 2.5 * cm)
            p.curveTo(W * 0.75, 4.5 * cm, W * 0.25, 1.5 * cm, 0, 3.0 * cm)
            p.close()
            canv.drawPath(p, fill=1, stroke=0)

            # Thin INDIGO_LIGHT left accent bar
            canv.setFillColor(hc(_Brand.INDIGO_LIGHT))
            canv.rect(0, 2.5 * cm, 5, H - 2.5 * cm, fill=1, stroke=0)

            # Bottom branding in white strip
            canv.setFont("Helvetica", 8)
            canv.setFillColor(hc(_Brand.GRAY_500))
            canv.drawString(
                LM,
                0.7 * cm,
                f"{author.upper()}  ·  STRICTLY CONFIDENTIAL  ·  {now_str}",
            )
            canv.restoreState()

        def _on_content(canv, doc):
            canv.saveState()
            # Left gradient accent bar
            steps = 60
            c1 = hc(_Brand.INDIGO)
            c2 = hc(_Brand.VIOLET)
            for i in range(steps):
                t = i / steps
                r = c1.red + t * (c2.red - c1.red)
                g = c1.green + t * (c2.green - c1.green)
                b = c1.blue + t * (c2.blue - c1.blue)
                canv.setFillColorRGB(r, g, b)
                canv.rect(0, i * H / steps, 4, H / steps + 1, fill=1, stroke=0)

            # Top header rule + text
            canv.setStrokeColor(hc(_Brand.INDIGO_PALE))
            canv.setLineWidth(0.4)
            canv.line(LM + 0.2 * cm, H - TM + 0.5 * cm, W - RM, H - TM + 0.5 * cm)
            canv.setFont("Helvetica", 7.5)
            canv.setFillColor(hc(_Brand.GRAY_500))
            canv.drawString(LM + 0.2 * cm, H - TM + 0.65 * cm, title.upper()[:60])
            canv.drawRightString(W - RM, H - TM + 0.65 * cm, author.upper())

            # Bottom footer rule + page number
            canv.line(LM + 0.2 * cm, BM - 0.5 * cm, W - RM, BM - 0.5 * cm)
            canv.setFont("Helvetica", 7.5)
            canv.setFillColor(hc(_Brand.GRAY_500))
            canv.drawString(LM + 0.2 * cm, BM - 0.75 * cm, f"{title}  ·  {now_str}")
            # Page badge
            canv.setFillColor(hc(_Brand.INDIGO))
            canv.roundRect(
                W - RM - 0.9 * cm,
                BM - 0.85 * cm,
                0.9 * cm,
                0.5 * cm,
                3,
                fill=1,
                stroke=0,
            )
            canv.setFont("Helvetica-Bold", 8)
            canv.setFillColor(hc(_Brand.WHITE))
            canv.drawCentredString(W - RM - 0.45 * cm, BM - 0.67 * cm, str(doc.page))
            canv.restoreState()

        # ── Document setup ────────────────────────────────────────────
        cover_fr = Frame(
            0,
            0,
            W,
            H,
            leftPadding=LM,
            rightPadding=RM,
            topPadding=TM + 2.5 * cm,
            bottomPadding=BM + 3 * cm,
        )
        body_fr = Frame(
            LM + 0.3 * cm,
            BM,
            W - LM - RM - 0.3 * cm,
            H - TM - BM,
            leftPadding=0,
            rightPadding=0,
            topPadding=0.3 * cm,
            bottomPadding=0.3 * cm,
        )

        tpl_cover = PageTemplate("Cover", [cover_fr], onPage=_on_cover)
        tpl_content = PageTemplate("Content", [body_fr], onPage=_on_content)

        doc = BaseDocTemplate(
            str(file_path),
            pagesize=A4,
            leftMargin=LM,
            rightMargin=RM,
            topMargin=TM,
            bottomMargin=BM,
            pageTemplates=[tpl_cover, tpl_content],
        )

        story = []

        # ── COVER ─────────────────────────────────────────────────────
        story.append(NextPageTemplate("Cover"))
        story.append(Spacer(1, 4.0 * cm))
        story.append(Paragraph(title, sCoverTitle))
        if subtitle:
            story.append(Spacer(1, 0.2 * cm))
            story.append(Paragraph(subtitle, sCoverSub))
        story.append(Spacer(1, 1.2 * cm))
        story.append(
            HRFlowable(
                width="55%", thickness=1.5, color=hc(_Brand.INDIGO_LIGHT), spaceAfter=10
            )
        )
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(f"Prepared by <b>{author}</b>  ·  {now_str}", sCoverMeta)
        )
        story.append(
            Paragraph(
                f"{len(sections)} Chapters  ·  Executive-Grade Analysis", sCoverMeta
            )
        )

        # ── STAT CARDS (cover context panel) ─────────────────────────
        story.append(Spacer(1, 1.8 * cm))
        card_w, card_h = 3.3 * cm, 2.0 * cm
        stat_data = [
            (str(len(sections)), "Chapters"),
            (str(sum(len(s.get("content", "").split()) for s in sections)), "Words"),
            (now_str.split()[0][:3], "Month"),
            ("AI", "Powered"),
        ]
        card_cells = [
            StatCard(v, l, hc(c), card_w, card_h)
            for (v, l), c in zip(stat_data, _Brand.CHART)
        ]
        card_tbl = Table([card_cells], colWidths=[card_w + 0.4 * cm] * len(card_cells))
        card_tbl.setStyle(
            TableStyle(
                [
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 4),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        story.append(card_tbl)

        # ── TOC PAGE ─────────────────────────────────────────────────
        story.append(NextPageTemplate("Content"))
        story.append(PageBreak())

        # TOC header gradient bar
        toc_bar = Table(
            [
                [
                    Paragraph(
                        "TABLE OF CONTENTS",
                        ps(
                            "TH",
                            "Normal",
                            fontSize=14,
                            textColor=hc(_Brand.WHITE),
                            fontName="Helvetica-Bold",
                        ),
                    )
                ]
            ],
            colWidths=[CW],
        )
        toc_bar.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, -1), hc(_Brand.INDIGO_DARK)),
                    ("LEFTPADDING", (0, 0), (-1, -1), 10),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ]
            )
        )
        story.append(toc_bar)
        story.append(Spacer(1, 0.4 * cm))

        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            brief = (
                sec.get("content", "")[:70].replace("&", "&amp;").replace("<", "&lt;")
            )
            row = Table(
                [
                    [Paragraph(f"{i:02d}", sTocNum), Paragraph(h, sTocTxt)],
                    [None, Paragraph(brief + "…", sTocSub)],
                ],
                colWidths=[0.9 * cm, CW - 0.9 * cm],
            )
            row.setStyle(
                TableStyle(
                    [
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("SPAN", (0, 0), (0, 1)),
                        ("LEFTPADDING", (0, 0), (-1, -1), 0),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 0),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 1),
                    ]
                )
            )
            story.append(
                KeepTogether(
                    [
                        row,
                        HRFlowable(
                            width="100%",
                            thickness=0.4,
                            color=hc(_Brand.GRAY_300),
                            spaceAfter=2,
                        ),
                    ]
                )
            )

        # ── SECTION PAGES ─────────────────────────────────────────────
        # Build one bar chart from word counts per section
        sec_labels = [s.get("heading", "")[:18] for s in sections]
        sec_words = [len(s.get("content", "").split()) for s in sections]

        for idx, sec in enumerate(sections):
            story.append(PageBreak())
            heading = sec.get("heading", f"Section {idx+1}")
            content = sec.get("content", "")

            # Section header bar
            hdr = Table(
                [
                    [
                        Paragraph(
                            f"{idx+1:02d}",
                            ps(
                                "SN",
                                "Normal",
                                fontSize=18,
                                textColor=hc(_Brand.WHITE),
                                fontName="Helvetica-Bold",
                                alignment=TA_CENTER,
                            ),
                        ),
                        Paragraph(
                            heading,
                            ps(
                                "SH",
                                "Normal",
                                fontSize=18,
                                textColor=hc(_Brand.WHITE),
                                fontName="Helvetica-Bold",
                            ),
                        ),
                    ]
                ],
                colWidths=[1.1 * cm, CW - 1.1 * cm],
            )
            hdr.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (0, 0), hc(_Brand.INDIGO)),
                        ("BACKGROUND", (1, 0), (1, 0), hc(_Brand.INDIGO_DARK)),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("LEFTPADDING", (0, 0), (0, 0), 6),
                        ("RIGHTPADDING", (0, 0), (0, 0), 6),
                        ("TOPPADDING", (0, 0), (-1, -1), 8),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                        ("LEFTPADDING", (1, 0), (1, 0), 12),
                    ]
                )
            )
            story.append(KeepTogether([hdr]))
            story.append(Spacer(1, 0.35 * cm))

            paras = [p.strip() for p in content.split("\n\n") if p.strip()]

            for pi, para_text in enumerate(paras):
                safe = (
                    para_text.replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                )
                style = sBodyLead if pi == 0 else sBody
                story.append(Paragraph(safe, style))

                # Every 3rd paragraph: insert a callout box with first sentence
                if pi > 0 and pi % 3 == 0:
                    first = para_text.split(".")[0].strip()
                    if len(first) > 25:
                        safe_f = (
                            first.replace("&", "&amp;")
                            .replace("<", "&lt;")
                            .replace(">", "&gt;")
                        )
                        box = CalloutBox(
                            text=f"<b>Key Insight:</b> <i>{safe_f}.</i>",
                            width=CW,
                            bg=hc(_Brand.INDIGO_BG),
                            accent=hc(_Brand.INDIGO),
                            style=sCallout,
                        )
                        story.append(Spacer(1, 0.15 * cm))
                        story.append(box)
                        story.append(Spacer(1, 0.25 * cm))

            # Embed a chart on section 1 (bar: word counts) and section 3 (pie: distribution)
            story.append(Spacer(1, 0.4 * cm))
            if idx == 0 and len(sec_labels) >= 2:
                try:
                    chart = _chart_image(
                        "bar",
                        sec_labels,
                        sec_words,
                        title="Words per Section",
                        w_cm=13,
                        h_cm=7,
                    )
                    story.append(
                        KeepTogether(
                            [
                                chart,
                                Paragraph(
                                    "Figure 1 — Content distribution across all document sections.",
                                    sCaption,
                                ),
                            ]
                        )
                    )
                    story.append(Spacer(1, 0.3 * cm))
                except Exception as ce:
                    logger.warning("Chart embed failed: %s", ce)

            elif idx == 2 and len(sec_labels) >= 2:
                try:
                    chart = _chart_image(
                        "pie",
                        sec_labels,
                        sec_words,
                        title="Section Distribution",
                        w_cm=11,
                        h_cm=7,
                    )
                    story.append(
                        KeepTogether(
                            [
                                chart,
                                Paragraph(
                                    "Figure 2 — Proportional breakdown of document sections.",
                                    sCaption,
                                ),
                            ]
                        )
                    )
                    story.append(Spacer(1, 0.3 * cm))
                except Exception as ce:
                    logger.warning("Chart embed failed: %s", ce)

            elif idx == len(sections) - 1 and len(sec_labels) >= 3:
                try:
                    # Line chart: simulated trend (cumulative word count)
                    cumulative = []
                    running = 0
                    for w in sec_words:
                        running += w
                        cumulative.append(running)
                    chart = _chart_image(
                        "line",
                        sec_labels,
                        cumulative,
                        title="Cumulative Content Growth",
                        w_cm=13,
                        h_cm=6,
                    )
                    story.append(
                        KeepTogether(
                            [
                                chart,
                                Paragraph(
                                    "Figure 3 — Cumulative content growth across document sections.",
                                    sCaption,
                                ),
                            ]
                        )
                    )
                    story.append(Spacer(1, 0.3 * cm))
                except Exception as ce:
                    logger.warning("Line chart embed failed: %s", ce)

            # Section end decorative rule
            story.append(
                HRFlowable(
                    width="35%",
                    thickness=2.5,
                    color=hc(_Brand.INDIGO_LIGHT),
                    spaceAfter=4,
                    hAlign="LEFT",
                )
            )

        doc.build(story)

        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"PDF generated successfully.\nTitle: {title}\n"
            f"Sections: {len(sections)}\nFile: {rel}\n"
            f"Size: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_pdf failed: %s", exc, exc_info=True)
        return f"PDF generation failed: {exc}"


def make_generate_pdf_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_pdf,
        name="generate_pdf",
        description=(
            "Generate a diamond-level PDF: gradient cover with stat cards, embedded bar/pie charts, "
            "callout boxes, numbered section badges, styled TOC, branded header/footer. "
            "Accepts title, sections [{'heading','content'}], subtitle, author. Returns file path."
        ),
        args_schema=GeneratePDFInput,
        return_direct=False,
    )


# =============================================================================
# 💎 2. DIAMOND POWERPOINT GENERATOR
# =============================================================================


class SlideSpec(BaseModel):
    title: str
    bullets: List[str] = Field(default_factory=list)
    content: str = Field(default="")
    notes: str = Field(default="")


class GeneratePPTInput(BaseModel):
    title: str
    subtitle: str = Field(default="Generated by SYNAPSE AI")
    slides: List[SlideSpec] = Field(...)
    author: str = Field(default="SYNAPSE AI")
    user_id: str = Field(default="anonymous")


def _generate_ppt(
    title: str,
    slides: List[Dict[str, Any]],
    subtitle: str = "Generated by SYNAPSE AI",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    try:
        from lxml import etree
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from pptx.util import Emu, Inches, Pt

        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.pptx"
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        prs.core_properties.author = author
        prs.core_properties.title = title
        prs.core_properties.subject = subtitle

        SW, SH = prs.slide_width, prs.slide_height
        blank = prs.slide_layouts[6]
        now_str = datetime.now(timezone.utc).strftime("%B %Y")

        def rgb(h):
            return _Brand.rgb(h)

        # ── XML gradient background helper ─────────────────────────
        def _grad_bg(slide, hex1, hex2, angle=135):
            bg = slide.background
            bgPr = bg._element.find(qn("p:bgPr"))
            if bgPr is None:
                bgPr = etree.SubElement(bg._element, qn("p:bgPr"))
            for child in list(bgPr):
                bgPr.remove(child)
            grad = etree.SubElement(bgPr, qn("a:gradFill"))
            gsLst = etree.SubElement(grad, qn("a:gsLst"))
            for pos, hx in [(0, hex1), (100, hex2)]:
                gs = etree.SubElement(gsLst, qn("a:gs"))
                gs.set("pos", str(pos * 1000))
                sr = etree.SubElement(gs, qn("a:srgbClr"))
                sr.set("val", hx.lstrip("#"))
            lin = etree.SubElement(grad, qn("a:lin"))
            lin.set("ang", str(int(angle * 60000)))
            lin.set("scaled", "0")

        def _solid_bg(slide, h):
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = rgb(h)

        # ── Shape helpers ──────────────────────────────────────────
        def _rect(slide, l, t, w, h, fill, radius=0, line=False):
            s = slide.shapes.add_shape(1 if radius == 0 else 5, l, t, w, h)
            s.fill.solid()
            s.fill.fore_color.rgb = rgb(fill)
            if not line:
                s.line.fill.background()
            return s

        def _tb(
            slide,
            text,
            l,
            t,
            w,
            h,
            sz=18,
            bold=False,
            italic=False,
            color=_Brand.DARK,
            align=PP_ALIGN.LEFT,
            wrap=True,
            font="Calibri",
        ):
            tb = slide.shapes.add_textbox(l, t, w, h)
            tf = tb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(sz)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
            run.font.color.rgb = rgb(color)
            return tb

        def _bullets_tf(
            slide,
            items,
            l,
            t,
            w,
            h,
            sz=15,
            color=_Brand.GRAY_700,
            icon="▸",
            icon_color=_Brand.INDIGO,
        ):
            tb = slide.shapes.add_textbox(l, t, w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.space_before = Pt(5)
                para.space_after = Pt(5)
                # icon
                ir = para.add_run()
                ir.text = f"{icon}  "
                ir.font.size = Pt(sz - 1)
                ir.font.color.rgb = rgb(icon_color)
                ir.font.bold = True
                ir.font.name = "Calibri"
                # text
                tr = para.add_run()
                tr.text = str(item)
                tr.font.size = Pt(sz)
                tr.font.color.rgb = rgb(color)
                tr.font.name = "Calibri"
            return tb

        def _footer(slide, num, total):
            _rect(slide, 0, SH - Inches(0.45), SW, Inches(0.45), _Brand.DARK)
            _tb(
                slide,
                author.upper(),
                Inches(0.35),
                SH - Inches(0.43),
                Inches(4),
                Inches(0.41),
                sz=8,
                bold=True,
                color=_Brand.GRAY_300,
                align=PP_ALIGN.LEFT,
            )
            _tb(
                slide,
                title,
                Inches(4),
                SH - Inches(0.43),
                SW - Inches(5.5),
                Inches(0.41),
                sz=8,
                color=_Brand.GRAY_500,
                align=PP_ALIGN.CENTER,
            )
            _tb(
                slide,
                f"{num} / {total}",
                SW - Inches(1.4),
                SH - Inches(0.43),
                Inches(1.2),
                Inches(0.4),
                sz=9,
                bold=True,
                color=_Brand.INDIGO_PALE,
                align=PP_ALIGN.RIGHT,
            )

        def _embed_chart_png(slide, png_bytes, l, t, w, h):
            """Embed a PNG chart image on a slide from bytes."""
            buf = io.BytesIO(png_bytes)
            slide.shapes.add_picture(buf, l, t, w, h)

        total = len(slides) + 3

        # ══════════════════════════════════════════════════════════
        # SLIDE 1 — TITLE (gradient + geometry)
        # ══════════════════════════════════════════════════════════
        s1 = prs.slides.add_slide(blank)
        _grad_bg(s1, _Brand.DARK, _Brand.INDIGO, angle=150)

        # Decorative translucent circles via shape + XML alpha
        for cx, cy, cr, alpha in [
            (SW - Inches(1), -Inches(0.5), Inches(5.5), "6000"),
            (SW + Inches(1), Inches(1), Inches(8), "3500"),
            (-Inches(0.5), SH - Inches(1), Inches(4), "5000"),
        ]:
            c = s1.shapes.add_shape(9, cx - cr, cy - cr, cr * 2, cr * 2)
            spPr = c._element.find(qn("p:spPr"))
            # build solidFill with alpha
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            sfill = etree.Element(f"{{{ns_a}}}solidFill")
            srgb = etree.SubElement(sfill, f"{{{ns_a}}}srgbClr")
            srgb.set("val", _Brand.INDIGO_MID.lstrip("#"))
            alph = etree.SubElement(srgb, f"{{{ns_a}}}alpha")
            alph.set("val", alpha)
            existing = spPr.find(f"{{{ns_a}}}solidFill")
            if existing is not None:
                spPr.remove(existing)
            gradFill = spPr.find(f"{{{ns_a}}}gradFill")
            if gradFill is not None:
                spPr.remove(gradFill)
            spPr.insert(0, sfill)
            c.line.fill.background()

        # Left vertical accent bar (gradient via XML is complex, use two rects)
        _rect(s1, 0, 0, Inches(0.12), SH // 2, _Brand.INDIGO_LIGHT)
        _rect(s1, 0, SH // 2, Inches(0.12), SH // 2, _Brand.VIOLET_LIGHT)

        # Bottom white strip
        _rect(s1, 0, SH - Inches(1.1), SW, Inches(1.1), _Brand.WHITE)

        # Thin indigo line above white strip
        _rect(s1, 0, SH - Inches(1.12), SW, Inches(0.04), _Brand.INDIGO_LIGHT)

        # Title + subtitle
        _tb(
            s1,
            title,
            Inches(0.7),
            Inches(1.6),
            Inches(9),
            Inches(2.2),
            sz=46,
            bold=True,
            color=_Brand.WHITE,
            align=PP_ALIGN.LEFT,
        )
        _tb(
            s1,
            subtitle,
            Inches(0.7),
            Inches(3.9),
            Inches(9),
            Inches(0.8),
            sz=19,
            color=_Brand.INDIGO_PALE,
            align=PP_ALIGN.LEFT,
        )

        # Divider line
        _rect(
            s1,
            Inches(0.7),
            Inches(3.72),
            Inches(5.5),
            Inches(0.04),
            _Brand.INDIGO_LIGHT,
        )

        # Meta in white strip
        _tb(
            s1,
            author,
            Inches(0.7),
            SH - Inches(1.0),
            Inches(5),
            Inches(0.45),
            sz=11,
            bold=True,
            color=_Brand.GRAY_700,
        )
        _tb(
            s1,
            now_str,
            Inches(5.5),
            SH - Inches(1.0),
            Inches(4),
            Inches(0.45),
            sz=11,
            color=_Brand.GRAY_500,
        )

        # ══════════════════════════════════════════════════════════
        # SLIDE 2 — AGENDA (two columns)
        # ══════════════════════════════════════════════════════════
        s2 = prs.slides.add_slide(blank)
        _solid_bg(s2, _Brand.WHITE)
        _rect(s2, 0, 0, Inches(0.12), SH, _Brand.INDIGO)
        _rect(s2, Inches(0.12), 0, SW - Inches(0.12), Inches(1.2), _Brand.DARK)
        _tb(
            s2,
            "AGENDA & OVERVIEW",
            Inches(0.5),
            Inches(0.15),
            SW - Inches(1),
            Inches(0.95),
            sz=30,
            bold=True,
            color=_Brand.WHITE,
        )

        all_titles = [
            (sl.get("title") if isinstance(sl, dict) else getattr(sl, "title", ""))
            for sl in slides
        ]
        mid = math.ceil(len(all_titles) / 2)
        _bullets_tf(
            s2,
            all_titles[:mid],
            Inches(0.5),
            Inches(1.4),
            Inches(5.9),
            SH - Inches(2.0),
            sz=15,
            color=_Brand.WHITE,
            icon="◆",
            icon_color=_Brand.INDIGO_LIGHT,
        )
        if all_titles[mid:]:
            _bullets_tf(
                s2,
                all_titles[mid:],
                Inches(6.9),
                Inches(1.4),
                Inches(5.9),
                SH - Inches(2.0),
                sz=15,
                color=_Brand.WHITE,
                icon="◆",
                icon_color=_Brand.INDIGO_LIGHT,
            )
        _footer(s2, 2, total)

        # ══════════════════════════════════════════════════════════
        # STATS DASHBOARD SLIDE (slide 3)
        # ══════════════════════════════════════════════════════════
        s3 = prs.slides.add_slide(blank)
        _grad_bg(s3, _Brand.INDIGO_BG, _Brand.WHITE, angle=160)
        _rect(s3, 0, 0, Inches(0.12), SH, _Brand.INDIGO)
        _rect(s3, Inches(0.12), 0, SW - Inches(0.12), Inches(1.2), _Brand.INDIGO_DARK)
        _tb(
            s3,
            "AT A GLANCE",
            Inches(0.5),
            Inches(0.15),
            SW - Inches(1),
            Inches(0.95),
            sz=28,
            bold=True,
            color=_Brand.WHITE,
        )

        # Stat cards (4 across)
        stat_defs = [
            (str(len(slides)), "Topics Covered", _Brand.INDIGO),
            (
                str(
                    sum(
                        len(
                            (
                                sl.get("content", "")
                                if isinstance(sl, dict)
                                else getattr(sl, "content", "")
                            ).split()
                        )
                        for sl in slides
                    )
                ),
                "Total Words",
                _Brand.VIOLET,
            ),
            (now_str.split()[0][:3].upper(), "Month", _Brand.CYAN),
            ("AI", "Powered By", _Brand.EMERALD),
        ]
        card_w, card_h = Inches(2.8), Inches(1.6)
        gap = Inches(0.35)
        start_x = Inches(0.6)
        for ci, (val, lbl, col) in enumerate(stat_defs):
            cx = start_x + ci * (card_w + gap)
            cy = Inches(1.55)
            # card bg
            bg_shp = s3.shapes.add_shape(5, cx, cy, card_w, card_h)
            bg_shp.fill.solid()
            bg_shp.fill.fore_color.rgb = rgb(_Brand.WHITE)
            bg_shp.line.color.rgb = rgb(col)
            bg_shp.line.width = Pt(1.5)
            # top accent
            _rect(s3, cx, cy, card_w, Inches(0.08), col)
            # value
            _tb(
                s3,
                val,
                cx,
                cy + Inches(0.15),
                card_w,
                Inches(0.85),
                sz=34,
                bold=True,
                color=col,
                align=PP_ALIGN.CENTER,
            )
            # label
            _tb(
                s3,
                lbl.upper(),
                cx,
                cy + Inches(1.05),
                card_w,
                Inches(0.45),
                sz=9,
                bold=True,
                color=_Brand.GRAY_500,
                align=PP_ALIGN.CENTER,
            )

        # Embed a bar chart below stats
        try:
            sec_labels = [
                (sl.get("title") if isinstance(sl, dict) else getattr(sl, "title", ""))[
                    :15
                ]
                for sl in slides
            ]
            sec_words = [
                len(
                    (
                        sl.get("content", "")
                        if isinstance(sl, dict)
                        else getattr(sl, "content", "")
                    ).split()
                )
                for sl in slides
            ]
            if any(w > 0 for w in sec_words):
                png = _chart_bytes(
                    "bar",
                    sec_labels,
                    sec_words,
                    "Content Distribution by Section",
                    width_in=10,
                    height_in=3.2,
                    dark_bg=False,
                )
                _embed_chart_png(
                    s3, png, Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.6)
                )
        except Exception as ce:
            logger.warning("Stats chart: %s", ce)

        _footer(s3, 3, total)

        # ══════════════════════════════════════════════════════════
        # CONTENT SLIDES
        # ══════════════════════════════════════════════════════════
        for sn, slide_data in enumerate(slides, start=4):
            sd = slide_data if isinstance(slide_data, dict) else slide_data.__dict__
            s_title = sd.get("title", "Slide")
            s_bullets = sd.get("bullets", [])
            s_content = sd.get("content", "")
            s_notes = sd.get("notes", "")

            # Derive bullets from content if missing
            if not s_bullets and s_content:
                paras = [p.strip() for p in s_content.split("\n\n") if p.strip()]
                s_bullets = [
                    (p.split(".")[0].strip() + ".")
                    for p in paras[:6]
                    if p.split(".")[0].strip()
                ]

            sl = prs.slides.add_slide(blank)
            _solid_bg(sl, _Brand.WHITE)

            # Left gradient accent: two colour rects approximation
            _rect(sl, 0, 0, Inches(0.12), SH // 2, _Brand.INDIGO)
            _rect(sl, 0, SH // 2, Inches(0.12), SH // 2, _Brand.VIOLET)

            # Header bar (dark)
            _rect(sl, Inches(0.12), 0, SW - Inches(0.12), Inches(1.18), _Brand.DARK)

            # Section title in header
            _tb(
                sl,
                s_title,
                Inches(0.5),
                Inches(0.14),
                SW - Inches(2.2),
                Inches(0.92),
                sz=26,
                bold=True,
                color=_Brand.WHITE,
            )

            # Number badge (right of header)
            _rect(sl, SW - Inches(1.25), 0, Inches(1.25), Inches(1.18), _Brand.INDIGO)
            _tb(
                sl,
                str(sn - 3),
                SW - Inches(1.25),
                Inches(0.1),
                Inches(1.25),
                Inches(1.0),
                sz=30,
                bold=True,
                color=_Brand.WHITE,
                align=PP_ALIGN.CENTER,
            )

            # Two-column layout: bullets left (wide), accent right panel
            content_w = SW - Inches(3.9)
            # Add white background rectangle for content area
            _rect(
                sl,
                Inches(0.42),
                Inches(1.18),
                content_w,
                SH - Inches(1.63),
                _Brand.WHITE,
            )
            _bullets_tf(
                sl,
                s_bullets,
                Inches(0.42),
                Inches(1.35),
                content_w,
                SH - Inches(2.1),
                sz=15,
                color=_Brand.DARK,
                icon="▶",
                icon_color=_Brand.INDIGO_MID,
            )

            # Right decorative panel
            _rect(
                sl,
                SW - Inches(3.3),
                Inches(1.18),
                Inches(3.3),
                SH - Inches(1.63),
                _Brand.INDIGO_BG,
            )
            # Panel label
            _tb(
                sl,
                s_title[:28],
                SW - Inches(3.2),
                Inches(1.35),
                Inches(3.1),
                Inches(0.9),
                sz=11,
                bold=True,
                color=_Brand.INDIGO_DARK,
                align=PP_ALIGN.CENTER,
            )

            # Mini pie chart in right panel (section word distribution)
            try:
                if sec_words and any(w > 0 for w in sec_words):
                    mini_png = _chart_bytes(
                        "pie", sec_labels, sec_words, "", width_in=2.8, height_in=2.8
                    )
                    _embed_chart_png(
                        sl,
                        mini_png,
                        SW - Inches(3.2),
                        Inches(2.4),
                        Inches(3.0),
                        Inches(3.0),
                    )
            except Exception as ce:
                logger.warning("Mini pie: %s", ce)

            # Separator line between columns
            _rect(
                sl,
                SW - Inches(3.35),
                Inches(1.18),
                Inches(0.03),
                SH - Inches(1.63),
                _Brand.GRAY_300,
            )

            if s_notes:
                sl.notes_slide.notes_text_frame.text = s_notes

            _footer(sl, sn - 3, len(slides))

        # ══════════════════════════════════════════════════════════
        # CHART SLIDE — dedicated full-page bar chart
        # ══════════════════════════════════════════════════════════
        sc = prs.slides.add_slide(blank)
        _solid_bg(sc, _Brand.WHITE)
        _rect(sc, 0, 0, Inches(0.12), SH, _Brand.INDIGO)
        _rect(sc, Inches(0.12), 0, SW - Inches(0.12), Inches(1.18), _Brand.INDIGO_DARK)
        _tb(
            sc,
            "DATA OVERVIEW",
            Inches(0.5),
            Inches(0.14),
            SW - Inches(1.5),
            Inches(0.92),
            sz=28,
            bold=True,
            color=_Brand.WHITE,
        )
        try:
            full_png = _chart_bytes(
                "bar",
                sec_labels,
                sec_words,
                "Content Volume by Topic",
                width_in=12,
                height_in=5.5,
                dark_bg=True,
            )
            _embed_chart_png(
                sc, full_png, Inches(0.5), Inches(1.3), Inches(6.1), Inches(5.7)
            )
        except Exception as ce:
            logger.warning("Full chart: %s", ce)

        # Second chart: line chart (cumulative trend)
        try:
            cumulative = []
            running = 0
            for w in sec_words:
                running += w
                cumulative.append(running)
            if any(c > 0 for c in cumulative):
                line_png = _chart_bytes(
                    "line",
                    sec_labels,
                    cumulative,
                    "Cumulative Content Growth",
                    width_in=5.5,
                    height_in=3.0,
                    dark_bg=True,
                )
                _embed_chart_png(
                    sc, line_png, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.8)
                )
        except Exception as ce:
            logger.warning("Line chart ppt: %s", ce)

        _footer(sc, len(slides) + 1, total)

        # ══════════════════════════════════════════════════════════
        # CLOSING SLIDE
        # ══════════════════════════════════════════════════════════
        sf = prs.slides.add_slide(blank)
        _grad_bg(sf, _Brand.DARK, _Brand.VIOLET, angle=155)
        _rect(sf, 0, 0, Inches(0.12), SH, _Brand.INDIGO_LIGHT)
        _tb(
            sf,
            "Thank You",
            Inches(0.7),
            Inches(1.8),
            Inches(10),
            Inches(1.8),
            sz=56,
            bold=True,
            color=_Brand.WHITE,
            align=PP_ALIGN.LEFT,
        )
        _rect(
            sf, Inches(0.7), Inches(3.75), Inches(6), Inches(0.04), _Brand.INDIGO_LIGHT
        )
        _tb(
            sf,
            title,
            Inches(0.7),
            Inches(3.9),
            Inches(10),
            Inches(0.75),
            sz=17,
            color=_Brand.INDIGO_PALE,
        )
        _tb(
            sf,
            f"Prepared by {author}  ·  {now_str}",
            Inches(0.7),
            SH - Inches(1.15),
            Inches(8),
            Inches(0.5),
            sz=11,
            color=_Brand.GRAY_300,
        )

        prs.save(str(file_path))

        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"PowerPoint generated successfully.\nTitle: {title}\n"
            f"Slides: {len(prs.slides)} (title+agenda+stats+content+chart+closing)\n"
            f"File: {rel}\nSize: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_ppt failed: %s", exc, exc_info=True)
        return f"PowerPoint generation failed: {exc}"


def make_generate_ppt_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_ppt,
        name="generate_ppt",
        description=(
            "Generate a diamond-level 16:9 PowerPoint: gradient title slide, 2-col agenda, "
            "stats dashboard with bar chart, icon-bullet content slides with mini pie charts, "
            "full-page chart slide, closing slide. Returns file path."
        ),
        args_schema=GeneratePPTInput,
        return_direct=False,
    )


# =============================================================================
# 💎 3. DIAMOND WORD / DOCX GENERATOR
# =============================================================================


class WordSection(BaseModel):
    heading: str
    content: str
    level: int = Field(default=1, ge=1, le=3)


class GenerateWordDocInput(BaseModel):
    title: str
    sections: List[WordSection]
    author: str = Field(default="SYNAPSE AI")
    add_toc: bool = Field(default=True)
    user_id: str = Field(default="anonymous")


def _generate_word_doc(
    title: str,
    sections: List[Dict[str, Any]],
    author: str = "SYNAPSE AI",
    add_toc: bool = True,
    user_id: str = "anonymous",
) -> str:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Cm, Inches, Pt, RGBColor

        def rgb(h):
            h = h.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.docx"
        now_str = datetime.now(timezone.utc).strftime("%B %d, %Y")

        doc = Document()
        doc.core_properties.author = author
        doc.core_properties.title = title
        doc.core_properties.subject = f"Generated by {author}"

        # Page layout A4, comfortable margins
        for sec in doc.sections:
            sec.page_width = Cm(21)
            sec.page_height = Cm(29.7)
            sec.top_margin = Cm(2.4)
            sec.bottom_margin = Cm(2.4)
            sec.left_margin = Cm(2.8)
            sec.right_margin = Cm(2.5)

        # ── XML helpers ──────────────────────────────────────────────
        def _shade(para, fill_hex):
            pPr = para._p.get_or_add_pPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), fill_hex.lstrip("#"))
            pPr.append(shd)

        def _border(para, sides, color_hex, size=6):
            pPr = para._p.get_or_add_pPr()
            pBdr = pPr.find(qn("w:pBdr"))
            if pBdr is None:
                pBdr = OxmlElement("w:pBdr")
                pPr.append(pBdr)
            for side in sides:
                el = OxmlElement(f"w:{side}")
                el.set(qn("w:val"), "single")
                el.set(qn("w:sz"), str(size))
                el.set(qn("w:space"), "4")
                el.set(qn("w:color"), color_hex.lstrip("#"))
                pBdr.append(el)

        def _run(
            para, text, bold=False, italic=False, color=None, size=None, font="Calibri"
        ):
            r = para.add_run(text)
            r.font.name = font
            r.font.bold = bold
            r.font.italic = italic
            if color:
                r.font.color.rgb = rgb(color)
            if size:
                r.font.size = Pt(size)
            return r

        def _spacer(doc, pts=4):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after = Pt(pts)

        def _hr(doc, color=_Brand.GRAY_300, sz=4):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(1)
            p.paragraph_format.space_after = Pt(1)
            _border(p, ["bottom"], color, sz)

        def _callout(doc, text, bg, accent):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(5)
            p.paragraph_format.space_after = Pt(5)
            p.paragraph_format.left_indent = Cm(0.4)
            p.paragraph_format.right_indent = Cm(0.4)
            _shade(p, bg)
            _border(p, ["left"], accent, 24)
            r = p.add_run(text)
            r.font.name = "Calibri"
            r.font.size = Pt(10.5)
            r.font.italic = True
            r.font.color.rgb = rgb(_Brand.DARK)

        def _insert_toc(doc):
            p = doc.add_paragraph()
            for el_def in [
                ("w:fldChar", {"w:fldCharType": "begin"}),
                ("w:instrText", {}, ' TOC \\o "1-3" \\h \\z \\u '),
                ("w:fldChar", {"w:fldCharType": "separate"}),
                ("w:fldChar", {"w:fldCharType": "end"}),
            ]:
                tag = el_def[0]
                attrs = el_def[1] if len(el_def) > 1 else {}
                text = el_def[2] if len(el_def) > 2 else None
                r = OxmlElement("w:r")
                el = OxmlElement(tag)
                for k, v in attrs.items():
                    el.set(qn(k), v)
                if text:
                    el.text = text
                    el.set(qn("xml:space"), "preserve")
                r.append(el)
                p._p.append(r)
            p.paragraph_format.space_after = Pt(0)

        def _embed_chart(doc, png_bytes, width_cm=14, caption=""):
            """Add a chart image to the document."""
            buf = io.BytesIO(png_bytes)
            doc.add_picture(buf, width=Cm(width_cm))
            if caption:
                cp = doc.add_paragraph(caption)
                cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cp.paragraph_format.space_before = Pt(2)
                cp.paragraph_format.space_after = Pt(6)
                for r in cp.runs:
                    r.font.size = Pt(8.5)
                    r.font.color.rgb = rgb(_Brand.GRAY_500)
                    r.font.italic = True

        # ── COVER ────────────────────────────────────────────────────
        # Title shaded block
        ct = doc.add_paragraph()
        ct.alignment = WD_ALIGN_PARAGRAPH.LEFT
        ct.paragraph_format.space_before = Pt(60)
        ct.paragraph_format.space_after = Pt(0)
        _shade(ct, _Brand.INDIGO_DARK.lstrip("#"))
        _border(ct, ["left"], _Brand.INDIGO_LIGHT, 36)
        r = ct.add_run("  " + title)
        r.font.name = "Calibri"
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = rgb(_Brand.WHITE)

        # Subtitle bar
        sb = doc.add_paragraph()
        sb.alignment = WD_ALIGN_PARAGRAPH.LEFT
        sb.paragraph_format.space_before = Pt(0)
        sb.paragraph_format.space_after = Pt(0)
        _shade(sb, _Brand.INDIGO.lstrip("#"))
        r2 = sb.add_run(f"  {author}  ·  {now_str}  ·  Comprehensive Analysis")
        r2.font.name = "Calibri"
        r2.font.size = Pt(11)
        r2.font.color.rgb = rgb(_Brand.INDIGO_PALE)

        # Stats row table on cover
        _spacer(doc, 10)
        stats = [
            (str(len(sections)), "Chapters"),
            (
                str(
                    sum(
                        (
                            len(s.get("content", "").split())
                            if isinstance(s, dict)
                            else len(getattr(s, "content", "").split())
                        )
                        for s in sections
                    )
                ),
                "Words",
            ),
            (now_str.split(",")[0], "Date"),
            ("SYNAPSE", "Platform"),
        ]
        tbl = doc.add_table(rows=1, cols=len(stats))
        tbl.style = "Table Grid"
        for ci, (val, lbl) in enumerate(stats):
            cell = tbl.cell(0, ci)
            cell.paragraphs[0].clear()
            # shade cell
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), _Brand.INDIGO_BG.lstrip("#"))
            tcPr.append(shd)
            p_val = cell.add_paragraph()
            p_val.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _run(
                p_val,
                val,
                bold=True,
                color=_Brand.CHART[ci % len(_Brand.CHART)],
                size=20,
            )
            p_lbl = cell.add_paragraph()
            p_lbl.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _run(p_lbl, lbl.upper(), size=8, color=_Brand.GRAY_500)

        _spacer(doc, 8)

        # Cover chart — bar of section word counts
        sec_labels = []
        sec_words = []
        for s in sections:
            h = (
                s.get("heading", "")
                if isinstance(s, dict)
                else getattr(s, "heading", "")
            )
            c = (
                s.get("content", "")
                if isinstance(s, dict)
                else getattr(s, "content", "")
            )
            sec_labels.append(h[:18])
            sec_words.append(len(c.split()))

        try:
            if any(w > 0 for w in sec_words):
                png = _chart_bytes(
                    "hbar",
                    sec_labels,
                    sec_words,
                    "Content Volume by Chapter",
                    width_in=6,
                    height_in=3.5,
                )
                _embed_chart(
                    doc,
                    png,
                    width_cm=13,
                    caption="Figure 1 — Content distribution across chapters.",
                )
        except Exception as ce:
            logger.warning("Cover chart: %s", ce)

        doc.add_page_break()

        # ── TOC ──────────────────────────────────────────────────────
        if add_toc:
            th = doc.add_paragraph()
            _shade(th, _Brand.INDIGO_DARK.lstrip("#"))
            _run(th, "  TABLE OF CONTENTS", bold=True, color=_Brand.WHITE, size=14)
            _spacer(doc, 6)
            _insert_toc(doc)
            _spacer(doc, 6)
            doc.add_page_break()

        # ── SECTIONS ─────────────────────────────────────────────────
        for si, section_data in enumerate(sections):
            if isinstance(section_data, dict):
                heading = section_data.get("heading", "Section")
                content = section_data.get("content", "")
                level = int(section_data.get("level", 1))
            else:
                heading = getattr(section_data, "heading", "Section")
                content = getattr(section_data, "content", "")
                level = int(getattr(section_data, "level", 1))

            # Section badge header
            badge = doc.add_paragraph()
            badge.paragraph_format.space_before = Pt(2 if si else 0)
            badge.paragraph_format.space_after = Pt(0)
            _shade(badge, _Brand.INDIGO.lstrip("#"))
            _border(badge, ["left"], _Brand.INDIGO_LIGHT, 32)
            _run(
                badge,
                f"  {si+1:02d}  {heading.upper()}  ",
                bold=True,
                color=_Brand.WHITE,
                size=13,
            )

            _spacer(doc, 4)

            paras = [p.strip() for p in content.split("\n\n") if p.strip()]
            for pi, para_text in enumerate(paras):
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                p.paragraph_format.space_before = Pt(0)
                p.paragraph_format.space_after = Pt(8)
                if pi == 0:
                    _border(p, ["left"], _Brand.INDIGO_PALE, 12)
                    _run(p, para_text, size=11.5, color=_Brand.GRAY_900)
                else:
                    _run(p, para_text, size=11, color=_Brand.GRAY_700)

                if pi > 0 and pi % 3 == 0:
                    first = para_text.split(".")[0].strip()
                    if len(first) > 30:
                        _callout(
                            doc,
                            f"Key Insight: {first}.",
                            _Brand.INDIGO_BG,
                            _Brand.INDIGO,
                        )

            # Embed a pie chart after section 2
            if si == 1 and any(w > 0 for w in sec_words):
                try:
                    png2 = _chart_bytes(
                        "pie",
                        sec_labels,
                        sec_words,
                        "Section Distribution",
                        width_in=5.5,
                        height_in=4,
                    )
                    _embed_chart(
                        doc,
                        png2,
                        width_cm=11,
                        caption="Figure 2 — Proportional breakdown of document sections.",
                    )
                except Exception as ce:
                    logger.warning("Pie chart: %s", ce)

            # Section end rule
            _hr(doc, _Brand.INDIGO_PALE, 4)
            _spacer(doc, 6)

        # ── FOOTER ──────────────────────────────────────────────────
        for sec in doc.sections:
            fp = sec.footer.paragraphs[0]
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            fp.clear()
            _border(fp, ["top"], _Brand.GRAY_300, 4)
            _run(fp, f"{title}  ·  {author}  ·  Page ", size=8, color=_Brand.GRAY_500)
            # PAGE field
            fR = OxmlElement("w:r")
            for tag, attrs, txt in [
                ("w:fldChar", {"w:fldCharType": "begin"}, None),
                ("w:instrText", {"xml:space": "preserve"}, " PAGE "),
                ("w:fldChar", {"w:fldCharType": "end"}, None),
            ]:
                el = OxmlElement(tag)
                for k, v in attrs.items():
                    el.set(qn(k), v)
                if txt:
                    el.text = txt
                fR.append(el)
            fp._p.append(fR)

        doc.save(str(file_path))

        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"Word document generated successfully.\nTitle: {title}\n"
            f"Sections: {len(sections)}\nFile: {rel}\n"
            f"Size: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_word_doc failed: %s", exc, exc_info=True)
        return f"Word document generation failed: {exc}"


def make_generate_word_doc_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_word_doc,
        name="generate_word_doc",
        description=(
            "Generate a diamond-level Word document: shaded cover with stats table and bar chart, "
            "numbered section badges, callout boxes, embedded pie chart, TOC field, page-number footer. "
            "Returns file path of .docx."
        ),
        args_schema=GenerateWordDocInput,
        return_direct=False,
    )


# =============================================================================
# 💎 4. DIAMOND MARKDOWN GENERATOR
# =============================================================================


class GenerateMarkdownInput(BaseModel):
    title: str
    sections: List[Dict[str, str]]
    author: str = Field(default="SYNAPSE AI")
    user_id: str = Field(default="anonymous")


def _generate_markdown(
    title: str,
    sections: List[Dict[str, str]],
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    try:
        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.md"
        now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        now_human = datetime.now(timezone.utc).strftime("%B %d, %Y")

        ICONS = [
            "🔍",
            "📌",
            "💡",
            "📊",
            "⚡",
            "🎯",
            "🔬",
            "🌐",
            "🚀",
            "✅",
            "🧠",
            "🔑",
            "📈",
            "🏆",
            "🔧",
        ]

        total_words = sum(len(s.get("content", "").split()) for s in sections)
        avg_words = total_words // max(len(sections), 1)

        L = []

        # ── YAML Front Matter ─────────────────────────────────────
        L += [
            "---",
            f'title: "{title}"',
            f'author: "{author}"',
            f'date: "{now_human}"',
            f'generated: "{now_str}"',
            f"sections: {len(sections)}",
            f"total_words: {total_words}",
            f"tags: [synapse-ai, report, analysis, research]",
            f'version: "2.0"',
            f'classification: "Confidential"',
            f'status: "final"',
            "---",
            "",
        ]

        # ── Hero banner ───────────────────────────────────────────
        L += [
            f"# {title}",
            "",
            "<!-- SYNAPSE AI — Diamond Level Document -->",
            "",
            f"**{author}** &nbsp;|&nbsp; {now_human} &nbsp;|&nbsp; "
            f"Executive-Grade Research &nbsp;|&nbsp; `v2.0`",
            "",
            "---",
            "",
        ]

        # ── Badge row ─────────────────────────────────────────────
        badge_base = "https://img.shields.io/badge"
        author_enc = author.replace(" ", "%20").replace("-", "--")
        L += [
            f"[![Author]({badge_base}/Author-{author_enc}-4F46E5?style=for-the-badge&logo=data:image/svg+xml;base64,PHN2Zy8+)]() "
            f"[![Sections]({badge_base}/Sections-{len(sections)}-7C3AED?style=for-the-badge)]() "
            f"[![Words]({badge_base}/Words-{total_words:,}-06B6D4?style=for-the-badge)]() "
            f"[![Status]({badge_base}/Status-Final-10B981?style=for-the-badge)]() "
            f"[![Version]({badge_base}/Version-2.0-F59E0B?style=for-the-badge)]()",
            "",
            "---",
            "",
        ]

        # ── Stats cards (ASCII art) ───────────────────────────────
        L += [
            "## 📊 Document Snapshot",
            "",
            "```",
            "┌─────────────────┬─────────────────┬─────────────────┬─────────────────┐",
            f"│   📚 SECTIONS   │   📝 WORDS      │   📅 DATE       │   🤖 ENGINE     │",
            "├─────────────────┼─────────────────┼─────────────────┼─────────────────┤",
            f"│  {str(len(sections)).center(15)}│  {str(f'{total_words:,}').center(15)}│  {now_human[:10].center(15)}│  {'SYNAPSE AI'.center(15)}│",
            "├─────────────────┼─────────────────┼─────────────────┼─────────────────┤",
            f"│   Avg {avg_words:,} w/sec  │  Research Grade │  {datetime.now(timezone.utc).strftime('%Y').center(15)}│  Diamond Level  │",
            "└─────────────────┴─────────────────┴─────────────────┴─────────────────┘",
            "```",
            "",
            "---",
            "",
        ]

        # ── Executive overview blockquote ─────────────────────────
        if sections:
            first_para = sections[0].get("content", "").split("\n\n")[0].strip()[:280]
            if first_para:
                L += [
                    "> ### 🔎 Executive Overview",
                    ">",
                ]
                words = first_para.split()
                line_buf, cc = [], 0
                for w in words:
                    if cc + len(w) + 1 > 78:
                        L.append(f"> {' '.join(line_buf)}")
                        line_buf, cc = [w], len(w)
                    else:
                        line_buf.append(w)
                        cc += len(w) + 1
                if line_buf:
                    L.append(f"> {' '.join(line_buf)}")
                L += [">", ""]

        # ── TOC ───────────────────────────────────────────────────
        L += ["## 📚 Table of Contents", ""]
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            icon = ICONS[(i - 1) % len(ICONS)]
            anch = (
                h.lower()
                .replace(" ", "-")
                .replace("/", "")
                .replace("&", "")
                .replace("(", "")
                .replace(")", "")
                .replace(",", "")
                .replace(".", "")
            )
            wc = len(sec.get("content", "").split())
            L.append(f"{i}. [{icon} {h}](#{anch}) _{wc:,} words_")
        L += ["", "---", ""]

        # ── Content sections ──────────────────────────────────────
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            body = sec.get("content", "")
            icon = ICONS[(i - 1) % len(ICONS)]
            wc = len(body.split())

            L += [
                f"## {icon} {h}",
                "",
                f"> *Section {i} of {len(sections)} &nbsp;·&nbsp; "
                f"{wc:,} words &nbsp;·&nbsp; {author}*",
                "",
            ]

            paras = [p.strip() for p in body.split("\n\n") if p.strip()]
            for pi, para in enumerate(paras):
                # First para: emphasis lead
                if pi == 0:
                    L += [f"**{para}**", ""]
                else:
                    L += [para, ""]

                # Every 3rd para: admonition callout
                if pi > 0 and pi % 3 == 0:
                    first_s = para.split(".")[0].strip()
                    if len(first_s) > 20:
                        L += [
                            "> [!TIP]",
                            f"> 💡 **Key Insight:** {first_s}.",
                            "",
                        ]

            # Section progress bar (ASCII)
            filled = round((i / len(sections)) * 20)
            bar = "█" * filled + "░" * (20 - filled)
            L += [
                f"```",
                f"Progress: [{bar}] {i}/{len(sections)} sections",
                f"```",
                "",
            ]

            if i < len(sections):
                L += ["---", ""]

        # ── Summary stats table ────────────────────────────────────
        L += [
            "---",
            "",
            "## 📋 Section Summary",
            "",
            "| # | Section | Words | Status | Coverage |",
            "|:-:|---------|------:|:------:|:--------:|",
        ]
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            icon = ICONS[(i - 1) % len(ICONS)]
            wc = len(sec.get("content", "").split())
            pct = round(wc / max(total_words, 1) * 100, 1)
            bars = "▓" * round(pct / 5) + "░" * (20 - round(pct / 5))
            L.append(f"| {i} | {icon} {h} | {wc:,} | ✅ Complete | `{bars}` {pct}% |")

        L += ["", "---", ""]

        # ── Document metadata footer ───────────────────────────────
        L += [
            "## 📎 Document Metadata",
            "",
            "| Field | Value |",
            "|-------|-------|",
            f"| **Title** | {title} |",
            f"| **Author** | {author} |",
            f"| **Generated** | {now_str} |",
            f"| **Total Sections** | {len(sections)} |",
            f"| **Total Words** | {total_words:,} |",
            f"| **Avg Words/Section** | {avg_words:,} |",
            f"| **Format** | Extended Markdown v2.0 |",
            f"| **Classification** | Confidential |",
            f"| **Status** | Final |",
            "",
            "---",
            "",
            f"<sub>*Generated by **{author}** · {now_str} · "
            f"Diamond Level Document Engine · All rights reserved.*</sub>",
            "",
        ]

        content = "\n".join(L)
        file_path.write_text(content, encoding="utf-8")

        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"Markdown generated successfully.\nTitle: {title}\n"
            f"Sections: {len(sections)}\nTotal words: {total_words:,}\n"
            f"File: {rel}\nSize: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_markdown failed: %s", exc, exc_info=True)
        return f"Markdown generation failed: {exc}"


def make_generate_markdown_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_markdown,
        name="generate_markdown",
        description=(
            "Generate a diamond-level Markdown: YAML front-matter, shield badges, ASCII stat cards, "
            "executive overview, emoji TOC with word counts, admonition callouts, ASCII progress bars, "
            "summary coverage table. Returns file path of .md."
        ),
        args_schema=GenerateMarkdownInput,
        return_direct=False,
    )


# =============================================================================
# 💎 5. DIAMOND HTML GENERATOR  (standalone, zero-dependency)
# =============================================================================


class GenerateHTMLInput(BaseModel):
    title: str
    sections: List[Dict[str, str]]
    subtitle: str = Field(default="")
    author: str = Field(default="SYNAPSE AI")
    user_id: str = Field(default="anonymous")


def _generate_html(
    title: str,
    sections: List[Dict[str, str]],
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    try:
        import base64
        import html as html_lib

        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.html"
        now_str = datetime.now(timezone.utc).strftime("%B %d, %Y")
        now_iso = datetime.now(timezone.utc).isoformat()

        total_words = sum(len(s.get("content", "").split()) for s in sections)

        def esc(t):
            return html_lib.escape(str(t))

        # ── Embedded chart PNGs as base64 ─────────────────────────
        sec_labels = [s.get("heading", "")[:20] for s in sections]
        sec_words = [len(s.get("content", "").split()) for s in sections]

        bar_b64 = pie_b64 = ""
        try:
            if any(w > 0 for w in sec_words):
                bar_b64 = _chart_b64(
                    "bar",
                    sec_labels,
                    sec_words,
                    "Content Distribution",
                    9,
                    4,
                    dark_bg=True,
                )
                pie_b64 = _chart_b64(
                    "pie",
                    sec_labels,
                    sec_words,
                    "Section Breakdown",
                    7,
                    5,
                    dark_bg=False,
                )
        except Exception as ce:
            logger.warning("HTML charts: %s", ce)

        # ── Colour token CSS vars ──────────────────────────────────
        css_vars = f"""
  --indigo:       {_Brand.INDIGO};
  --indigo-dark:  {_Brand.INDIGO_DARK};
  --indigo-mid:   {_Brand.INDIGO_MID};
  --indigo-light: {_Brand.INDIGO_LIGHT};
  --indigo-pale:  {_Brand.INDIGO_PALE};
  --indigo-bg:    {_Brand.INDIGO_BG};
  --violet:       {_Brand.VIOLET};
  --violet-light: {_Brand.VIOLET_LIGHT};
  --cyan:         {_Brand.CYAN};
  --emerald:      {_Brand.EMERALD};
  --amber:        {_Brand.AMBER};
  --rose:         {_Brand.ROSE};
  --dark:         {_Brand.DARK};
  --gray-900:     {_Brand.GRAY_900};
  --gray-700:     {_Brand.GRAY_700};
  --gray-500:     {_Brand.GRAY_500};
  --gray-300:     {_Brand.GRAY_300};
  --gray-100:     {_Brand.GRAY_100};"""

        # ── Sidebar nav items ──────────────────────────────────────
        ICONS = [
            "🔍",
            "📌",
            "💡",
            "📊",
            "⚡",
            "🎯",
            "🔬",
            "🌐",
            "🚀",
            "✅",
            "🧠",
            "🔑",
            "📈",
            "🏆",
            "🔧",
        ]
        nav_items = ""
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            icon = ICONS[(i - 1) % len(ICONS)]
            sid = f"sec-{i}"
            nav_items += f"""
        <a href="#{sid}" class="nav-item" onclick="setActive(this)">
          <span class="nav-icon">{icon}</span>
          <span class="nav-text">{esc(h)}</span>
        </a>"""

        # ── Stat cards ────────────────────────────────────────────
        stat_cards_html = ""
        stats = [
            (str(len(sections)), "Chapters", _Brand.INDIGO, "📚"),
            (f"{total_words:,}", "Total Words", _Brand.VIOLET, "📝"),
            (now_str.split()[0], "Month", _Brand.CYAN, "📅"),
            ("AI", "Powered", _Brand.EMERALD, "🤖"),
        ]
        for val, lbl, col, icon in stats:
            stat_cards_html += f"""
      <div class="stat-card" style="border-top:4px solid {col}">
        <div class="stat-icon">{icon}</div>
        <div class="stat-value" style="color:{col}">{esc(val)}</div>
        <div class="stat-label">{esc(lbl)}</div>
      </div>"""

        # ── Section content ───────────────────────────────────────
        sections_html = ""
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            content = sec.get("content", "")
            icon = ICONS[(i - 1) % len(ICONS)]
            sid = f"sec-{i}"
            wc = len(content.split())
            paras = [p.strip() for p in content.split("\n\n") if p.strip()]

            paras_html = ""
            for pi, para in enumerate(paras):
                safe_para = esc(para)
                if pi == 0:
                    paras_html += f'<p class="lead-para">{safe_para}</p>'
                else:
                    paras_html += f"<p>{safe_para}</p>"
                # Callout every 3rd paragraph
                if pi > 0 and pi % 3 == 0:
                    first_s = para.split(".")[0].strip()
                    if len(first_s) > 25:
                        paras_html += f"""
          <div class="callout">
            <div class="callout-icon">💡</div>
            <div class="callout-content">
              <strong>Key Insight:</strong> {esc(first_s)}.
            </div>
          </div>"""

            sections_html += f"""
      <section id="{sid}" class="doc-section">
        <div class="section-header">
          <div class="section-badge">{i:02d}</div>
          <div class="section-title-block">
            <h2>{icon} {esc(h)}</h2>
            <span class="section-meta">Section {i} of {len(sections)} &nbsp;·&nbsp; {wc:,} words</span>
          </div>
        </div>
        <div class="section-body">
          {paras_html}
        </div>
      </section>"""

        # ── Chart section ─────────────────────────────────────────
        charts_html = ""
        if bar_b64:
            charts_html += f"""
      <section id="charts" class="doc-section charts-section">
        <div class="section-header">
          <div class="section-badge" style="background:var(--violet)">📊</div>
          <div class="section-title-block">
            <h2>📊 Data Visualisation</h2>
            <span class="section-meta">Content analytics &amp; distribution</span>
          </div>
        </div>
        <div class="charts-grid">
          <div class="chart-card">
            <h3>Content Volume by Section</h3>
            <img src="data:image/png;base64,{bar_b64}" alt="Bar chart" />
          </div>
          <div class="chart-card chart-card--light">
            <h3>Section Distribution</h3>
            <img src="data:image/png;base64,{pie_b64}" alt="Pie chart" />
          </div>
        </div>
      </section>"""

        # ── TOC ───────────────────────────────────────────────────
        toc_html = ""
        for i, sec in enumerate(sections, 1):
            h = sec.get("heading", f"Section {i}")
            icon = ICONS[(i - 1) % len(ICONS)]
            wc = len(sec.get("content", "").split())
            sid = f"sec-{i}"
            toc_html += f"""
          <a href="#{sid}" class="toc-item">
            <span class="toc-num">{i:02d}</span>
            <span class="toc-icon">{icon}</span>
            <span class="toc-title">{esc(h)}</span>
            <span class="toc-words">{wc:,} w</span>
          </a>"""

        # ══════════════════════════════════════════════════════════
        # FULL HTML DOCUMENT
        # ══════════════════════════════════════════════════════════
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8"/>
  <meta name="viewport" content="width=device-width,initial-scale=1.0"/>
  <meta name="generator" content="{esc(author)}"/>
  <meta name="created" content="{now_iso}"/>
  <title>{esc(title)}</title>
  <style>
    :root {{{css_vars}
    }}

    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}

    body {{
      font-family: 'Segoe UI', system-ui, -apple-system, sans-serif;
      background: var(--gray-100);
      color: var(--gray-900);
      display: flex;
      min-height: 100vh;
      line-height: 1.7;
    }}

    /* ── SIDEBAR ── */
    .sidebar {{
      width: 280px;
      min-width: 280px;
      background: var(--dark);
      height: 100vh;
      position: sticky;
      top: 0;
      display: flex;
      flex-direction: column;
      overflow-y: auto;
      z-index: 100;
    }}
    .sidebar-brand {{
      padding: 28px 20px 20px;
      background: linear-gradient(135deg, var(--indigo-dark), var(--violet));
      border-bottom: 1px solid rgba(255,255,255,0.08);
    }}
    .sidebar-brand .brand-logo {{
      font-size: 11px; font-weight: 700; letter-spacing: 2px;
      color: var(--indigo-pale); text-transform: uppercase; margin-bottom: 8px;
    }}
    .sidebar-brand .brand-title {{
      font-size: 15px; font-weight: 700; color: #fff;
      line-height: 1.3;
      display: -webkit-box; -webkit-line-clamp: 3;
      -webkit-box-orient: vertical; overflow: hidden;
    }}
    .sidebar-brand .brand-meta {{
      margin-top: 8px; font-size: 11px; color: var(--indigo-pale); opacity: 0.75;
    }}
    .sidebar-nav {{ padding: 12px 0; flex: 1; }}
    .nav-section-label {{
      padding: 16px 20px 6px;
      font-size: 10px; font-weight: 700; letter-spacing: 1.5px;
      color: var(--indigo-light); text-transform: uppercase;
    }}
    .nav-item {{
      display: flex; align-items: flex-start; gap: 10px;
      padding: 10px 20px; text-decoration: none;
      color: rgba(255,255,255,0.65); font-size: 13px;
      transition: all 0.2s; border-left: 3px solid transparent;
      line-height: 1.4;
    }}
    .nav-item:hover {{
      background: rgba(255,255,255,0.06);
      color: #fff; border-left-color: var(--indigo-light);
    }}
    .nav-item.active {{
      background: rgba(79,70,229,0.25);
      color: #fff; border-left-color: var(--indigo-mid);
    }}
    .nav-icon {{ font-size: 15px; flex-shrink: 0; margin-top: 1px; }}
    .sidebar-footer {{
      padding: 16px 20px;
      border-top: 1px solid rgba(255,255,255,0.08);
      font-size: 10px; color: rgba(255,255,255,0.35);
    }}

    /* ── MAIN ── */
    .main {{ flex: 1; display: flex; flex-direction: column; min-width: 0; }}

    /* ── HERO ── */
    .hero {{
      background: linear-gradient(135deg, var(--dark) 0%, var(--indigo-dark) 50%, var(--violet) 100%);
      padding: 60px 56px 48px;
      position: relative; overflow: hidden;
    }}
    .hero::before {{
      content: ""; position: absolute;
      width: 500px; height: 500px; border-radius: 50%;
      background: rgba(255,255,255,0.04);
      right: -100px; top: -150px;
    }}
    .hero::after {{
      content: ""; position: absolute;
      width: 300px; height: 300px; border-radius: 50%;
      background: rgba(99,102,241,0.15);
      right: 180px; top: 30px;
    }}
    .hero-eyebrow {{
      font-size: 11px; font-weight: 700; letter-spacing: 3px;
      color: var(--indigo-pale); text-transform: uppercase;
      margin-bottom: 16px; position: relative; z-index: 1;
    }}
    .hero h1 {{
      font-size: clamp(28px, 4vw, 48px); font-weight: 800;
      color: #fff; line-height: 1.15; max-width: 700px;
      position: relative; z-index: 1; margin-bottom: 16px;
    }}
    .hero-sub {{
      font-size: 17px; color: var(--indigo-pale);
      max-width: 560px; position: relative; z-index: 1;
    }}
    .hero-divider {{
      width: 60px; height: 4px;
      background: linear-gradient(90deg, var(--indigo-light), var(--violet-light));
      border-radius: 2px; margin: 24px 0 20px; position: relative; z-index: 1;
    }}
    .hero-meta {{ font-size: 12px; color: rgba(255,255,255,0.5); position: relative; z-index: 1; }}

    /* ── STAT CARDS ── */
    .stats-bar {{ background: #fff; padding: 28px 56px; display: flex; gap: 20px; flex-wrap: wrap;
      border-bottom: 1px solid var(--gray-300); }}
    .stat-card {{
      flex: 1; min-width: 140px; padding: 18px 20px; border-radius: 10px;
      background: var(--gray-100); text-align: center;
    }}
    .stat-icon {{ font-size: 22px; margin-bottom: 6px; }}
    .stat-value {{ font-size: 26px; font-weight: 800; margin-bottom: 4px; }}
    .stat-label {{ font-size: 10px; font-weight: 700; letter-spacing: 1.5px;
      text-transform: uppercase; color: var(--gray-500); }}

    /* ── TOC ── */
    .toc-section {{ background: #fff; padding: 36px 56px; border-bottom: 1px solid var(--gray-300); }}
    .toc-section h2 {{ font-size: 14px; font-weight: 700; letter-spacing: 2px;
      text-transform: uppercase; color: var(--indigo); margin-bottom: 16px; }}
    .toc-grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(300px, 1fr)); gap: 8px; }}
    .toc-item {{
      display: flex; align-items: center; gap: 10px; padding: 10px 14px;
      text-decoration: none; border-radius: 8px; transition: all 0.2s;
      background: var(--indigo-bg); color: var(--dark);
    }}
    .toc-item:hover {{ background: var(--indigo-pale); transform: translateX(3px); }}
    .toc-num {{ font-size: 11px; font-weight: 800; color: var(--indigo);
      min-width: 26px; }}
    .toc-icon {{ font-size: 16px; }}
    .toc-title {{ flex: 1; font-size: 13.5px; font-weight: 500; }}
    .toc-words {{ font-size: 11px; color: var(--gray-500); white-space: nowrap; }}

    /* ── CONTENT ── */
    .content-area {{ padding: 0 56px 56px; }}

    .doc-section {{
      margin-top: 48px; background: #fff;
      border-radius: 12px; overflow: hidden;
      box-shadow: 0 1px 3px rgba(0,0,0,0.06), 0 4px 16px rgba(0,0,0,0.04);
    }}
    .section-header {{
      display: flex; align-items: stretch;
      background: linear-gradient(90deg, var(--indigo-dark), var(--indigo));
      min-height: 70px;
    }}
    .section-badge {{
      display: flex; align-items: center; justify-content: center;
      min-width: 70px; font-size: 22px; font-weight: 800; color: #fff;
      background: rgba(0,0,0,0.2); flex-shrink: 0;
    }}
    .section-title-block {{ padding: 16px 24px; }}
    .section-title-block h2 {{ font-size: 20px; font-weight: 700; color: #fff;
      margin-bottom: 4px; }}
    .section-meta {{ font-size: 12px; color: var(--indigo-pale); }}
    .section-body {{ padding: 28px 32px; }}
    .section-body p {{ margin-bottom: 16px; font-size: 15px; line-height: 1.8;
      color: var(--gray-700); }}
    .lead-para {{
      font-size: 16.5px !important; font-weight: 600 !important;
      color: var(--gray-900) !important;
      border-left: 4px solid var(--indigo-pale);
      padding-left: 16px; margin-bottom: 20px !important;
    }}

    /* ── CALLOUT ── */
    .callout {{
      display: flex; gap: 12px; align-items: flex-start;
      background: var(--indigo-bg); border-radius: 8px;
      border-left: 4px solid var(--indigo);
      padding: 14px 18px; margin: 12px 0 20px;
    }}
    .callout-icon {{ font-size: 18px; flex-shrink: 0; margin-top: 2px; }}
    .callout-content {{ font-size: 14px; line-height: 1.6; color: var(--dark); }}

    /* ── CHARTS ── */
    .charts-section {{ background: var(--dark) !important; }}
    .charts-section .section-header {{
      background: linear-gradient(90deg, #0f0c2e, var(--violet)) !important;
    }}
    .charts-section .section-body {{ background: var(--dark); padding: 32px; }}
    .charts-grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 24px; }}
    .chart-card {{
      background: rgba(255,255,255,0.05); border-radius: 10px;
      padding: 20px; text-align: center;
      border: 1px solid rgba(255,255,255,0.08);
    }}
    .chart-card--light {{ background: rgba(255,255,255,0.92); }}
    .chart-card h3 {{ font-size: 13px; font-weight: 700; letter-spacing: 1px;
      text-transform: uppercase; color: var(--indigo-pale); margin-bottom: 16px; }}
    .chart-card--light h3 {{ color: var(--indigo-dark); }}
    .chart-card img {{ max-width: 100%; border-radius: 6px; }}

    /* ── FOOTER ── */
    .doc-footer {{
      background: var(--dark); color: rgba(255,255,255,0.45);
      padding: 32px 56px; margin-top: 48px;
      display: flex; justify-content: space-between; align-items: center;
      flex-wrap: wrap; gap: 12px; font-size: 12px;
    }}
    .footer-brand {{ color: var(--indigo-light); font-weight: 700; }}

    /* ── SCROLL PROGRESS ── */
    #progress-bar {{
      position: fixed; top: 0; left: 0; height: 3px; width: 0%;
      background: linear-gradient(90deg, var(--indigo), var(--violet));
      z-index: 9999; transition: width 0.1s;
    }}

    /* ── PRINT ── */
    @media print {{
      .sidebar {{ display: none; }}
      .main {{ margin: 0; }}
      .doc-section {{ page-break-inside: avoid; box-shadow: none; }}
    }}

    @media (max-width: 768px) {{
      .sidebar {{ display: none; }}
      .hero, .stats-bar, .toc-section, .content-area {{ padding-left: 20px; padding-right: 20px; }}
      .charts-grid {{ grid-template-columns: 1fr; }}
    }}

    /* ── THEME SWITCHER ── */
    .theme-toggle {{
      position: fixed; bottom: 24px; right: 24px; z-index: 9999;
      width: 48px; height: 48px; border-radius: 50%;
      background: linear-gradient(135deg, var(--indigo), var(--violet));
      border: none; cursor: pointer; box-shadow: 0 4px 16px rgba(79,70,229,0.4);
      display: flex; align-items: center; justify-content: center;
      font-size: 20px; transition: transform 0.2s, box-shadow 0.2s;
      color: white;
    }}
    .theme-toggle:hover {{ transform: scale(1.1); box-shadow: 0 6px 24px rgba(79,70,229,0.6); }}

    /* ── DARK MODE ── */
    body.dark-mode {{
      background: #0f0c2e;
      color: #e2e8f0;
    }}
    body.dark-mode .main {{ background: #0f0c2e; }}
    body.dark-mode .stats-bar {{ background: #1a1740; border-color: #2d2a6e; }}
    body.dark-mode .stat-card {{ background: #1e1b4b; }}
    body.dark-mode .stat-value {{ filter: brightness(1.2); }}
    body.dark-mode .toc-section {{ background: #1a1740; border-color: #2d2a6e; }}
    body.dark-mode .toc-item {{ background: #1e1b4b; color: #c7d2fe; }}
    body.dark-mode .toc-item:hover {{ background: #2d2a6e; }}
    body.dark-mode .toc-words {{ color: #818cf8; }}
    body.dark-mode .content-area {{ background: #0f0c2e; }}
    body.dark-mode .doc-section {{ background: #1a1740; box-shadow: 0 2px 16px rgba(0,0,0,0.4); }}
    body.dark-mode .section-body p {{ color: #cbd5e1; }}
    body.dark-mode .lead-para {{ color: #e2e8f0 !important; border-color: #4f46e5; }}
    body.dark-mode .callout {{ background: #1e1b4b; border-color: #6366f1; }}
    body.dark-mode .callout-content {{ color: #c7d2fe; }}
    body.dark-mode .section-meta {{ color: #818cf8; }}
    body.dark-mode .chart-card {{ background: rgba(255,255,255,0.04); border-color: rgba(255,255,255,0.1); }}
  </style>
</head>
<body>

<div id="progress-bar"></div>
<button class="theme-toggle" onclick="toggleTheme()" title="Toggle dark/light mode" id="theme-btn">🌙</button>

<!-- SIDEBAR -->
<aside class="sidebar">
  <div class="sidebar-brand">
    <div class="brand-logo">SYNAPSE AI</div>
    <div class="brand-title">{esc(title)}</div>
    <div class="brand-meta">{esc(author)} · {esc(now_str)}</div>
  </div>
  <nav class="sidebar-nav">
    <div class="nav-section-label">Contents</div>
    {nav_items}
    {"<div class='nav-section-label' style='margin-top:12px'>Analytics</div><a href='#charts' class='nav-item' onclick='setActive(this)'><span class='nav-icon'>📊</span><span class='nav-text'>Data Visualisation</span></a>" if bar_b64 else ""}
  </nav>
  <div class="sidebar-footer">
    Generated by {esc(author)}<br/>
    {esc(now_str)} · Diamond Level
  </div>
</aside>

<!-- MAIN -->
<main class="main">

  <!-- HERO -->
  <header class="hero">
    <div class="hero-eyebrow">{esc(author)} &nbsp;·&nbsp; Executive Report</div>
    <h1>{esc(title)}</h1>
    {"<p class='hero-sub'>" + esc(subtitle) + "</p>" if subtitle else ""}
    <div class="hero-divider"></div>
    <p class="hero-meta">{esc(now_str)} &nbsp;·&nbsp; {len(sections)} chapters &nbsp;·&nbsp; {total_words:,} words</p>
  </header>

  <!-- STAT CARDS -->
  <div class="stats-bar">
    {stat_cards_html}
  </div>

  <!-- TOC -->
  <div class="toc-section">
    <h2>📚 Table of Contents</h2>
    <div class="toc-grid">
      {toc_html}
    </div>
  </div>

  <!-- CONTENT -->
  <div class="content-area">
    {sections_html}
    {charts_html}
  </div>

  <!-- FOOTER -->
  <footer class="doc-footer">
    <span>
      <span class="footer-brand">{esc(author)}</span> &nbsp;·&nbsp;
      {esc(title)}
    </span>
    <span>{esc(now_str)} &nbsp;·&nbsp; Diamond Level &nbsp;·&nbsp; All rights reserved</span>
  </footer>

</main>

<script>
  // Theme toggle
  function toggleTheme() {{
    const body = document.body;
    const btn  = document.getElementById('theme-btn');
    body.classList.toggle('dark-mode');
    btn.textContent = body.classList.contains('dark-mode') ? '☀️' : '🌙';
    localStorage.setItem('synapse-theme', body.classList.contains('dark-mode') ? 'dark' : 'light');
  }}
  // Restore saved theme
  (function() {{
    if (localStorage.getItem('synapse-theme') === 'dark') {{
      document.body.classList.add('dark-mode');
      const btn = document.getElementById('theme-btn');
      if (btn) btn.textContent = '☀️';
    }}
  }})();

  // Scroll progress bar
  window.addEventListener('scroll', () => {{
    const el = document.getElementById('progress-bar');
    const pct = window.scrollY / (document.body.scrollHeight - window.innerHeight) * 100;
    if (el) el.style.width = Math.min(100, pct) + '%';
  }});

  // Active nav item
  function setActive(el) {{
    document.querySelectorAll('.nav-item').forEach(n => n.classList.remove('active'));
    el.classList.add('active');
  }}

  // Highlight nav on scroll
  const sections = document.querySelectorAll('.doc-section');
  const navItems = document.querySelectorAll('.nav-item');
  const observer = new IntersectionObserver(entries => {{
    entries.forEach(e => {{
      if (e.isIntersecting) {{
        const id = e.target.id;
        navItems.forEach(n => {{
          n.classList.toggle('active', n.getAttribute('href') === '#' + id);
        }});
      }}
    }});
  }}, {{ rootMargin: '-30% 0px -60% 0px' }});
  sections.forEach(s => observer.observe(s));
</script>

</body>
</html>"""

        file_path.write_text(html, encoding="utf-8")
        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"HTML document generated successfully.\nTitle: {title}\n"
            f"Sections: {len(sections)}\nFile: {rel}\n"
            f"Size: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_html failed: %s", exc, exc_info=True)
        return f"HTML generation failed: {exc}"


def _generate_html_page_from_prompt(
    title: str,
    prompt: str,
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
    openrouter_key: str = "",
    gemini_key: str = "",
    model_override: str = "",
) -> str:
    """
    Generate a real, fully self-contained HTML/CSS/JS web page based on the user's
    prompt (portfolio, landing page, e-commerce, etc.).

    When an LLM key is available the LLM produces the actual HTML code directly.
    Falls back to a clean, user-described placeholder page if no key is configured.
    Returns the same result string as _generate_html (with "Path: …" line).
    """
    import os
    import re

    try:
        file_path = _doc_dir(user_id) / f"{uuid.uuid4().hex}.html"

        html_content = None

        # ── 1. LLM-generated HTML ────────────────────────────────────────────
        system_prompt = (
            "You are an expert front-end developer. Write a COMPLETE, fully self-contained HTML file.\n\n"
            "RULES (non-negotiable):\n"
            "1. Output ONLY raw HTML — no markdown fences, no explanation, no preamble.\n"
            "2. Start with <!DOCTYPE html> and end with </html>.\n"
            "3. All CSS in a <style> tag in <head>. All JS in <script> before </body>.\n"
            "4. Only allowed external: Google Fonts CDN, cdn.jsdelivr.net (Alpine.js / Chart.js).\n\n"
            "DESIGN:\n"
            "- Match the user's page type, style and colour scheme exactly.\n"
            "- CSS custom properties for colours. Mobile-first grid/flex layout.\n"
            "- Smooth transitions, keyframe animations, hover effects, sticky nav with blur.\n"
            "- Real, meaningful placeholder content — no lorem ipsum.\n"
            "- Semantic HTML5 (<header><nav><main><section><footer>).\n"
            "- picsum.photos / placehold.co for images.\n\n"
            "CRITICAL: Return ONLY the raw HTML starting with <!DOCTYPE html>."
        )

        user_msg = (
            f"Title: {title} | Type: {subtitle or 'Web Page'}\n\n"
            f"Requirements:\n{prompt[:1500]}\n\n"
            "Generate the complete HTML file. Start with <!DOCTYPE html>."
        )

        raw_html = None
        try:
            if openrouter_key:
                from langchain_core.messages import HumanMessage, SystemMessage
                from langchain_openai import ChatOpenAI

                # Choose model — prefer user override, else affordable default
                model = model_override or os.environ.get(
                    "OPENROUTER_MODEL", "openai/gpt-4o-mini"
                )

                # Use a conservative token limit to keep costs low.
                # Good HTML pages fit well within 2000 output tokens.
                # Users can raise this via HTML_GEN_MAX_TOKENS env var.
                max_tok = int(os.environ.get("HTML_GEN_MAX_TOKENS", "2000"))

                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=openrouter_key,
                    openai_api_base=os.environ.get(
                        "OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"
                    ),
                    temperature=0.7,
                    max_tokens=max_tok,
                )
                try:
                    resp = llm.invoke(
                        [
                            SystemMessage(content=system_prompt),
                            HumanMessage(content=user_msg),
                        ]
                    )
                    raw_html = (
                        resp.content
                        if isinstance(resp.content, str)
                        else str(resp.content)
                    )
                    logger.info(
                        "HTML page: OpenRouter generated %d chars", len(raw_html)
                    )
                except Exception as primary_exc:
                    logger.warning(
                        "HTML LLM primary call failed (%s): %s", model, primary_exc
                    )
                    # If the chosen model is expensive and fails with 402, auto-retry with gpt-4o-mini
                    if "402" in str(primary_exc) and model != "openai/gpt-4o-mini":
                        logger.info(
                            "HTML page: retrying with gpt-4o-mini (insufficient credits for %s)",
                            model,
                        )
                        llm_fallback = ChatOpenAI(
                            model="openai/gpt-4o-mini",
                            openai_api_key=openrouter_key,
                            openai_api_base=os.environ.get(
                                "OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"
                            ),
                            temperature=0.7,
                            max_tokens=4096,
                        )
                        try:
                            resp = llm_fallback.invoke(
                                [
                                    SystemMessage(content=system_prompt),
                                    HumanMessage(content=user_msg),
                                ]
                            )
                            raw_html = (
                                resp.content
                                if isinstance(resp.content, str)
                                else str(resp.content)
                            )
                            logger.info(
                                "HTML page: gpt-4o-mini fallback generated %d chars",
                                len(raw_html),
                            )
                        except Exception as fallback_exc:
                            logger.warning(
                                "HTML LLM fallback also failed: %s", fallback_exc
                            )
                            raw_html = None
                    else:
                        raise  # re-raise so outer handler catches it properly

            elif gemini_key:
                from langchain_core.messages import HumanMessage, SystemMessage
                from langchain_google_genai import ChatGoogleGenerativeAI

                llm = ChatGoogleGenerativeAI(
                    model=os.environ.get("GEMINI_MODEL", "gemini-1.5-flash-latest"),
                    google_api_key=gemini_key,
                    temperature=0.7,
                    max_output_tokens=4096,
                    convert_system_message_to_human=True,
                )
                resp = llm.invoke(
                    [
                        SystemMessage(content=system_prompt),
                        HumanMessage(content=user_msg),
                    ]
                )
                raw_html = (
                    resp.content if isinstance(resp.content, str) else str(resp.content)
                )
                logger.info("HTML page: Gemini generated %d chars", len(raw_html))

        except Exception as llm_exc:
            logger.warning("HTML LLM generation failed: %s — using fallback", llm_exc)
            raw_html = None

        if raw_html:
            # Strip any markdown fences the LLM may have added
            cleaned = re.sub(r"```html\s*", "", raw_html, flags=re.IGNORECASE)
            cleaned = re.sub(r"```\s*$", "", cleaned).strip()
            # Verify it looks like HTML
            if "<!DOCTYPE html>" in cleaned.upper() or "<html" in cleaned.lower():
                html_content = cleaned

        # ── 2. Fallback: clean placeholder page ──────────────────────────────
        if not html_content:
            logger.info("HTML page: using fallback placeholder for '%s'", title)
            now_str = datetime.now(timezone.utc).strftime("%B %d, %Y")
            # Derive a style hint from the prompt
            prompt_lower = prompt.lower()
            if any(
                k in prompt_lower for k in ["dark", "luxury", "neon", "cyber", "black"]
            ):
                bg, text, accent = "#0f0f0f", "#f5f5f5", "#6366f1"
            elif any(
                k in prompt_lower
                for k in ["pastel", "soft", "minimal", "clean", "white"]
            ):
                bg, text, accent = "#fafafa", "#1a1a2e", "#6366f1"
            else:
                bg, text, accent = "#0f172a", "#f8fafc", "#6366f1"

            subtitle_display = subtitle or "Web Page"
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title}</title>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet" />
  <style>
    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
    :root {{
      --accent: {accent};
      --bg: {bg};
      --text: {text};
    }}
    body {{
      font-family: 'Inter', sans-serif;
      background: var(--bg);
      color: var(--text);
      min-height: 100vh;
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      padding: 40px 20px;
    }}
    .container {{
      max-width: 700px;
      width: 100%;
      text-align: center;
    }}
    .badge {{
      display: inline-block;
      background: var(--accent);
      color: #fff;
      font-size: 11px;
      font-weight: 700;
      letter-spacing: 2px;
      text-transform: uppercase;
      padding: 6px 16px;
      border-radius: 999px;
      margin-bottom: 28px;
    }}
    h1 {{
      font-size: clamp(2rem, 6vw, 3.5rem);
      font-weight: 800;
      line-height: 1.15;
      margin-bottom: 20px;
      background: linear-gradient(135deg, var(--text), var(--accent));
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      background-clip: text;
    }}
    .subtitle {{
      font-size: 1.1rem;
      opacity: 0.65;
      margin-bottom: 40px;
      line-height: 1.6;
    }}
    .prompt-box {{
      background: rgba(99,102,241,0.08);
      border: 1px solid rgba(99,102,241,0.25);
      border-radius: 16px;
      padding: 28px;
      margin-bottom: 40px;
      text-align: left;
      font-size: 0.95rem;
      line-height: 1.7;
      opacity: 0.85;
    }}
    .prompt-box h3 {{
      font-size: 0.75rem;
      letter-spacing: 1.5px;
      text-transform: uppercase;
      color: var(--accent);
      margin-bottom: 12px;
      font-weight: 700;
    }}
    .cta-btn {{
      display: inline-flex;
      align-items: center;
      gap: 8px;
      background: var(--accent);
      color: #fff;
      font-size: 1rem;
      font-weight: 600;
      padding: 14px 32px;
      border-radius: 12px;
      border: none;
      cursor: pointer;
      transition: transform 0.2s, box-shadow 0.2s;
      text-decoration: none;
    }}
    .cta-btn:hover {{
      transform: translateY(-2px);
      box-shadow: 0 12px 32px rgba(99,102,241,0.35);
    }}
    .note {{
      margin-top: 48px;
      font-size: 0.8rem;
      opacity: 0.4;
    }}
  </style>
</head>
<body>
  <div class="container">
    <div class="badge">{subtitle_display}</div>
    <h1>{title}</h1>
    <p class="subtitle">Your AI-generated page is ready. Add your OpenRouter or Gemini API key in Settings to get a fully custom-designed HTML page from your prompt.</p>
    <div class="prompt-box">
      <h3>Your Request</h3>
      <p>{prompt[:500]}{'…' if len(prompt) > 500 else ''}</p>
    </div>
    <a class="cta-btn" href="#">Get Started →</a>
    <p class="note">Generated by SYNAPSE AI · {now_str}</p>
  </div>
  <script>
    document.querySelector('.cta-btn').addEventListener('click', e => {{
      e.preventDefault();
      alert('Configure your AI API key in SYNAPSE Settings to generate a fully custom page!');
    }});
  </script>
</body>
</html>"""

        file_path.write_text(html_content, encoding="utf-8")
        sz = file_path.stat().st_size
        rel = _rel(file_path)
        return (
            f"HTML page generated successfully.\nTitle: {title}\n"
            f"File: {rel}\nSize: {sz:,} bytes\nPath: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("_generate_html_page_from_prompt failed: %s", exc, exc_info=True)
        return f"HTML generation failed: {exc}"


def make_generate_html_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_html,
        name="generate_html",
        description=(
            "Generate a diamond-level standalone HTML5 document with glassmorphism sidebar nav, "
            "gradient hero header, stat cards, emoji TOC grid, section cards with callouts, "
            "embedded Matplotlib bar/pie charts as base64, scroll progress bar, "
            "and a branded footer. Zero external dependencies — fully self-contained. "
            "Returns file path of .html."
        ),
        args_schema=GenerateHTMLInput,
        return_direct=False,
    )
