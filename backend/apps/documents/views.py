"""
backend.apps.documents.views
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
REST API views for Document Studio.

Endpoints (Phase 5.2 + 5.3):
  POST  /api/v1/documents/generate/          — generate a document via agent tools
  POST  /api/v1/documents/generate-project/  — generate a project scaffold (.zip)
  GET   /api/v1/documents/                   — list user's documents
  GET   /api/v1/documents/{id}/              — retrieve document metadata
  GET   /api/v1/documents/{id}/download/     — stream the file for download
  DELETE /api/v1/documents/{id}/             — delete a document

Phase 5.2 — Document Generation (Week 14)
Phase 5.3 — Project Builder (Week 15)
"""

from __future__ import annotations

import logging
import mimetypes
import os
import sys
from datetime import datetime, timezone
from pathlib import Path

from apps.core.pagination import StandardPagination

from django.conf import settings
from django.http import FileResponse, Http404
from rest_framework import status
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.request import Request
from rest_framework.response import Response
from rest_framework.views import APIView

from .models import GeneratedDocument
from .serializers import (
    DocumentGenerateSerializer,
    GeneratedDocumentListSerializer,
    GeneratedDocumentSerializer,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# QA-16: Document LLM Builder — extracted from DocumentGenerateView
# ---------------------------------------------------------------------------


class _DocumentLLMBuilder:
    """
    Handles LLM provider selection and raw text invocation for document generation.

    Extracted from DocumentGenerateView._expand_prompt_to_sections() to:
      1. Separate concerns (prompt construction vs. LLM invocation)
      2. Make the LLM call independently testable
      3. Eliminate 45 lines of duplicated inline LLM invocation code in views.py
    """

    @staticmethod
    def invoke(
        system: str,
        user_msg: str,
        openrouter_key: str = "",
        gemini_key: str = "",
        model_override: str = "",
    ) -> str | None:
        """
        Invoke the LLM and return the raw string response, or None on failure.

        Provider selection (in order of preference):
          1. Vercel AI Gateway (AI_GATEWAY_API_KEY)
          2. Groq             (GROQ_API_KEY)
          3. OpenRouter       (passed-in user/env key)
          4. Gemini direct    (passed-in user/env key)
        """
        import os

        try:
            from langchain_core.messages import HumanMessage, SystemMessage
        except ImportError:
            logger.error("langchain-core not installed — cannot generate LLM sections")
            return None

        try:
            from langchain_openai import ChatOpenAI
        except ImportError:
            ChatOpenAI = None  # type: ignore

        def _try_openai_compat(
            label: str, api_key: str, base_url: str, model: str
        ) -> str | None:
            """Helper: try one OpenAI-compatible provider, return raw text or None."""
            if not (api_key and ChatOpenAI):
                return None
            try:
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    openai_api_base=base_url,
                    temperature=0.7,
                    max_tokens=16000,
                )
                resp = llm.invoke(
                    [SystemMessage(content=system), HumanMessage(content=user_msg)]
                )
                raw = (
                    resp.content if isinstance(resp.content, str) else str(resp.content)
                )
                logger.info(
                    "_DocumentLLMBuilder: %s response (%d chars)", label, len(raw)
                )
                return raw
            except Exception as exc:
                logger.warning("_DocumentLLMBuilder: %s failed (%s)", label, exc)
                return None

        # ── 1. Vercel AI Gateway ──────────────────────────────────────────────
        gateway_key = (os.environ.get("AI_GATEWAY_API_KEY") or "").strip()
        if gateway_key and not gateway_key.startswith("your-"):
            raw = _try_openai_compat(
                "AI Gateway",
                gateway_key,
                "https://ai-gateway.vercel.sh/v1",
                model_override
                or os.environ.get("AI_GATEWAY_MODEL", "openai/gpt-4o-mini"),
            )
            if raw:
                return raw

        # ── 2. Groq (fast inference) ──────────────────────────────────────────
        groq_key = (os.environ.get("GROQ_API_KEY") or "").strip()
        if groq_key and not groq_key.startswith("your-"):
            raw = _try_openai_compat(
                "Groq",
                groq_key,
                "https://api.groq.com/openai/v1",
                model_override
                or os.environ.get("GROQ_MODEL", "llama-3.3-70b-versatile"),
            )
            if raw:
                return raw

        # ── 3. OpenRouter (passed-in or env key) ──────────────────────────────
        if openrouter_key:
            raw = _try_openai_compat(
                "OpenRouter",
                openrouter_key,
                os.environ.get("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"),
                model_override
                or os.environ.get("OPENROUTER_MODEL", "openai/gpt-4o-mini"),
            )
            if raw:
                return raw

        # ── Google Gemini (fallback) ──────────────────────────────────────────
        if gemini_key:
            try:
                from langchain_google_genai import ChatGoogleGenerativeAI

                llm = ChatGoogleGenerativeAI(
                    model=model_override
                    or os.environ.get("GEMINI_MODEL", "gemini-1.5-flash-latest"),
                    google_api_key=gemini_key,
                    temperature=0.7,
                    max_output_tokens=8192,
                    convert_system_message_to_human=True,
                )
                resp = llm.invoke(
                    [SystemMessage(content=system), HumanMessage(content=user_msg)]
                )
                raw = (
                    resp.content if isinstance(resp.content, str) else str(resp.content)
                )
                logger.info("_DocumentLLMBuilder: Gemini response (%d chars)", len(raw))
                return raw
            except Exception as exc:
                logger.warning("_DocumentLLMBuilder: Gemini failed (%s)", exc)

        # ── No provider available ─────────────────────────────────────────────
        if not (openrouter_key or gemini_key or gateway_key or groq_key):
            logger.warning(
                "_DocumentLLMBuilder: No LLM API key configured. "
                "Set AI_GATEWAY_API_KEY, GROQ_API_KEY, OPENROUTER_API_KEY, or "
                "GEMINI_API_KEY — or have the user save a key in Settings → AI Engine."
            )
        return None


# Ensure ai_engine is importable regardless of whether we're running inside
# Docker (PYTHONPATH=/app:/ai_engine_pkg) or locally (project root on path).
# In Docker: /ai_engine_pkg is on PYTHONPATH and ./ai_engine is mounted there.
# Locally: we add the project root (two levels up from backend/) so that
# `import ai_engine` resolves to <project_root>/ai_engine/.
def _ensure_ai_engine_on_path() -> None:
    try:
        import ai_engine  # noqa: F401 — already importable, nothing to do
    except ImportError:
        # Walk up from this file: views.py → documents/ → apps/ → backend/ → project_root
        project_root = Path(__file__).resolve().parents[3]
        if str(project_root) not in sys.path:
            sys.path.insert(0, str(project_root))
            logger.debug("ai_engine: added %s to sys.path", project_root)


_ensure_ai_engine_on_path()

# MIME type map for document types
_MIME_MAP = {
    "pdf": "application/pdf",
    "ppt": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "word": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "markdown": "text/markdown",
    "html": "text/html",
    "project": "application/zip",
}

_EXT_MAP = {
    "pdf": ".pdf",
    "ppt": ".pptx",
    "word": ".docx",
    "markdown": ".md",
    "html": ".html",
    "project": ".zip",
}


# ---------------------------------------------------------------------------
# Generate
# ---------------------------------------------------------------------------


class DocumentGenerateView(APIView):
    """
    POST /api/v1/documents/generate/

    Generates a document using the appropriate agent tool:
      - PDF   → generate_pdf tool
      - PPT   → generate_ppt tool
      - Word  → generate_word_doc tool
      - Markdown → generate_markdown tool

    If 'sections' is provided in the request, they are passed directly to the tool.
    Otherwise, a minimal single-section document is generated from the prompt.
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request) -> Response:
        serializer = DocumentGenerateSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        data = serializer.validated_data
        doc_type = data["doc_type"]
        title = data["title"]
        prompt = data["prompt"]
        sections = data.get("sections") or []
        subtitle = data.get("subtitle", "")
        author = data.get("author", "") or "SYNAPSE AI"  # coerce empty string → default
        user_id = str(request.user.id)
        content_types = data.get("content_types") or []
        source_item_ids = data.get("source_item_ids") or []

        # ── HTML pages: generate the actual HTML/CSS/JS directly from the prompt ──
        if doc_type == "html":
            openrouter_key, gemini_key = self._get_llm_keys(request.user)
            try:
                from ai_engine.agents.doc_tools import _generate_html_page_from_prompt

                result_str = _generate_html_page_from_prompt(
                    title=title,
                    prompt=prompt,
                    subtitle=subtitle,
                    author=author,
                    user_id=user_id,
                    openrouter_key=openrouter_key,
                    gemini_key=gemini_key,
                    model_override=data.get("model", ""),
                )
            except Exception as exc:
                logger.error("HTML page generation failed: %s", exc)
                return Response(
                    {"error": f"HTML page generation failed: {exc}"},
                    status=status.HTTP_500_INTERNAL_SERVER_ERROR,
                )
            file_path_str = ""
            for line in result_str.splitlines():
                if line.startswith("Path:"):
                    file_path_str = line.replace("Path:", "").strip()
                    break
            sources_used = []
        else:
            # Build sections from prompt if not provided — use LLM + RAG context
            if not sections:
                sections, sources_used = self._expand_prompt_to_sections(
                    prompt,
                    title,
                    doc_type,
                    user=request.user,
                    model_override=data.get("model", ""),
                    content_types=content_types,
                    source_item_ids=source_item_ids,
                )
            else:
                sources_used = []

            try:
                result_str, file_path_str = self._call_tool(
                    doc_type=doc_type,
                    title=title,
                    sections=sections,
                    subtitle=subtitle,
                    author=author,
                    user_id=user_id,
                )
            except Exception as exc:
                logger.error("Document generation failed: %s", exc)
                return Response(
                    {"error": f"Document generation failed: {exc}"},
                    status=status.HTTP_500_INTERNAL_SERVER_ERROR,
                )

        if "failed" in result_str.lower():
            return Response(
                {"error": result_str},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

        # Parse file path from result string
        abs_path = Path(file_path_str) if file_path_str else None
        rel_path = ""
        file_size = 0
        if abs_path and abs_path.exists():
            media_root = Path(settings.MEDIA_ROOT)
            try:
                rel_path = str(abs_path.relative_to(media_root))
            except ValueError:
                rel_path = str(abs_path)
            file_size = abs_path.stat().st_size

        # Persist record
        doc_obj = GeneratedDocument.objects.create(
            user=request.user,
            title=title,
            doc_type=doc_type,
            file_path=rel_path,
            file_size_bytes=file_size,
            agent_prompt=prompt,
            sources_metadata=sources_used,
            metadata={
                "subtitle": subtitle,
                "author": author,
                "section_count": len(sections),
                "tool_result": result_str[:500],
                "rag_sources_count": len(sources_used),
            },
        )

        return Response(
            GeneratedDocumentSerializer(doc_obj, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )

    @staticmethod
    def _get_llm_keys(user=None) -> tuple[str, str]:
        """
        Resolve OpenRouter and Gemini API keys.
        Priority: user.preferences > environment variables.
        Always returns (openrouter_key, gemini_key) — OpenRouter is preferred.
        """
        import os

        openrouter_key = ""
        gemini_key = ""

        # 1. Try user preferences first (keys saved via Settings page)
        if user is not None:
            try:
                # Always refresh from DB so we pick up keys saved in this session
                user.refresh_from_db(fields=["preferences"])
                prefs = getattr(user, "preferences", {}) or {}
                if isinstance(prefs, dict):
                    openrouter_key = prefs.get("openrouter_api_key", "").strip()
                    gemini_key = prefs.get("gemini_api_key", "").strip()
            except Exception:
                pass

        # 2. Fall back to environment variables
        if not openrouter_key:
            openrouter_key = os.environ.get("OPENROUTER_API_KEY", "").strip()
        if not gemini_key:
            gemini_key = os.environ.get("GEMINI_API_KEY", "").strip()

        return openrouter_key, gemini_key

    @staticmethod
    def _retrieve_rag_context(
        prompt: str,
        user,
        content_types: list,
        source_item_ids: list,
    ) -> tuple:
        """
        Retrieve relevant content from the user's Synapse knowledge base.

        Two modes:
          1. source_item_ids provided → fetch those specific items directly from DB
          2. otherwise → run hybrid vector search via SynapseRetriever

        Returns (context_text: str, sources: list[dict])
        """
        sources = []
        context_parts = []

        # ── Mode 1: pinned specific items ──────────────────────────────────────
        if source_item_ids:
            type_to_model = {
                "article": ("apps.articles.models", "Article"),
                "paper": ("apps.papers.models", "ResearchPaper"),
                "repository": ("apps.repositories.models", "Repository"),
                "video": ("apps.videos.models", "Video"),
            }
            for item in source_item_ids:
                item_id = item.get("id")
                item_type = (
                    item.get("type", "").lower().rstrip("s")
                )  # normalise "articles" → "article"
                if not item_id or item_type not in type_to_model:
                    continue
                module_path, model_name = type_to_model[item_type]
                try:
                    import importlib

                    mod = importlib.import_module(module_path)
                    Model = getattr(mod, model_name)
                    obj = Model.objects.get(pk=item_id)
                    title_val = (
                        getattr(obj, "title", "")
                        or getattr(obj, "name", "")
                        or str(obj)
                    )
                    summary_val = (
                        getattr(obj, "summary", "")
                        or getattr(obj, "abstract", "")
                        or getattr(obj, "description", "")
                        or getattr(obj, "body", "")[:1000]
                        or ""
                    )
                    url_val = (
                        getattr(obj, "url", "") or getattr(obj, "html_url", "") or ""
                    )
                    context_parts.append(
                        f"[{item_type.upper()}] {title_val}\n{summary_val}"
                    )
                    sources.append(
                        {
                            "id": str(item_id),
                            "content_type": item_type,
                            "title": title_val,
                            "url": url_val,
                            "snippet": summary_val[:300],
                        }
                    )
                except Exception as exc:
                    logger.warning(
                        "Could not fetch pinned item %s (%s): %s",
                        item_id,
                        item_type,
                        exc,
                    )

        # ── Mode 2: vector search ──────────────────────────────────────────────
        else:
            try:
                from ai_engine.rag.retriever import SynapseRetriever

                retriever = SynapseRetriever(
                    k=4,
                    mode="hybrid",
                    use_reranker=False,  # skip reranker for speed during doc generation
                    content_types=content_types
                    or list(
                        __import__(
                            "ai_engine.rag.retriever", fromlist=["COLLECTION_NAMES"]
                        ).COLLECTION_NAMES.keys()
                    ),
                )
                docs = retriever.invoke(f"{prompt} {title}")
                for doc in docs[:8]:  # cap at 8 sources
                    meta = doc.metadata or {}
                    title_v = meta.get("title", "")
                    url_v = meta.get("source", "") or meta.get("url", "")
                    ct = meta.get("content_type", "unknown")
                    snippet = doc.page_content[:300]
                    context_parts.append(
                        f"[{ct.upper()}] {title_v}\n{doc.page_content}"
                    )
                    sources.append(
                        {
                            "id": str(meta.get("id", "")),
                            "content_type": ct,
                            "title": title_v,
                            "url": url_v,
                            "snippet": snippet,
                            "similarity_score": meta.get("rerank_score")
                            or meta.get("rrf_score")
                            or meta.get("similarity_score", 0.0),
                        }
                    )
                logger.info(
                    "RAG retrieved %d sources for document generation", len(sources)
                )
            except Exception as exc:
                logger.warning("RAG retrieval failed (non-fatal): %s", exc)

        context_text = "\n\n---\n\n".join(context_parts) if context_parts else ""
        return context_text, sources

    @staticmethod
    def _expand_prompt_to_sections(
        prompt: str,
        title: str,
        doc_type: str,
        user=None,
        model_override: str = "",
        content_types: list = None,
        source_item_ids: list = None,
    ) -> tuple:
        """
        Use OpenRouter (preferred) or Gemini to expand a free-text prompt into
        a list of structured sections: [{"heading": str, "content": str}, ...].

        Now RAG-enhanced: retrieves relevant content from the user's Synapse
        knowledge base (saved articles, papers, repos, videos) and injects it
        as grounding context into the LLM prompt before generating sections.

        Returns (sections: list, sources_used: list[dict])
        """
        import json
        import re

        # ── Step 1: Retrieve RAG context from user's knowledge base ───────────
        context_text, sources_used = DocumentGenerateView._retrieve_rag_context(
            prompt=prompt,
            user=user,
            content_types=content_types or [],
            source_item_ids=source_item_ids or [],
        )

        # ── Step 2: Build LLM prompt — inject context if available ────────────
        if context_text:
            rag_block = (
                "KNOWLEDGE BASE CONTEXT\n"
                "The following content has been retrieved from the user's personal Synapse "
                "knowledge base (their saved articles, research papers, repositories, and videos). "
                "Use this as your PRIMARY source material. Ground your document in this real content. "
                "Reference specific titles, findings, and details from these sources.\n\n"
                f"{context_text}\n\n"
                "END OF KNOWLEDGE BASE CONTEXT\n"
            )
        else:
            rag_block = ""

        system = (
            "You are a senior research analyst and strategy consultant with expertise across industries, "
            "markets, and disciplines — equivalent to a Principal at McKinsey, BCG, or Bain. "
            "Your writing is authoritative, evidence-based, and deeply analytical.\n\n"
            + (rag_block if rag_block else "")
            + "TASK: Write a comprehensive, research-grade professional document on the given topic.\n\n"
            "STRUCTURAL REQUIREMENTS:\n"
            "- Generate exactly 6-8 major sections with clear, specific headings.\n"
            "- Each section: minimum 4 full paragraphs, each paragraph 5-7 sentences.\n"
            "- Paragraphs separated by double newlines (\\n\\n).\n\n"
            "CONTENT REQUIREMENTS:\n"
            "- Write specifically about the EXACT topic given — not a generic template.\n"
            + (
                "- Draw on and reference the knowledge base content provided above.\n"
                "- Cite specific items from the knowledge base by their title where relevant.\n"
                if rag_block
                else ""
            )
            + "- Include real facts, statistics, named organisations, countries, dates, and figures.\n"
            "- Provide expert analysis: causation, trends, strategic implications, future scenarios.\n"
            "- Vary paragraph types: analytical, comparative, descriptive, prescriptive.\n"
            "- Zero generic filler — every sentence must add specific, substantive value.\n\n"
            "OUTPUT FORMAT:\n"
            "Return ONLY a valid JSON array. Each element must have exactly two keys:\n"
            "  'heading': string (the section title)\n"
            "  'content': string (the full section text with \\n\\n between paragraphs)\n"
            "CRITICAL JSON RULES:\n"
            "  - Do NOT use backslash escapes inside content strings (no \\t, \\e, \\s etc.)\n"
            "  - Separate paragraphs with the literal characters \\n\\n (double newline)\n"
            "  - Do not use any special characters that would break JSON parsing\n"
            "  - No markdown fences (```), no preamble, no explanation outside the array\n"
            "  - Return ONLY the raw JSON array starting with [ and ending with ]"
        )
        user_msg = (
            f"Topic: {title}\n"
            f"Document type: {doc_type}\n"
            f"User instruction: {prompt}\n\n"
            + (
                "Use the knowledge base context provided in the system prompt as your primary source. "
                if rag_block
                else ""
            )
            + "Write a deeply researched, expert-level document SPECIFICALLY about the topic above. "
            "Do NOT use generic templates or placeholder text. Every sentence must be about this "
            "specific topic with real, concrete, topic-specific information. "
            "Return the JSON sections array now."
        )

        openrouter_key, gemini_key = DocumentGenerateView._get_llm_keys(user)

        # QA-16: LLM invocation delegated to _DocumentLLMBuilder to separate
        # concerns — this method handles RAG + prompt construction + JSON parsing;
        # _DocumentLLMBuilder handles provider selection and raw invocation.
        raw = _DocumentLLMBuilder.invoke(
            system=system,
            user_msg=user_msg,
            openrouter_key=openrouter_key,
            gemini_key=gemini_key,
            model_override=model_override,
        )

        if raw:
            # Extract the JSON array from the response — handles markdown fences,
            # preamble text, and other surrounding content the LLM may add.
            try:
                # Strip markdown code fences (```json ... ``` or ``` ... ```)
                clean = re.sub(r"```(?:json)?\s*", "", raw).strip()
                clean = re.sub(r"```\s*$", "", clean).strip()

                # Find the first '[' and last ']' to isolate the JSON array
                start = clean.index("[")
                end = clean.rindex("]") + 1
                json_str = clean[start:end]

                # Sanitise the JSON string character-by-character:
                # - Replace bare newlines/tabs inside string values with \n/\t
                # - Fix invalid escape sequences (e.g. \e, \a, \T) → double-escape them
                VALID_ESCAPES = set('"\\/ bfnrtu')
                sanitised = []
                in_string = False
                escape_next = False
                for ch in json_str:
                    if escape_next:
                        # Check if this is a valid JSON escape character
                        if ch in VALID_ESCAPES:
                            sanitised.append(ch)
                        else:
                            # Invalid escape: double-escape the backslash
                            sanitised.append("\\")
                            sanitised.append(ch)
                        escape_next = False
                        continue
                    if ch == "\\" and in_string:
                        escape_next = True
                        sanitised.append(ch)
                        continue
                    if ch == '"':
                        in_string = not in_string
                        sanitised.append(ch)
                        continue
                    if in_string and ch == "\n":
                        sanitised.append("\\n")
                        continue
                    if in_string and ch == "\r":
                        sanitised.append("\\r")
                        continue
                    if in_string and ch == "\t":
                        sanitised.append("\\t")
                        continue
                    sanitised.append(ch)
                json_str = "".join(sanitised)

                # Try parsing; if truncated, attempt to repair by trimming to last complete object
                sections = None
                try:
                    sections = json.loads(json_str)
                except json.JSONDecodeError as exc:
                    logger.warning("JSON parse failed (%s) — attempting repair", exc)
                    # Repair strategy: find last complete {...} block and close the array
                    try:
                        # Find the last complete JSON object (ends with "}")
                        last_complete = json_str.rfind("},")
                        if last_complete == -1:
                            last_complete = json_str.rfind("}")
                        if last_complete > 100:
                            repaired = json_str[: last_complete + 1].rstrip() + "\n]"
                            sections = json.loads(repaired)
                            logger.info(
                                "JSON repaired: recovered %d sections", len(sections)
                            )
                    except json.JSONDecodeError:
                        # Last resort: extract individual objects with regex
                        try:
                            objects = re.findall(
                                r'\{\s*"heading"\s*:\s*"([^"]+)"\s*,\s*"content"\s*:\s*"((?:[^"\\]|\\.)*)"\s*\}',
                                json_str,
                                re.DOTALL,
                            )
                            if objects:
                                sections = [
                                    {"heading": h, "content": c.replace("\\n", "\n\n")}
                                    for h, c in objects
                                ]
                                logger.info(
                                    "Regex fallback: recovered %d sections",
                                    len(sections),
                                )
                        except Exception:
                            pass

                if isinstance(sections, list) and sections:
                    valid_sections = [
                        s
                        for s in sections
                        if isinstance(s, dict)
                        and s.get("heading", "").strip()
                        and s.get("content", "").strip()
                        and len(s.get("content", "")) > 50
                    ]
                    if valid_sections:
                        logger.info(
                            "LLM expanded prompt into %d sections", len(valid_sections)
                        )
                        return valid_sections, sources_used
            except (ValueError, Exception) as exc:
                logger.warning(
                    "Failed to parse LLM sections JSON: %s — raw: %s", exc, raw[:200]
                )

        # Smart fallback: generate structured sections using the prompt as a topic guide.
        # This produces a real, readable document even without an LLM.
        logger.info(
            "Generating structured fallback sections from prompt (no LLM or parse failed)"
        )
        return (
            DocumentGenerateView._build_fallback_sections(prompt, title, doc_type),
            sources_used,
        )

    @staticmethod
    def _build_fallback_sections(prompt: str, title: str, doc_type: str) -> list:
        """
        Build a structured, content-rich set of document sections without an LLM.
        Uses the prompt and title to produce a coherent, meaningful document skeleton
        with real prose content rather than placeholder text.
        """
        # Derive a clean topic description from the title and prompt
        topic = title.strip() if title.strip() else prompt[:80].strip()

        sections = [
            {
                "heading": "Executive Summary",
                "content": (
                    f"This document provides a comprehensive analysis and overview of {topic}. "
                    f"The following report has been prepared to deliver a thorough examination of the subject, "
                    f"covering key concepts, current developments, and future outlook. "
                    f"It is intended to serve as a definitive reference for understanding the landscape, "
                    f"challenges, and opportunities associated with this topic.\n\n"
                    f"{prompt.strip()}\n\n"
                    f"The insights presented in this document draw from a wide range of considerations "
                    f"including technical, strategic, economic, and societal perspectives. "
                    f"Readers will gain a structured understanding of the topic from multiple dimensions, "
                    f"enabling informed decision-making and strategic planning."
                ),
            },
            {
                "heading": "Introduction and Background",
                "content": (
                    f"The subject of {topic} has gained significant prominence in recent years, "
                    f"driven by rapid changes in technology, market dynamics, and global trends. "
                    f"Understanding the historical context and foundational principles is essential "
                    f"for any meaningful engagement with this topic. This section outlines the origins, "
                    f"evolution, and key milestones that have shaped the current state of affairs.\n\n"
                    f"The landscape surrounding {topic} is complex and multifaceted, involving a wide "
                    f"array of stakeholders, technologies, and regulatory frameworks. Over the past decade, "
                    f"we have witnessed transformative shifts that have redefined the boundaries of what is "
                    f"possible. These shifts have created both unprecedented opportunities and significant challenges "
                    f"for individuals, organisations, and governments alike.\n\n"
                    f"A thorough grounding in the background of {topic} is therefore indispensable for "
                    f"anyone seeking to navigate this domain with clarity and confidence. The sections that follow "
                    f"build upon this foundation to deliver actionable insights and forward-looking perspectives."
                ),
            },
            {
                "heading": "Key Concepts and Core Principles",
                "content": (
                    f"At the heart of {topic} lie a set of core concepts and principles that define "
                    f"the field and guide its practitioners. These foundational ideas form the intellectual "
                    f"scaffolding upon which all further analysis rests. A clear understanding of these concepts "
                    f"is critical for distinguishing signal from noise and for making sound judgements.\n\n"
                    f"Among the most important principles is the recognition that {topic} operates within "
                    f"a dynamic system of interdependencies. No single element can be understood in isolation; "
                    f"rather, each component influences and is influenced by the broader ecosystem. This systems-level "
                    f"thinking is essential for anticipating second-order effects and unintended consequences.\n\n"
                    f"Furthermore, the principles underpinning {topic} are not static — they evolve in response "
                    f"to new evidence, technological advances, and changing societal values. Staying abreast of "
                    f"these developments requires continuous learning and a willingness to challenge established assumptions. "
                    f"Practitioners who embrace this mindset are best positioned to lead and innovate in the field."
                ),
            },
            {
                "heading": "Current State and Recent Developments",
                "content": (
                    f"The current state of {topic} reflects a period of rapid evolution and transformation. "
                    f"Recent developments have accelerated progress in ways that were difficult to foresee even "
                    f"a few years ago, creating a landscape that is both exciting and challenging to navigate. "
                    f"Key players — whether they are organisations, governments, or individuals — are actively "
                    f"adapting their strategies to keep pace with these changes.\n\n"
                    f"Among the most notable recent trends is the convergence of multiple technologies and disciplines "
                    f"that are reshaping the domain of {topic}. This convergence is creating new possibilities "
                    f"for innovation, collaboration, and value creation. At the same time, it is raising important "
                    f"questions about governance, equity, and sustainability that demand careful consideration.\n\n"
                    f"Data and evidence from recent developments underscore the scale and pace of change. "
                    f"Adoption rates, investment volumes, and output metrics all point to a domain experiencing "
                    f"significant momentum. Understanding these trends in their proper context is essential "
                    f"for making sense of where things stand today and where they are headed."
                ),
            },
            {
                "heading": "Challenges and Considerations",
                "content": (
                    f"Despite the many opportunities presented by {topic}, there are also substantial challenges "
                    f"that must be acknowledged and addressed. These challenges span technical, organisational, "
                    f"ethical, and regulatory dimensions, and they require thoughtful and collaborative responses "
                    f"from all stakeholders involved.\n\n"
                    f"On the technical side, the complexity inherent in {topic} means that solutions are rarely "
                    f"straightforward. Integration with existing systems, ensuring reliability and security, "
                    f"and managing the costs of implementation are perennial concerns that demand rigorous "
                    f"engineering and project management. Failure to address these technical hurdles can undermine "
                    f"even the most well-conceived initiatives.\n\n"
                    f"Beyond the technical realm, there are significant human and organisational factors to consider. "
                    f"Change management, skills development, and cultural alignment are often the decisive factors "
                    f"in whether an initiative succeeds or fails. Leaders must invest in people as much as in "
                    f"technology, fostering an environment of trust, learning, and adaptability to navigate "
                    f"the challenges that lie ahead."
                ),
            },
            {
                "heading": "Opportunities and Strategic Recommendations",
                "content": (
                    f"Notwithstanding the challenges, {topic} presents a wealth of opportunities for those "
                    f"who are prepared to engage with it strategically and thoughtfully. The organisations and "
                    f"individuals who stand to benefit most are those who move early, invest wisely, and build "
                    f"the capabilities needed to capitalise on emerging trends.\n\n"
                    f"From a strategic perspective, the most important step is to develop a clear and coherent "
                    f"vision for how {topic} fits into broader goals and objectives. This vision should be "
                    f"grounded in a realistic assessment of current capabilities and resources, while also "
                    f"being sufficiently ambitious to drive meaningful progress. A phased approach — starting "
                    f"with high-impact, lower-risk initiatives and scaling over time — is often the most effective path.\n\n"
                    f"Investment in research, partnerships, and talent development is equally important. "
                    f"The most successful actors in this domain consistently prioritise learning and collaboration, "
                    f"recognising that no single organisation has all the answers. By building strong ecosystems "
                    f"of partners, suppliers, and communities, they are able to accelerate progress and share "
                    f"both the risks and rewards of innovation."
                ),
            },
            {
                "heading": "Future Outlook and Conclusion",
                "content": (
                    f"Looking ahead, the trajectory of {topic} points toward continued growth, disruption, "
                    f"and transformation. The forces driving change show no signs of abating; if anything, "
                    f"they are likely to intensify as new technologies mature and adoption broadens. "
                    f"Those who engage proactively with this evolution will be far better positioned than "
                    f"those who wait for the landscape to stabilise before acting.\n\n"
                    f"The long-term outlook for {topic} is broadly positive, with significant potential "
                    f"to create value, improve outcomes, and address some of the most pressing challenges "
                    f"facing society. However, realising this potential will require sustained commitment, "
                    f"smart investment, and a willingness to navigate uncertainty with resilience and creativity. "
                    f"The stakes are high, but so are the rewards for those who rise to the occasion.\n\n"
                    f"In conclusion, this document has sought to provide a comprehensive and balanced perspective "
                    f"on {topic}, covering its background, key concepts, current state, challenges, opportunities, "
                    f"and future outlook. The analysis presented here is intended to serve as a starting point "
                    f"for deeper exploration and informed action. We encourage readers to engage critically "
                    f"with the ideas presented and to contribute their own expertise and insights to the ongoing "
                    f"conversation about how best to navigate and shape this important domain."
                ),
            },
        ]

        return sections

    @staticmethod
    def _call_tool(
        doc_type: str,
        title: str,
        sections: list,
        subtitle: str,
        author: str,
        user_id: str,
    ) -> tuple[str, str]:
        """
        Call the appropriate doc generation tool and return (result_str, abs_file_path).
        """
        from ai_engine.agents.doc_tools import (
            _generate_html,
            _generate_markdown,
            _generate_pdf,
            _generate_ppt,
            _generate_word_doc,
        )

        if doc_type == "pdf":
            result = _generate_pdf(
                title=title,
                sections=sections,
                subtitle=subtitle,
                author=author,
                user_id=user_id,
            )
        elif doc_type == "ppt":
            # Convert sections → slides with rich bullet points derived from content
            slides = []
            for s in sections:
                heading = s.get("heading", "Slide")
                content = s.get("content", "")
                # Extract first sentence of each paragraph as a bullet point
                paras = [p.strip() for p in content.split("\n\n") if p.strip()]
                bullets = []
                for para in paras[:5]:
                    first = para.split(".")[0].strip()
                    if first and len(first) > 10:
                        bullets.append(first + ".")
                # Fallback: truncate content into 3 bullets
                if not bullets and content:
                    words = content.split()
                    chunk = 12
                    for i in range(0, min(len(words), chunk * 3), chunk):
                        bullets.append(" ".join(words[i : i + chunk]) + "…")
                slides.append(
                    {
                        "title": heading,
                        "bullets": bullets,
                        "content": content,
                        "notes": paras[0][:300] if paras else "",
                    }
                )
            result = _generate_ppt(
                title=title,
                slides=slides,
                subtitle=subtitle or "Generated by SYNAPSE AI",
                author=author,
                user_id=user_id,
            )
        elif doc_type == "word":
            # Add level=1 to each section if missing
            word_sections = [
                {
                    "heading": s.get("heading", "Section"),
                    "content": s.get("content", ""),
                    "level": s.get("level", 1),
                }
                for s in sections
            ]
            result = _generate_word_doc(
                title=title,
                sections=word_sections,
                author=author,
                add_toc=True,
                user_id=user_id,
            )
        elif doc_type == "markdown":
            result = _generate_markdown(
                title=title,
                sections=sections,
                author=author,
                user_id=user_id,
            )
        elif doc_type == "html":
            result = _generate_html(
                title=title,
                sections=sections,
                subtitle=subtitle or "",
                author=author,
                user_id=user_id,
            )
        else:
            raise ValueError(f"Unsupported doc_type: {doc_type}")

        # Extract file path from result string ("Path: /abs/path")
        file_path = ""
        for line in result.splitlines():
            if line.startswith("Path:"):
                file_path = line.replace("Path:", "").strip()
                break

        return result, file_path


# ---------------------------------------------------------------------------
# Generate Project (Phase 5.3)
# ---------------------------------------------------------------------------


class ProjectGenerateView(APIView):
    """
    POST /api/v1/documents/generate-project/

    Generates a project scaffold (.zip) using the create_project agent tool.
    Supported project_type values: django, fastapi, nextjs, datascience, react_lib
    Optional features list: ['auth', 'testing', 'ci_cd']
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request) -> Response:
        from .serializers import ProjectGenerateSerializer

        serializer = ProjectGenerateSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        data = serializer.validated_data
        project_type = data["project_type"]
        name = data["name"]
        features = data.get("features", [])
        description = data.get("description", "") or ""
        user_id = str(request.user.id)

        # ── HTML Template: generate actual HTML/CSS/JS directly from the prompt ─
        if project_type == "html_template":
            prompt = description or (
                f"Create a stunning, production-ready HTML template called '{name}'. "
                "Use modern HTML5, CSS3 (inline <style>), and vanilla JavaScript. "
                "Make it visually impressive, fully responsive, and self-contained in a single HTML file. "
                "Include beautiful gradients, micro-interactions, smooth animations, and premium typography. "
                "Use Google Fonts and CDN libraries only — no build tools required."
            )
            openrouter_key, gemini_key = DocumentGenerateView._get_llm_keys(
                request.user
            )
            try:
                from ai_engine.agents.doc_tools import _generate_html_page_from_prompt

                result_str = _generate_html_page_from_prompt(
                    title=name,
                    prompt=prompt,
                    subtitle="HTML Template",
                    author="SYNAPSE AI",
                    user_id=user_id,
                    openrouter_key=openrouter_key,
                    gemini_key=gemini_key,
                )
            except Exception as exc:
                logger.error("HTML template generation failed: %s", exc)
                return Response(
                    {"error": f"HTML template generation failed: {exc}"},
                    status=status.HTTP_500_INTERNAL_SERVER_ERROR,
                )
            file_path_str = ""
            for line in result_str.splitlines():
                if line.startswith("Path:"):
                    file_path_str = line.replace("Path:", "").strip()
                    break
            sections = []

            abs_path = Path(file_path_str) if file_path_str else None
            rel_path = ""
            file_size = 0
            if abs_path and abs_path.exists():
                media_root = Path(settings.MEDIA_ROOT)
                try:
                    rel_path = str(abs_path.relative_to(media_root))
                except ValueError:
                    rel_path = str(abs_path)
                file_size = abs_path.stat().st_size

            doc_obj = GeneratedDocument.objects.create(
                user=request.user,
                title=name,
                doc_type="html",
                file_path=rel_path,
                file_size_bytes=file_size,
                agent_prompt=prompt,
                metadata={
                    "project_type": "html_template",
                    "project_name": name,
                    "section_count": len(sections),
                    "sections": sections,
                    "description": description,
                },
            )
            return Response(
                GeneratedDocumentSerializer(doc_obj, context={"request": request}).data,
                status=status.HTTP_201_CREATED,
            )

        # ── Standard project scaffold ─────────────────────────────────────────
        try:
            from ai_engine.agents.project_tools import _create_project

            result_str = _create_project(
                project_type=project_type,
                name=name,
                features=features,
                description=description,
                user_id=user_id,
            )
        except Exception as exc:
            logger.error("Project generation failed: %s", exc)
            return Response(
                {"error": f"Project generation failed: {exc}"},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

        if (
            "failed" in result_str.lower()
            or "unknown project_type" in result_str.lower()
        ):
            return Response({"error": result_str}, status=status.HTTP_400_BAD_REQUEST)

        # Parse file path from result string ("Path: /abs/path")
        abs_path_str = ""
        for line in result_str.splitlines():
            if line.startswith("Path:"):
                abs_path_str = line.replace("Path:", "").strip()
                break

        abs_path = Path(abs_path_str) if abs_path_str else None
        rel_path = ""
        file_size = 0
        if abs_path and abs_path.exists():
            media_root = Path(settings.MEDIA_ROOT)
            try:
                rel_path = str(abs_path.relative_to(media_root))
            except ValueError:
                rel_path = str(abs_path)
            file_size = abs_path.stat().st_size

        title = f"{name} ({project_type} scaffold)"
        doc_obj = GeneratedDocument.objects.create(
            user=request.user,
            title=title,
            doc_type="project",
            file_path=rel_path,
            file_size_bytes=file_size,
            agent_prompt=description or f"Generate {project_type} project: {name}",
            metadata={
                "project_type": project_type,
                "project_name": name,
                "features": features,
                "description": description,
                "tool_result": result_str[:500],
            },
        )

        return Response(
            GeneratedDocumentSerializer(doc_obj, context={"request": request}).data,
            status=status.HTTP_201_CREATED,
        )


# ---------------------------------------------------------------------------
# List
# ---------------------------------------------------------------------------


class DocumentListView(APIView):
    """GET /api/v1/documents/ — list the authenticated user's documents."""

    permission_classes = [IsAuthenticated]

    def get(self, request: Request) -> Response:
        qs = GeneratedDocument.objects.filter(user=request.user).order_by("-created_at")

        doc_type_filter = request.query_params.get("doc_type")
        if doc_type_filter:
            qs = qs.filter(doc_type=doc_type_filter)

        paginator = StandardPagination()
        page = paginator.paginate_queryset(qs, request)
        # Use the full serializer so download_url is included in list responses
        serializer = GeneratedDocumentSerializer(
            page, many=True, context={"request": request}
        )
        return paginator.get_paginated_response(serializer.data)


# ---------------------------------------------------------------------------
# Detail + Delete
# ---------------------------------------------------------------------------


class DocumentDetailView(APIView):
    """
    GET    /api/v1/documents/{id}/ — retrieve document metadata
    DELETE /api/v1/documents/{id}/ — delete document record + file
    """

    permission_classes = [IsAuthenticated]

    def _get_doc(self, doc_id, user) -> GeneratedDocument | None:
        try:
            return GeneratedDocument.objects.get(id=doc_id, user=user)
        except GeneratedDocument.DoesNotExist:
            return None

    def get(self, request: Request, doc_id: str) -> Response:
        doc = self._get_doc(doc_id, request.user)
        if not doc:
            return Response(
                {"error": "Document not found."}, status=status.HTTP_404_NOT_FOUND
            )
        return Response(
            GeneratedDocumentSerializer(doc, context={"request": request}).data
        )

    def delete(self, request: Request, doc_id: str) -> Response:
        doc = self._get_doc(doc_id, request.user)
        if not doc:
            return Response(
                {"error": "Document not found."}, status=status.HTTP_404_NOT_FOUND
            )

        # Delete physical file
        if doc.file_path:
            abs_path = Path(settings.MEDIA_ROOT) / doc.file_path
            if abs_path.exists():
                try:
                    abs_path.unlink()
                except Exception as exc:
                    logger.warning("Could not delete file %s: %s", abs_path, exc)

        doc.delete()
        return Response(status=status.HTTP_204_NO_CONTENT)


# ---------------------------------------------------------------------------
# Download
# ---------------------------------------------------------------------------


class DocumentDownloadView(APIView):
    """GET /api/v1/documents/{id}/download/ — stream the file as a download."""

    permission_classes = [IsAuthenticated]

    def get(self, request: Request, doc_id: str) -> FileResponse:
        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404("Document not found.")

        if not doc.file_path:
            return Response(
                {"error": "No file available for this document."},
                status=status.HTTP_404_NOT_FOUND,
            )

        stored = Path(doc.file_path)
        # If file_path is already absolute, use it directly; otherwise join with MEDIA_ROOT
        if stored.is_absolute():
            abs_path = stored
        else:
            abs_path = Path(settings.MEDIA_ROOT) / stored

        if not abs_path.exists():
            logger.error(
                "Download failed — file missing on disk. doc_id=%s file_path=%s resolved=%s MEDIA_ROOT=%s",
                doc_id,
                doc.file_path,
                abs_path,
                settings.MEDIA_ROOT,
            )
            return Response(
                {"error": f"File not found on server. Expected at: {abs_path}"},
                status=status.HTTP_404_NOT_FOUND,
            )

        content_type = _MIME_MAP.get(doc.doc_type, "application/octet-stream")
        ext = _EXT_MAP.get(doc.doc_type, "")
        filename = f"{doc.title[:50].replace('/', '_')}{ext}"

        # Open file via context manager; FileResponse closes it when done streaming
        fh = open(abs_path, "rb")  # noqa: WPS515 – FileResponse takes ownership
        response = FileResponse(
            fh,
            content_type=content_type,
            as_attachment=True,
            filename=filename,
        )
        return response


class DocumentPreviewView(APIView):
    """
    GET /api/v1/documents/{id}/preview/
    Returns a PNG thumbnail/preview image of the document.
    For PDF: renders first page via Pillow + ReportLab.
    For HTML: renders a styled preview card.
    For all others: returns a branded cover card image.
    """

    permission_classes = [IsAuthenticated]

    def get(self, request: Request, doc_id: str) -> Response:
        from django.http import HttpResponse

        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        try:
            png_bytes = self._render_preview(doc)
            return HttpResponse(png_bytes, content_type="image/png")
        except Exception as exc:
            logger.error("Preview generation failed: %s", exc)
            return Response({"error": str(exc)}, status=500)

    @staticmethod
    def _render_preview(doc) -> bytes:
        """
        Generate a real first-page visual preview of the document.
        - PDF  → PyMuPDF renders actual first page (falls back to cover image)
        - PPTX → python-pptx + Pillow renders title slide shapes
        - DOCX → python-docx + Pillow renders first page text
        - HTML → Pillow renders styled HTML layout preview
        - MD   → Pillow renders markdown content preview
        Falls back to branded cover image if rendering fails.
        """
        abs_path = None
        if doc.file_path:
            abs_path = Path(settings.MEDIA_ROOT) / doc.file_path

        doc_type = doc.doc_type
        title = doc.title
        subtitle = (doc.metadata or {}).get("subtitle", "")
        author = (doc.metadata or {}).get("author", "SYNAPSE AI")
        sections = (doc.metadata or {}).get("sections", [])

        try:
            if doc_type == "pdf" and abs_path and abs_path.exists():
                return DocumentPreviewView._render_pdf_preview(str(abs_path))
            elif doc_type == "ppt" and abs_path and abs_path.exists():
                return DocumentPreviewView._render_pptx_preview(
                    str(abs_path), title, subtitle
                )
            elif doc_type in ("word", "docx") and abs_path and abs_path.exists():
                return DocumentPreviewView._render_docx_preview(
                    str(abs_path), title, sections
                )
            elif doc_type == "html" and abs_path and abs_path.exists():
                return DocumentPreviewView._render_html_preview(
                    str(abs_path), title, sections
                )
            elif doc_type == "markdown" and abs_path and abs_path.exists():
                return DocumentPreviewView._render_markdown_preview(
                    str(abs_path), title, sections
                )
        except Exception as exc:
            logger.warning(
                "Real preview rendering failed for %s (%s): %s", doc_type, doc.id, exc
            )

        # Fallback to branded cover image
        from ai_engine.agents.doc_tools import _cover_image_bytes

        return _cover_image_bytes(
            title=title,
            subtitle=subtitle or doc_type.upper() + " Document",
            author=author,
            doc_type=doc_type,
        )

    # ── PDF preview via PyMuPDF ───────────────────────────────────────────────
    @staticmethod
    def _render_pdf_preview(path: str) -> bytes:
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        page = doc[0]
        mat = fitz.Matrix(1.8, 1.8)  # 1.8× zoom → ~1080px wide
        pix = page.get_pixmap(matrix=mat, alpha=False)
        png = pix.tobytes("png")
        doc.close()
        # Crop to 1200×630 social-card ratio
        return DocumentPreviewView._crop_to_ratio(png, 1200, 630)

    # ── PPTX preview via python-pptx + Pillow ────────────────────────────────
    @staticmethod
    def _render_pptx_preview(path: str, title: str, subtitle: str) -> bytes:
        import io as _io

        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(path)
        slide = prs.slides[0]  # Title slide

        W, H = 1200, 630

        # ── Background gradient from slide background fill ────────────────
        # Try to read the gradient colours from the slide XML
        bg_c1 = (55, 48, 163)  # indigo-dark fallback
        bg_c2 = (124, 58, 237)  # violet fallback
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                # Try gradient stops
                from pptx.oxml.ns import qn

                gsLst = bg._element.find(".//" + qn("a:gsLst"))
                if gsLst is not None:
                    stops = gsLst.findall(qn("a:gs"))
                    if len(stops) >= 2:

                        def _hex_to_rgb(h):
                            h = h.lstrip("#")
                            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

                        s0 = stops[0].find(".//" + qn("a:srgbClr"))
                        s1 = stops[-1].find(".//" + qn("a:srgbClr"))
                        if s0 is not None:
                            bg_c1 = _hex_to_rgb(s0.get("val", "3730A3"))
                        if s1 is not None:
                            bg_c2 = _hex_to_rgb(s1.get("val", "7C3AED"))
        except Exception:
            pass

        img = Image.new("RGB", (W, H))
        draw = ImageDraw.Draw(img)

        # Draw gradient background
        for y in range(H):
            t = y / H
            r = int(bg_c1[0] + t * (bg_c2[0] - bg_c1[0]))
            g = int(bg_c1[1] + t * (bg_c2[1] - bg_c1[1]))
            b = int(bg_c1[2] + t * (bg_c2[2] - bg_c1[2]))
            draw.line([(0, y), (W, y)], fill=(r, g, b))

        # Draw shapes from slide
        FONT_BOLD = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        ]
        FONT_REG = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]
        import os as _os

        def _font(paths, size):
            for fp in paths:
                if _os.path.exists(fp):
                    try:
                        return ImageFont.truetype(fp, size)
                    except:
                        continue
            return ImageFont.load_default()

        f_title = _font(FONT_BOLD, 52)
        f_sub = _font(FONT_REG, 22)
        f_meta = _font(FONT_REG, 14)

        # Slide dimensions (EMU)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Map shape position to image coordinates
            sx = int(shape.left / slide_w * W)
            sy = int(shape.top / slide_h * H)
            sw = int(shape.width / slide_w * W)
            # Pick font based on text size in shape
            try:
                sz_pt = shape.text_frame.paragraphs[0].runs[0].font.size
                sz_pt = int(sz_pt.pt) if sz_pt else 18
            except:
                sz_pt = 18
            font = _font(
                FONT_BOLD if sz_pt > 20 else FONT_REG,
                max(12, min(48, int(sz_pt * W / 960))),
            )
            color = (255, 255, 255)
            try:
                rgb = shape.text_frame.paragraphs[0].runs[0].font.color.rgb
                color = (rgb.r, rgb.g, rgb.b)
            except:
                pass
            # Wrap text
            words = text.split()
            lines, cur = [], ""
            for w in words:
                test = (cur + " " + w).strip()
                if draw.textlength(test, font=font) > sw * 0.95:
                    if cur:
                        lines.append(cur)
                    cur = w
                else:
                    cur = test
            if cur:
                lines.append(cur)
            for li, line in enumerate(lines[:4]):
                draw.text(
                    (sx, sy + li * int(sz_pt * 1.3 * W / 960)),
                    line,
                    fill=color,
                    font=font,
                )

        # Left accent strip
        draw.rectangle([0, 0, 8, H], fill=(129, 140, 248))
        # Bottom strip
        draw.rectangle([0, H - 50, W, H], fill=(255, 255, 255))
        draw.text(
            (40, H - 36),
            f"SYNAPSE AI  ·  {title[:60]}",
            fill=(79, 70, 229),
            font=_font(FONT_BOLD, 13),
        )

        buf = _io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        return buf.read()

    # ── DOCX preview via python-docx + Pillow ────────────────────────────────
    @staticmethod
    def _render_docx_preview(path: str, title: str, sections: list) -> bytes:
        import io as _io
        import os as _os

        from docx import Document as DocxDoc
        from PIL import Image, ImageDraw, ImageFont

        doc = DocxDoc(path)
        W, H = 1200, 630

        # Blue theme for Word
        img = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Header bar
        draw.rectangle([0, 0, W, 90], fill=(29, 78, 216))
        draw.rectangle([0, 0, 10, H], fill=(59, 130, 246))

        FONT_BOLD = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        ]
        FONT_REG = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]

        def _font(paths, size):
            for fp in paths:
                if _os.path.exists(fp):
                    try:
                        return ImageFont.truetype(fp, size)
                    except:
                        continue
            return ImageFont.load_default()

        draw.text((28, 22), title[:55], fill=(255, 255, 255), font=_font(FONT_BOLD, 36))
        draw.text(
            (28, 68),
            "WORD DOCUMENT  ·  SYNAPSE AI",
            fill=(147, 197, 253),
            font=_font(FONT_REG, 14),
        )

        # Page content — render first paragraphs
        y = 110
        for para in doc.paragraphs[:20]:
            txt = para.text.strip()
            if not txt:
                continue
            is_heading = para.style.name.startswith("Heading") or (
                len(txt) < 60
                and para.runs
                and any(r.bold for r in para.runs if r.text.strip())
            )
            if is_heading:
                if y > 570:
                    break
                draw.rectangle([18, y + 2, 18 + 4, y + 24], fill=(29, 78, 216))
                draw.text(
                    (28, y), txt[:70], fill=(17, 24, 39), font=_font(FONT_BOLD, 16)
                )
                y += 30
            else:
                words = txt.split()
                line, lines_out = "", []
                for word in words:
                    test = line + " " + word if line else word
                    if len(test) > 95:
                        lines_out.append(line)
                        line = word
                    else:
                        line = test
                if line:
                    lines_out.append(line)
                for ln in lines_out[:3]:
                    if y > 590:
                        break
                    draw.text((28, y), ln, fill=(55, 65, 81), font=_font(FONT_REG, 13))
                    y += 18
                y += 4

        # Footer
        draw.rectangle([0, H - 40, W, H], fill=(243, 244, 246))
        draw.rectangle([0, H - 40, W, H - 39], fill=(209, 213, 219))
        draw.text(
            (28, H - 28),
            f"SYNAPSE AI  ·  {title[:50]}  ·  Page 1",
            fill=(107, 114, 128),
            font=_font(FONT_REG, 12),
        )

        buf = _io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        return buf.read()

    # ── HTML preview via Pillow ───────────────────────────────────────────────
    @staticmethod
    def _render_html_preview(path: str, title: str, sections: list) -> bytes:
        import io as _io
        import os as _os

        from PIL import Image, ImageDraw, ImageFont

        # Read and parse key content from the HTML
        html_content = Path(path).read_text(encoding="utf-8")

        W, H = 1200, 630
        img = Image.new("RGB", (W, H), (15, 12, 46))  # dark background
        draw = ImageDraw.Draw(img)

        FONT_BOLD = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        ]
        FONT_REG = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]

        def _font(paths, size):
            for fp in paths:
                if _os.path.exists(fp):
                    try:
                        return ImageFont.truetype(fp, size)
                    except:
                        continue
            return ImageFont.load_default()

        # Sidebar (dark)
        draw.rectangle([0, 0, 220, H], fill=(30, 27, 74))
        draw.rectangle([220, 0, 224, H], fill=(99, 102, 241))

        # Sidebar header
        draw.rectangle([0, 0, 220, 80], fill=(55, 48, 163))
        draw.text(
            (14, 14), "SYNAPSE AI", fill=(199, 210, 254), font=_font(FONT_BOLD, 12)
        )
        draw.text((14, 32), title[:22], fill=(255, 255, 255), font=_font(FONT_BOLD, 14))

        # Sidebar nav items
        icons = ["🔍", "📌", "💡", "📊", "⚡", "🎯", "🔬"]
        sec_names = (
            [s.get("heading", "")[:20] for s in sections[:7]] if sections else []
        )
        for i, (icon, name) in enumerate(zip(icons, sec_names)):
            sy = 90 + i * 38
            if i == 0:
                draw.rectangle([0, sy, 220, sy + 36], fill=(79, 70, 229, 160))
                draw.rectangle([0, sy, 4, sy + 36], fill=(129, 140, 248))
            draw.text(
                (14, sy + 10),
                f"{icon} {name[:18]}",
                fill=(200, 200, 220),
                font=_font(FONT_REG, 13),
            )

        # Main content area — gradient hero
        for y in range(120):
            t = y / 120
            r = int(30 + t * (55 - 30))
            g = int(27 + t * (48 - 27))
            b = int(74 + t * (163 - 74))
            draw.line([(224, y), (W, y)], fill=(r, g, b))

        # Hero title
        draw.text(
            (248, 18),
            "SYNAPSE AI  ·  Executive Report",
            fill=(129, 140, 248),
            font=_font(FONT_REG, 11),
        )
        words = title.split()
        lines_out, cur = [], ""
        for w in words:
            test = (cur + " " + w).strip()
            if len(test) > 28:
                lines_out.append(cur)
                cur = w
            else:
                cur = test
        if cur:
            lines_out.append(cur)
        for li, line in enumerate(lines_out[:2]):
            draw.text(
                (248, 36 + li * 42),
                line,
                fill=(255, 255, 255),
                font=_font(FONT_BOLD, 36),
            )

        # Stat cards row
        card_y = 145
        stats = [
            ("Chapters", str(len(sections)), "#4F46E5"),
            (
                "Words",
                str(sum(len(s.get("content", "").split()) for s in sections)),
                "#7C3AED",
            ),
            ("Format", "HTML", "#06B6D4"),
            ("Level", "Diamond", "#10B981"),
        ]
        for ci, (lbl, val, col) in enumerate(stats):
            cx = 248 + ci * 230
            r, g, b = int(col[1:3], 16), int(col[3:5], 16), int(col[5:7], 16)
            draw.rectangle([cx, card_y, cx + 210, card_y + 80], fill=(26, 23, 64))
            draw.rectangle([cx, card_y, cx + 210, card_y + 4], fill=(r, g, b))
            draw.text(
                (cx + 105 - len(val) * 10, card_y + 14),
                val,
                fill=(r, g, b),
                font=_font(FONT_BOLD, 26),
            )
            draw.text(
                (cx + 105 - len(lbl) * 5, card_y + 50),
                lbl.upper(),
                fill=(107, 114, 128),
                font=_font(FONT_REG, 11),
            )

        # Content section cards
        sec_y = 245
        for si, sec in enumerate(sections[:3]):
            sx = 248 + si * 313
            heading = sec.get("heading", "")[:24]
            content_snippet = sec.get("content", "")[:80].replace("\n", " ")
            draw.rectangle([sx, sec_y, sx + 298, sec_y + 150], fill=(26, 23, 64))
            draw.rectangle([sx, sec_y, sx + 298, sec_y + 4], fill=(79, 70, 229))
            # Section number
            draw.rectangle([sx, sec_y, sx + 36, sec_y + 36], fill=(55, 48, 163))
            draw.text(
                (sx + 8, sec_y + 8),
                f"{si+1:02d}",
                fill=(255, 255, 255),
                font=_font(FONT_BOLD, 16),
            )
            draw.text(
                (sx + 42, sec_y + 8),
                heading,
                fill=(199, 210, 254),
                font=_font(FONT_BOLD, 13),
            )
            # Content snippet
            words2 = content_snippet.split()
            lines2, cur2 = [], ""
            for w in words2:
                test = (cur2 + " " + w).strip()
                if len(test) > 35:
                    lines2.append(cur2)
                    cur2 = w
                else:
                    cur2 = test
            if cur2:
                lines2.append(cur2)
            for li, ln in enumerate(lines2[:4]):
                draw.text(
                    (sx + 10, sec_y + 45 + li * 22),
                    ln + "…" if li == 3 else ln,
                    fill=(148, 163, 184),
                    font=_font(FONT_REG, 12),
                )

        # Scroll progress bar at top
        draw.rectangle([224, 0, W, 3], fill=(79, 70, 229))
        draw.rectangle([224, 0, 224 + 300, 3], fill=(129, 140, 248))

        # Footer
        draw.rectangle([224, H - 36, W, H], fill=(15, 12, 46))
        draw.rectangle([224, H - 37, W, H - 36], fill=(55, 48, 163))
        draw.text(
            (244, H - 26),
            f"SYNAPSE AI  ·  {title[:50]}  ·  Diamond Level  ·  All rights reserved",
            fill=(75, 85, 99),
            font=_font(FONT_REG, 11),
        )

        buf = _io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        return buf.read()

    # ── Markdown preview via Pillow ───────────────────────────────────────────
    @staticmethod
    def _render_markdown_preview(path: str, title: str, sections: list) -> bytes:
        import io as _io
        import os as _os

        from PIL import Image, ImageDraw, ImageFont

        content = Path(path).read_text(encoding="utf-8")
        lines = content.split("\n")

        W, H = 1200, 630
        img = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        FONT_BOLD = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        ]
        FONT_MONO = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
        ]
        FONT_REG = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]

        def _font(paths, size):
            for fp in paths:
                if _os.path.exists(fp):
                    try:
                        return ImageFont.truetype(fp, size)
                    except:
                        continue
            return ImageFont.load_default()

        # Background
        draw.rectangle([0, 0, W, H], fill=(248, 249, 255))

        # Left accent bar (emerald for markdown)
        draw.rectangle([0, 0, 6, H], fill=(16, 185, 129))

        # Top badge bar
        draw.rectangle([0, 0, W, 56], fill=(5, 150, 105))
        draw.text(
            (20, 10), "SYNAPSE AI", fill=(255, 255, 255), font=_font(FONT_BOLD, 13)
        )
        draw.text(
            (20, 30), "MARKDOWN DOCUMENT", fill=(52, 211, 153), font=_font(FONT_REG, 12)
        )
        # Badges
        for bi, (lbl, col) in enumerate(
            [("FINAL", "#10B981"), ("v2.0", "#059669"), ("AI", "#047857")]
        ):
            bx = W - 200 + bi * 62
            r, g, b = int(col[1:3], 16), int(col[3:5], 16), int(col[5:7], 16)
            draw.rectangle([bx, 12, bx + 54, 36], fill=(r, g, b))
            draw.text(
                (bx + 8, 18), lbl, fill=(255, 255, 255), font=_font(FONT_BOLD, 12)
            )

        # Render markdown lines
        y = 72
        for raw_line in lines[:50]:
            if y > 590:
                break
            stripped = raw_line.strip()
            if not stripped or stripped.startswith("---") or stripped.startswith("```"):
                continue
            if stripped.startswith("# "):
                txt = stripped[2:][:70]
                draw.text((20, y), txt, fill=(17, 24, 39), font=_font(FONT_BOLD, 28))
                y += 38
            elif stripped.startswith("## "):
                txt = stripped[3:][:70]
                draw.rectangle([20, y + 2, 24, y + 22], fill=(16, 185, 129))
                draw.text((30, y), txt, fill=(5, 150, 105), font=_font(FONT_BOLD, 18))
                y += 26
            elif stripped.startswith("### "):
                txt = stripped[4:][:80]
                draw.text((30, y), txt, fill=(55, 65, 81), font=_font(FONT_BOLD, 14))
                y += 20
            elif stripped.startswith(">"):
                txt = stripped.lstrip("> ").strip()[:100]
                draw.rectangle([20, y, 24, y + 18], fill=(16, 185, 129))
                draw.rectangle([20, y, 500, y + 18], fill=(236, 253, 245))
                draw.text(
                    (30, y + 2), txt, fill=(5, 150, 105), font=_font(FONT_REG, 12)
                )
                y += 22
            elif stripped.startswith("!["):
                continue  # skip image tags
            elif stripped.startswith("|"):
                # Table row
                cells = [c.strip() for c in stripped.split("|")[1:-1]]
                cell_w = (W - 40) // max(len(cells), 1)
                draw.rectangle([20, y, W - 20, y + 22], fill=(243, 244, 246))
                for ci, cell in enumerate(cells[:6]):
                    draw.text(
                        (24 + ci * cell_w, y + 4),
                        cell[:18],
                        fill=(55, 65, 81),
                        font=_font(FONT_REG, 12),
                    )
                y += 24
            elif stripped.startswith("**"):
                txt = stripped.strip("*")[:100]
                draw.text((20, y), txt, fill=(17, 24, 39), font=_font(FONT_BOLD, 13))
                y += 18
            elif stripped.startswith("![") or stripped.startswith("http"):
                continue
            else:
                txt = stripped[:120]
                draw.text((20, y), txt, fill=(75, 85, 99), font=_font(FONT_REG, 12))
                y += 17

        # Footer
        draw.rectangle([0, H - 36, W, H], fill=(243, 244, 246))
        draw.rectangle([0, H - 37, W, H - 36], fill=(209, 213, 219))
        draw.text(
            (20, H - 26),
            f"SYNAPSE AI  ·  {title[:60]}  ·  Markdown v2.0",
            fill=(107, 114, 128),
            font=_font(FONT_REG, 12),
        )

        buf = _io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        return buf.read()

    # ── Helper: crop PNG bytes to target ratio ────────────────────────────────
    @staticmethod
    def _crop_to_ratio(png_bytes: bytes, target_w: int, target_h: int) -> bytes:
        import io as _io

        from PIL import Image

        img = Image.open(_io.BytesIO(png_bytes))
        src_w, src_h = img.size
        target_ratio = target_w / target_h
        src_ratio = src_w / src_h
        if src_ratio > target_ratio:
            # Too wide — crop sides
            new_w = int(src_h * target_ratio)
            left = (src_w - new_w) // 2
            img = img.crop((left, 0, left + new_w, src_h))
        elif src_ratio < target_ratio:
            # Too tall — crop bottom
            new_h = int(src_w / target_ratio)
            img = img.crop((0, 0, src_w, new_h))
        img = img.resize((target_w, target_h), Image.LANCZOS)
        buf = _io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        return buf.read()


class DocumentSectionRegenerateView(APIView):
    """
    POST /api/v1/documents/{id}/regenerate-section/
    Regenerate a single section of a document using the LLM.
    Body: { section_index: int, heading: str, instruction: str }
    Returns: { heading: str, content: str }
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request, doc_id: str) -> Response:
        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        heading = request.data.get("heading", "").strip()
        instruction = request.data.get("instruction", "").strip()
        if not heading:
            return Response({"error": "heading is required"}, status=400)
        if not instruction:
            instruction = f"Write a comprehensive section about: {heading}"

        # Use the LLM expander to regenerate just this one section
        prompt = (
            f"{instruction}\n\nDocument title: {doc.title}\nSection heading: {heading}"
        )
        try:
            sections = DocumentGenerateView._expand_prompt_to_sections(
                prompt=prompt,
                title=doc.title,
                doc_type=doc.doc_type,
                user=request.user,
            )
            # Take the first section's content (most relevant)
            best = next(
                (
                    s
                    for s in sections
                    if heading.lower() in s.get("heading", "").lower()
                ),
                sections[0] if sections else None,
            )
            if not best:
                return Response(
                    {"error": "Could not generate section content"}, status=500
                )
            return Response({"heading": heading, "content": best["content"]})
        except Exception as exc:
            logger.error("Section regen failed: %s", exc)
            return Response({"error": str(exc)}, status=500)


class DocumentSectionsUpdateView(APIView):
    """
    POST /api/v1/documents/{id}/update-sections/
    Save edited sections to the document metadata AND rebuild the document file.
    Body: { sections: [{heading, content}, ...] }
    Returns: updated GeneratedDocumentSerializer data with new download_url
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request, doc_id: str) -> Response:
        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        sections = request.data.get("sections", [])
        if not isinstance(sections, list) or not sections:
            return Response({"error": "sections must be a non-empty list"}, status=400)

        # Validate each section has heading + content
        for i, sec in enumerate(sections):
            if (
                not isinstance(sec, dict)
                or "heading" not in sec
                or "content" not in sec
            ):
                return Response(
                    {"error": f"Section {i} must have heading and content"}, status=400
                )

        # Delete old file if it exists
        if doc.file_path:
            old_path = Path(settings.MEDIA_ROOT) / doc.file_path
            if old_path.exists():
                try:
                    old_path.unlink()
                except Exception:
                    pass

        # Rebuild the document file with updated sections
        subtitle = doc.metadata.get("subtitle", "") if doc.metadata else ""
        author = (
            doc.metadata.get("author", "SYNAPSE AI") if doc.metadata else "SYNAPSE AI"
        )
        user_id = str(request.user.id)

        try:
            result_str, file_path_str = self._call_tool(
                doc_type=doc.doc_type,
                title=doc.title,
                sections=sections,
                subtitle=subtitle,
                author=author,
                user_id=user_id,
            )
        except Exception as exc:
            logger.error("Document rebuild failed: %s", exc)
            return Response({"error": f"Rebuild failed: {exc}"}, status=500)

        if "failed" in result_str.lower():
            return Response({"error": result_str}, status=500)

        # Update DB record
        abs_path = Path(file_path_str) if file_path_str else None
        rel_path = ""
        file_size = 0
        if abs_path and abs_path.exists():
            media_root = Path(settings.MEDIA_ROOT)
            try:
                rel_path = str(abs_path.relative_to(media_root))
            except ValueError:
                rel_path = str(abs_path)
            file_size = abs_path.stat().st_size

        # Create new version instead of overwriting the original
        next_version = (doc.version or 1) + 1
        new_doc = GeneratedDocument.objects.create(
            user=request.user,
            title=doc.title,
            doc_type=doc.doc_type,
            file_path=rel_path,
            file_size_bytes=file_size,
            agent_prompt=doc.agent_prompt,
            version=next_version,
            parent=doc.parent or doc,
            metadata={
                **doc.metadata,
                "sections": sections,
                "section_count": len(sections),
                "last_rebuilt": str(datetime.now(timezone.utc).isoformat()),
                "subtitle": (doc.metadata or {}).get("subtitle", ""),
                "author": (doc.metadata or {}).get("author", "SYNAPSE AI"),
            },
        )
        return Response(
            GeneratedDocumentSerializer(new_doc, context={"request": request}).data
        )

    @staticmethod
    def _call_tool(doc_type, title, sections, subtitle, author, user_id):
        return DocumentGenerateView._call_tool(
            doc_type=doc_type,
            title=title,
            sections=sections,
            subtitle=subtitle,
            author=author,
            user_id=user_id,
        )


class DocumentRegenerateAllView(APIView):
    """
    POST /api/v1/documents/{id}/regenerate-all/
    Use the LLM to regenerate ALL sections from the original prompt,
    rebuild the document file, update the DB record.
    Body: { instruction?: str }  (optional extra instruction appended to original prompt)
    Returns: updated GeneratedDocumentSerializer data
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request, doc_id: str) -> Response:
        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        instruction = request.data.get("instruction", "").strip()
        prompt = doc.agent_prompt or f"Write a comprehensive document about {doc.title}"
        if instruction:
            prompt = f"{prompt}\n\nAdditional instruction: {instruction}"

        # Regenerate all sections via LLM
        try:
            sections = DocumentGenerateView._expand_prompt_to_sections(
                prompt=prompt,
                title=doc.title,
                doc_type=doc.doc_type,
                user=request.user,
            )
        except Exception as exc:
            logger.error("Section expansion failed: %s", exc)
            return Response({"error": f"Section generation failed: {exc}"}, status=500)

        # Delete old file
        if doc.file_path:
            old_path = Path(settings.MEDIA_ROOT) / doc.file_path
            if old_path.exists():
                try:
                    old_path.unlink()
                except Exception:
                    pass

        # Rebuild file
        subtitle = doc.metadata.get("subtitle", "") if doc.metadata else ""
        author = (
            doc.metadata.get("author", "SYNAPSE AI") if doc.metadata else "SYNAPSE AI"
        )
        user_id = str(request.user.id)

        try:
            result_str, file_path_str = DocumentGenerateView._call_tool(
                doc_type=doc.doc_type,
                title=doc.title,
                sections=sections,
                subtitle=subtitle,
                author=author,
                user_id=user_id,
            )
        except Exception as exc:
            logger.error("Document rebuild failed: %s", exc)
            return Response({"error": f"Rebuild failed: {exc}"}, status=500)

        abs_path = Path(file_path_str) if file_path_str else None
        rel_path = ""
        file_size = 0
        if abs_path and abs_path.exists():
            media_root = Path(settings.MEDIA_ROOT)
            try:
                rel_path = str(abs_path.relative_to(media_root))
            except ValueError:
                rel_path = str(abs_path)
            file_size = abs_path.stat().st_size

        # Create new version instead of overwriting the original
        next_version = (doc.version or 1) + 1
        new_doc = GeneratedDocument.objects.create(
            user=request.user,
            title=doc.title,
            doc_type=doc.doc_type,
            file_path=rel_path,
            file_size_bytes=file_size,
            agent_prompt=doc.agent_prompt,
            version=next_version,
            parent=doc.parent or doc,
            metadata={
                **doc.metadata,
                "sections": sections,
                "section_count": len(sections),
                "last_rebuilt": str(datetime.now(timezone.utc).isoformat()),
            },
        )
        return Response(
            GeneratedDocumentSerializer(new_doc, context={"request": request}).data
        )


class DocumentVersionsView(APIView):
    """
    GET /api/v1/documents/{id}/versions/
    Returns all versions of a document (parent + all children).
    """

    permission_classes = [IsAuthenticated]

    def get(self, request: Request, doc_id: str) -> Response:
        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=request.user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        # Find the root parent
        root = doc.parent or doc
        # Get all versions: root + all children
        from django.db.models import Q

        all_versions = GeneratedDocument.objects.filter(
            Q(id=root.id) | Q(parent=root),
            user=request.user,
        ).order_by("version")

        return Response(
            GeneratedDocumentSerializer(
                all_versions, many=True, context={"request": request}
            ).data
        )


class DocumentGenerateStreamView(APIView):
    """
    POST /api/v1/documents/generate-stream/
    Server-Sent Events (SSE) streaming endpoint for document generation.
    Streams progress events as the document is being generated.

    SSE event format: data: {"step": "...", "message": "...", "progress": 0-100, "done": bool, "error"?: str}
    """

    permission_classes = [IsAuthenticated]

    def post(self, request: Request):
        import json as _json

        from django.http import StreamingHttpResponse

        # Validate input
        serializer = DocumentGenerateSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(serializer.errors, status=400)

        data = serializer.validated_data
        doc_type = data["doc_type"]
        title = data["title"]
        prompt = data["prompt"]
        subtitle = data.get("subtitle", "")
        author = data.get("author", "SYNAPSE AI")
        user = request.user
        user_id = str(user.id)
        sections = data.get("sections", [])
        model_override = data.get("model", "")

        # Resolve LLM keys HERE — before the generator starts — so the user
        # object is still fully bound to the request/DB session.
        _openrouter_key, _gemini_key = DocumentGenerateView._get_llm_keys(user)

        def event_stream():
            def emit(step, message, progress, done=False, error=None, extra=None):
                payload = {
                    "step": step,
                    "message": message,
                    "progress": progress,
                    "done": done,
                }
                if error:
                    payload["error"] = error
                if extra:
                    payload.update(extra)
                yield f"data: {_json.dumps(payload)}\n\n"

            # Use keys resolved before generator started
            openrouter_key = _openrouter_key
            gemini_key = _gemini_key

            try:
                # Step 1: Validate & start
                yield from emit(
                    "start",
                    f'Starting {doc_type.upper()} generation for "{title}"...',
                    5,
                )

                # Step 2: Report LLM key status
                openrouter_key, gemini_key = (
                    openrouter_key,
                    gemini_key,
                )  # already resolved
                has_llm = bool(openrouter_key or gemini_key)
                provider = (
                    "OpenRouter"
                    if openrouter_key
                    else ("Gemini" if gemini_key else "None")
                )

                if has_llm:
                    yield from emit("llm_check", f"AI engine ready: {provider}", 10)
                else:
                    yield from emit(
                        "llm_check",
                        "No AI key configured — using structured template",
                        10,
                    )

                # Step 3 & 4: Generate the document
                if doc_type == "html":
                    # HTML: LLM generates the actual page HTML/CSS/JS directly from the prompt
                    yield from emit(
                        "generating",
                        (
                            f"Designing your HTML page with AI ({provider})..."
                            if has_llm
                            else "Building HTML page..."
                        ),
                        20,
                    )
                    from ai_engine.agents.doc_tools import (
                        _generate_html_page_from_prompt,
                    )

                    result_str = _generate_html_page_from_prompt(
                        title=title,
                        prompt=prompt,
                        subtitle=subtitle,
                        author=author,
                        user_id=user_id,
                        openrouter_key=openrouter_key,
                        gemini_key=gemini_key,
                        model_override=model_override,
                    )
                    file_path_str = ""
                    for line in result_str.splitlines():
                        if line.startswith("Path:"):
                            file_path_str = line.replace("Path:", "").strip()
                            break
                    sections_list = []
                    yield from emit("building", "Finalising HTML page...", 75)
                else:
                    # All other types: expand prompt → sections → render
                    if not sections:
                        yield from emit(
                            "generating",
                            "Generating document structure and content with AI...",
                            20,
                        )
                        sections_list = DocumentGenerateView._expand_prompt_to_sections(
                            prompt=prompt,
                            title=title,
                            doc_type=doc_type,
                            user=user,
                            model_override=model_override,
                        )
                        yield from emit(
                            "sections_ready",
                            f'Generated {len(sections_list)} sections ({sum(len(s.get("content","").split()) for s in sections_list):,} words)',
                            50,
                        )
                    else:
                        sections_list = sections
                        yield from emit(
                            "sections_ready",
                            f"Using {len(sections_list)} provided sections",
                            50,
                        )

                    format_labels = {
                        "pdf": "Rendering premium PDF with charts and gradients",
                        "ppt": "Building PowerPoint with slides, charts and animations",
                        "word": "Composing Word document with styles and TOC",
                        "markdown": "Writing Markdown with badges and callouts",
                    }
                    yield from emit(
                        "building",
                        format_labels.get(
                            doc_type, f"Building {doc_type.upper()} document"
                        ),
                        65,
                    )

                    result_str, file_path_str = DocumentGenerateView._call_tool(
                        doc_type=doc_type,
                        title=title,
                        sections=sections_list,
                        subtitle=subtitle,
                        author=author,
                        user_id=user_id,
                    )

                if "failed" in result_str.lower():
                    yield from emit(
                        "error", result_str, 65, done=True, error=result_str
                    )
                    return

                yield from emit("file_ready", "Document file created successfully", 80)

                # Step 5: Save to database
                yield from emit("saving", "Saving to your document library...", 88)

                abs_path = Path(file_path_str) if file_path_str else None
                rel_path = ""
                file_size = 0
                if abs_path and abs_path.exists():
                    media_root = Path(settings.MEDIA_ROOT)
                    try:
                        rel_path = str(abs_path.relative_to(media_root))
                    except ValueError:
                        rel_path = str(abs_path)
                    file_size = abs_path.stat().st_size

                doc = GeneratedDocument.objects.create(
                    user=user,
                    title=title,
                    doc_type=doc_type,
                    file_path=rel_path,
                    file_size_bytes=file_size,
                    agent_prompt=prompt,
                    metadata={
                        "sections": sections_list,
                        "section_count": len(sections_list),
                        "subtitle": subtitle,
                        "author": author,
                        "provider": provider,
                    },
                )

                # Step 6: Done
                from rest_framework.request import Request as DRFRequest

                doc_data = GeneratedDocumentSerializer(
                    doc, context={"request": request}
                ).data
                yield from emit(
                    "complete",
                    f"Done! {doc_type.upper()} ready — {file_size/1024:.1f} KB, {len(sections_list)} sections",
                    100,
                    done=True,
                    extra={"document": doc_data},
                )

            except Exception as exc:
                logger.error("SSE stream generation error: %s", exc, exc_info=True)
                yield f'data: {_json.dumps({"step": "error", "message": str(exc), "progress": 0, "done": True, "error": str(exc)})}\n\n'

        response = StreamingHttpResponse(
            event_stream(),
            content_type="text/event-stream",
        )
        response["Cache-Control"] = "no-cache"
        response["X-Accel-Buffering"] = "no"
        response["Access-Control-Allow-Origin"] = "*"
        return response


class DocumentRenderView(APIView):
    """
    GET /api/v1/documents/{id}/render/?token=<jwt>
    Returns the document as a fully self-contained, scrollable HTML page
    for embedding in an iframe. Auth via ?token= query param (JWT).

    Supported:
      html     → serves the file directly (already self-contained)
      pdf      → returns an <embed> wrapper page
      markdown → parses .md → styled HTML with Synapse theme
      word     → python-docx → styled HTML
      ppt      → python-pptx → slide-by-slide HTML viewer
    """

    # Allow any — we handle auth manually so ?token= works in iframes
    permission_classes = []
    authentication_classes = []

    def get(self, request: Request, doc_id: str):
        from rest_framework_simplejwt.exceptions import TokenError
        from rest_framework_simplejwt.tokens import AccessToken

        from django.contrib.auth import get_user_model
        from django.http import HttpResponse, HttpResponseForbidden

        User = get_user_model()
        user = None

        # 1. Try ?token= query param (iframe-friendly)
        token_param = request.query_params.get("token", "").strip()
        if token_param:
            try:
                validated = AccessToken(token_param)
                user = User.objects.get(id=validated["user_id"])
            except (TokenError, User.DoesNotExist, Exception) as e:
                logger.warning(
                    "DocumentRenderView: token param invalid (%s) — trying Authorization header",
                    e,
                )
                # Don't 403 immediately — fall through to try the Bearer header

        # 2. Try Authorization header (standard DRF JWT)
        if user is None:
            auth_header = request.META.get("HTTP_AUTHORIZATION", "")
            if auth_header.startswith("Bearer "):
                try:
                    validated = AccessToken(auth_header[7:])
                    user = User.objects.get(id=validated["user_id"])
                except (TokenError, User.DoesNotExist, Exception):
                    pass

        # 3. Try refreshing via the refresh token in the cookie or header (future)
        if user is None:
            # Both methods failed — give a clear error
            logger.warning(
                "DocumentRenderView: all auth methods failed for doc %s", doc_id
            )
            return HttpResponseForbidden(
                "<html><body style='font:14px sans-serif;padding:40px;color:#dc2626;font-family:sans-serif'>"
                "<h3>⚠️ Session Expired</h3>"
                "<p>Your session has expired. Please <a href='/login' style='color:#4F46E5'>sign in again</a> "
                "and reopen the preview.</p>"
                "</body></html>",
                content_type="text/html",
            )

        try:
            doc = GeneratedDocument.objects.get(id=doc_id, user=user)
        except GeneratedDocument.DoesNotExist:
            raise Http404

        stored = Path(doc.file_path) if doc.file_path else None
        abs_path = None
        if stored:
            abs_path = (
                stored if stored.is_absolute() else Path(settings.MEDIA_ROOT) / stored
            )

        def _html_resp(html: str) -> HttpResponse:
            """Return an HTML response that is allowed to be framed by the same origin."""
            resp = HttpResponse(html, content_type="text/html; charset=utf-8")
            resp["X-Frame-Options"] = "SAMEORIGIN"
            resp["Content-Security-Policy"] = (
                "default-src * 'unsafe-inline' 'unsafe-eval' data: blob:; "
                "frame-ancestors 'self';"
            )
            return resp

        try:
            if doc.doc_type == "html" and abs_path and abs_path.exists():
                html = abs_path.read_text(encoding="utf-8")
                return _html_resp(html)

            elif doc.doc_type == "pdf" and abs_path and abs_path.exists():
                html = self._pdf_wrapper(doc, abs_path)
                return _html_resp(html)

            elif doc.doc_type == "markdown" and abs_path and abs_path.exists():
                html = self._render_markdown(doc, abs_path)
                return _html_resp(html)

            elif doc.doc_type in ("word", "docx") and abs_path and abs_path.exists():
                html = self._render_docx(doc, abs_path)
                return _html_resp(html)

            elif doc.doc_type == "ppt" and abs_path and abs_path.exists():
                html = self._render_pptx(doc, abs_path)
                return _html_resp(html)

        except Exception as exc:
            logger.error("DocumentRenderView error: %s", exc, exc_info=True)

        # Final fallback
        return _html_resp(
            f'<html><body style="font:16px sans-serif;padding:40px;color:#374151">'
            f"<p>Preview unavailable. "
            f'<a href="/api/v1/documents/{doc_id}/download/" style="color:#4F46E5">Download file</a></p>'
            f"</body></html>"
        )

    # ── PDF preview via PyMuPDF — renders all pages as images ────────────────
    @staticmethod
    def _pdf_wrapper(doc, abs_path: Path) -> str:
        """Render all PDF pages as PNG images embedded in a scrollable HTML viewer."""
        import base64

        import fitz  # PyMuPDF

        pdf = fitz.open(str(abs_path))
        pages_html = []

        # Render every page at 1.5× zoom (good quality, reasonable size)
        mat = fitz.Matrix(1.5, 1.5)
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode()
            pages_html.append(
                f'<div class="page" id="page-{page_num+1}">'
                f'<img src="data:image/png;base64,{b64}" '
                f'alt="Page {page_num+1}" loading="lazy"/>'
                f'<div class="page-label">Page {page_num+1} of {len(pdf)}</div>'
                f"</div>"
            )

        pdf.close()
        pages_body = "\n".join(pages_html)
        total = len(pages_html)

        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{
    background:#1a1a2e;
    font-family: system-ui, sans-serif;
    display: flex; flex-direction: column; min-height: 100vh;
  }}
  .topbar {{
    background: linear-gradient(90deg,#1E1B4B,#4F46E5);
    padding: 10px 20px;
    display: flex; align-items: center; gap: 12px;
    position: sticky; top: 0; z-index: 10;
    border-bottom: 1px solid rgba(255,255,255,.1);
  }}
  .topbar h1 {{
    color: #fff; font-size: 14px; font-weight: 600; flex: 1;
    white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
  }}
  .topbar .badge {{
    background: rgba(255,255,255,.15); color: #C7D2FE;
    font-size: 11px; font-weight: 700; padding: 3px 10px;
    border-radius: 20px; letter-spacing: 1px; white-space: nowrap;
  }}
  .nav {{ display: flex; align-items: center; gap: 8px; }}
  .nav button {{
    background: rgba(255,255,255,.1); border: 1px solid rgba(255,255,255,.2);
    color: #fff; padding: 4px 12px; border-radius: 6px;
    cursor: pointer; font-size: 13px; transition: background .2s;
  }}
  .nav button:hover {{ background: rgba(255,255,255,.2); }}
  .nav span {{ color: #818CF8; font-size: 12px; white-space: nowrap; }}
  .pages {{
    flex: 1; overflow-y: auto; padding: 24px 0;
    display: flex; flex-direction: column; align-items: center; gap: 16px;
  }}
  .page {{
    position: relative; max-width: 900px; width: 95%;
    border-radius: 4px; overflow: hidden;
    box-shadow: 0 8px 32px rgba(0,0,0,.5);
  }}
  .page img {{
    width: 100%; height: auto; display: block;
  }}
  .page-label {{
    position: absolute; bottom: 8px; right: 12px;
    background: rgba(0,0,0,.6); color: rgba(255,255,255,.7);
    font-size: 11px; padding: 2px 8px; border-radius: 10px;
  }}
  .progress-bar {{
    position: fixed; top: 0; left: 0; height: 3px;
    background: linear-gradient(90deg,#4F46E5,#7C3AED);
    transition: width .15s; z-index: 999;
  }}
</style>
</head>
<body>
<div id="pb" class="progress-bar" style="width:0%"></div>
<div class="topbar">
  <span class="badge">📄 PDF</span>
  <h1>{doc.title}</h1>
  <div class="nav">
    <button onclick="prevPage()">◀</button>
    <span id="pg-label">1 / {total}</span>
    <button onclick="nextPage()">▶</button>
  </div>
</div>
<div class="pages" id="pages">
{pages_body}
</div>
<script>
let cur = 1;
const total = {total};

function goTo(n) {{
  cur = Math.max(1, Math.min(total, n));
  document.getElementById('page-'+cur)?.scrollIntoView({{behavior:'smooth', block:'start'}});
  document.getElementById('pg-label').textContent = cur + ' / ' + total;
}}
function nextPage() {{ goTo(cur+1); }}
function prevPage() {{ goTo(cur-1); }}

window.addEventListener('keydown', e => {{
  if (e.key==='ArrowRight'||e.key==='ArrowDown') nextPage();
  if (e.key==='ArrowLeft'||e.key==='ArrowUp') prevPage();
}});

// Scroll progress + current page tracking
const pagesEl = document.getElementById('pages');
pagesEl.addEventListener('scroll', () => {{
  const {{scrollTop, scrollHeight, clientHeight}} = pagesEl;
  const pct = scrollTop / (scrollHeight - clientHeight) * 100;
  document.getElementById('pb').style.width = Math.min(100, pct) + '%';
  // Update current page indicator
  for (let i = 1; i <= total; i++) {{
    const el = document.getElementById('page-'+i);
    if (el) {{
      const r = el.getBoundingClientRect();
      if (r.top <= clientHeight * 0.5 && r.bottom >= 0) {{
        cur = i;
        document.getElementById('pg-label').textContent = i + ' / ' + total;
        break;
      }}
    }}
  }}
}});
</script>
</body>
</html>"""

    # ── Markdown → styled HTML ────────────────────────────────────────────────
    @staticmethod
    def _render_markdown(doc, abs_path: Path) -> str:
        import re as _re

        raw = abs_path.read_text(encoding="utf-8")

        # Strip YAML front-matter
        if raw.startswith("---"):
            end = raw.find("---", 3)
            if end != -1:
                raw = raw[end + 3 :].lstrip()

        # Simple Markdown → HTML (handles the most common patterns)
        lines = raw.split("\n")
        html_parts = []
        in_code = False
        in_table = False
        in_blockquote = False

        for line in lines:
            stripped = line.rstrip()

            # Code blocks
            if stripped.startswith("```"):
                if in_code:
                    html_parts.append("</code></pre>")
                    in_code = False
                else:
                    lang = stripped[3:].strip()
                    html_parts.append(f'<pre class="code"><code class="lang-{lang}">')
                    in_code = True
                continue
            if in_code:
                import html as _html_mod

                html_parts.append(_html_mod.escape(line) + "\n")
                continue

            # Blockquotes
            if stripped.startswith(">"):
                if not in_blockquote:
                    html_parts.append("<blockquote>")
                    in_blockquote = True
                content = stripped.lstrip("> ").strip()
                # Handle [!NOTE] / [!TIP] admonitions
                if content.startswith("[!"):
                    type_end = content.find("]")
                    admon_type = content[2:type_end] if type_end > 0 else "NOTE"
                    html_parts.append(
                        f'<p class="admon admon-{admon_type.lower()}">'
                        f"<strong>{admon_type}</strong></p>"
                    )
                else:
                    html_parts.append(f"<p>{_md_inline(content)}</p>")
                continue
            else:
                if in_blockquote:
                    html_parts.append("</blockquote>")
                    in_blockquote = False

            # Tables
            if "|" in stripped:
                if not in_table:
                    html_parts.append("<table>")
                    in_table = True
                if _re.match(r"^\|[-:\s|]+\|$", stripped):
                    continue  # separator row
                cells = [c.strip() for c in stripped.strip("|").split("|")]
                tag = "th" if not any("th" in p for p in html_parts[-3:]) else "td"
                html_parts.append(
                    "<tr>"
                    + "".join(f"<{tag}>{_md_inline(c)}</{tag}>" for c in cells)
                    + "</tr>"
                )
                continue
            else:
                if in_table:
                    html_parts.append("</table>")
                    in_table = False

            # Headings
            if stripped.startswith("# "):
                html_parts.append(f"<h1>{_md_inline(stripped[2:])}</h1>")
            elif stripped.startswith("## "):
                html_parts.append(f"<h2>{_md_inline(stripped[3:])}</h2>")
            elif stripped.startswith("### "):
                html_parts.append(f"<h3>{_md_inline(stripped[4:])}</h3>")
            elif stripped.startswith("#### "):
                html_parts.append(f"<h4>{_md_inline(stripped[5:])}</h4>")
            # Horizontal rules
            elif _re.match(r"^[-*_]{3,}$", stripped):
                html_parts.append("<hr/>")
            # Lists
            elif stripped.startswith("- ") or stripped.startswith("* "):
                html_parts.append(f"<li>{_md_inline(stripped[2:])}</li>")
            elif _re.match(r"^\d+\. ", stripped):
                item_text = _re.sub(r"^\d+\. ", "", stripped)
                html_parts.append(f"<li>{_md_inline(item_text)}</li>")
            # Images
            elif stripped.startswith("!["):
                pass  # skip image tags in text preview
            # Badges (shield.io)
            elif "img.shields.io" in stripped:
                pass
            # Empty line
            elif not stripped:
                html_parts.append("<br/>")
            else:
                html_parts.append(f"<p>{_md_inline(stripped)}</p>")

        if in_code:
            html_parts.append("</code></pre>")
        if in_table:
            html_parts.append("</table>")
        if in_blockquote:
            html_parts.append("</blockquote>")

        body = "\n".join(html_parts)
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<style>
  :root {{
    --indigo: #4F46E5; --violet: #7C3AED; --dark: #1E1B4B;
    --gray-900: #111827; --gray-700: #374151; --gray-100: #F3F4F6;
    --emerald: #10B981; --amber: #F59E0B;
  }}
  * {{ box-sizing:border-box; margin:0; padding:0; }}
  body {{
    font-family: 'Segoe UI', system-ui, sans-serif;
    background: #f8f9ff; color: var(--gray-700);
    line-height: 1.8; padding: 0;
  }}
  .topbar {{
    background: linear-gradient(90deg,var(--dark),var(--violet));
    padding: 14px 32px; position: sticky; top: 0; z-index: 10;
    display: flex; align-items: center; gap: 12px;
  }}
  .topbar .badge {{
    background: rgba(255,255,255,.15); color: #C7D2FE;
    font-size: 11px; font-weight: 700; letter-spacing: 1.5px;
    padding: 3px 10px; border-radius: 20px; text-transform: uppercase;
  }}
  .topbar h1 {{
    color: #fff; font-size: 15px; font-weight: 700; flex: 1;
    white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
  }}
  .content {{
    max-width: 820px; margin: 0 auto; padding: 40px 32px 80px;
  }}
  h1 {{ font-size: 2.2em; color: var(--dark); margin: 0.6em 0 0.3em;
        border-bottom: 3px solid var(--indigo); padding-bottom: 0.3em; }}
  h2 {{ font-size: 1.5em; color: var(--indigo); margin: 1.6em 0 0.4em;
        padding-left: 12px; border-left: 4px solid var(--indigo); }}
  h3 {{ font-size: 1.2em; color: var(--dark); margin: 1.2em 0 0.3em; }}
  h4 {{ font-size: 1em; color: var(--gray-700); margin: 1em 0 0.2em; }}
  p  {{ margin: 0.7em 0; color: var(--gray-700); }}
  a  {{ color: var(--indigo); text-decoration: none; }}
  a:hover {{ text-decoration: underline; }}
  strong {{ color: var(--gray-900); }}
  code {{ background: #EEF2FF; color: var(--violet); padding: 2px 6px;
          border-radius: 4px; font-size: 0.88em; font-family: monospace; }}
  pre.code {{ background: var(--dark); border-radius: 10px;
              padding: 20px; overflow-x: auto; margin: 1.2em 0; }}
  pre.code code {{ background: none; color: #C7D2FE;
                  font-size: 0.9em; padding: 0; }}
  blockquote {{
    border-left: 4px solid var(--indigo); background: #EEF2FF;
    padding: 14px 20px; border-radius: 0 8px 8px 0; margin: 1em 0;
  }}
  blockquote p {{ margin: 0.3em 0; color: var(--dark); }}
  .admon {{ font-weight: 700; color: var(--indigo); margin-bottom: 4px; }}
  .admon-tip {{ color: var(--emerald); }}
  .admon-warning {{ color: var(--amber); }}
  table {{ width: 100%; border-collapse: collapse; margin: 1.2em 0;
           border-radius: 8px; overflow: hidden;
           box-shadow: 0 1px 4px rgba(0,0,0,.08); }}
  th {{ background: var(--indigo); color: #fff; padding: 10px 14px;
        font-size: 13px; text-align: left; }}
  td {{ padding: 9px 14px; border-bottom: 1px solid #E5E7EB;
        font-size: 13px; }}
  tr:nth-child(even) td {{ background: #F9FAFB; }}
  li {{ margin: 0.4em 0 0.4em 1.5em; }}
  hr {{ border: none; border-top: 2px solid #E5E7EB; margin: 2em 0; }}
  br {{ display: block; margin: 0.3em; content: ""; }}
  .progress-bar {{
    position: fixed; top: 0; left: 0; height: 3px;
    background: linear-gradient(90deg,var(--indigo),var(--violet));
    transition: width .15s; z-index: 999;
  }}
</style>
</head>
<body>
<div id="pb" class="progress-bar" style="width:0%"></div>
<div class="topbar">
  <span class="badge">📋 Markdown</span>
  <h1>{doc.title}</h1>
</div>
<div class="content">
{body}
</div>
<script>
window.addEventListener('scroll',()=>{{
  const pct=window.scrollY/(document.body.scrollHeight-window.innerHeight)*100;
  document.getElementById('pb').style.width=Math.min(100,pct)+'%';
}});
</script>
</body>
</html>"""

    # ── DOCX → styled HTML ────────────────────────────────────────────────────
    @staticmethod
    def _render_docx(doc, abs_path: Path) -> str:
        import html as _html_mod

        from docx import Document as DocxDoc

        d = DocxDoc(str(abs_path))
        parts = []

        for para in d.paragraphs:
            text = para.text.strip()
            if not text:
                parts.append("<br/>")
                continue
            safe = _html_mod.escape(text)
            style = para.style.name if para.style else ""
            is_bold = any(r.bold for r in para.runs if r.text.strip())
            is_shaded = False
            try:
                pPr = para._p.find(
                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pPr"
                )
                if pPr is not None:
                    shd = pPr.find(
                        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd"
                    )
                    if shd is not None:
                        fill = shd.get(
                            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill",
                            "",
                        )
                        if fill and fill.lower() not in ("ffffff", "auto", ""):
                            is_shaded = True
            except Exception:
                pass

            if "Heading 1" in style or (is_bold and is_shaded and len(text) < 80):
                parts.append(f"<h2>{safe}</h2>")
            elif "Heading 2" in style or (is_bold and len(text) < 60):
                parts.append(f"<h3>{safe}</h3>")
            elif "Heading" in style:
                parts.append(f"<h4>{safe}</h4>")
            elif is_shaded and len(text) > 5:
                parts.append(f'<div class="badge-section">{safe}</div>')
            else:
                parts.append(f"<p>{safe}</p>")

        body = "\n".join(parts)
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<style>
  :root {{ --indigo:#4F46E5; --violet:#7C3AED; --dark:#1E1B4B;
           --gray-700:#374151; --gray-900:#111827; }}
  * {{ box-sizing:border-box; margin:0; padding:0; }}
  body {{
    font-family:'Segoe UI',system-ui,sans-serif;
    background:#fff; color:var(--gray-700); line-height:1.8;
  }}
  .topbar {{
    background:linear-gradient(90deg,#1d4ed8,#2563eb);
    padding:14px 32px; position:sticky; top:0; z-index:10;
    display:flex; align-items:center; gap:12px;
  }}
  .topbar .badge {{ background:rgba(255,255,255,.2); color:#BFDBFE;
    font-size:11px; font-weight:700; padding:3px 10px;
    border-radius:20px; letter-spacing:1.5px; }}
  .topbar h1 {{ color:#fff; font-size:15px; font-weight:700; flex:1;
    white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }}
  .page {{
    max-width: 800px; margin: 0 auto; padding: 40px 40px 80px;
    background:#fff; min-height:100vh;
    box-shadow: 0 0 40px rgba(0,0,0,.06);
  }}
  h2 {{ font-size:1.4em; color:#fff; background:linear-gradient(90deg,var(--indigo),var(--violet));
        padding:10px 16px; margin:1.5em -40px 0.8em; border-radius:0; }}
  h3 {{ font-size:1.1em; color:var(--indigo); margin:1.2em 0 0.4em;
        border-left:4px solid var(--indigo); padding-left:10px; }}
  h4 {{ font-size:1em; color:var(--dark); margin:1em 0 0.3em; }}
  p  {{ margin:0.7em 0; font-size:15px; }}
  .badge-section {{
    background:var(--indigo); color:#fff; font-weight:700;
    padding:8px 16px; margin:1.5em -40px 0.8em; font-size:13px;
    letter-spacing:0.5px;
  }}
  br {{ display:block; margin:0.3em; }}
  .progress-bar {{
    position:fixed; top:0; left:0; height:3px;
    background:linear-gradient(90deg,#1d4ed8,#7C3AED);
    transition:width .15s; z-index:999;
  }}
</style>
</head>
<body>
<div id="pb" class="progress-bar" style="width:0%"></div>
<div class="topbar">
  <span class="badge">📝 WORD</span>
  <h1>{doc.title}</h1>
</div>
<div class="page">
{body}
</div>
<script>
window.addEventListener('scroll',()=>{{
  const pct=window.scrollY/(document.body.scrollHeight-window.innerHeight)*100;
  document.getElementById('pb').style.width=Math.min(100,pct)+'%';
}});
</script>
</body>
</html>"""

    # ── PPTX → slide-by-slide HTML ────────────────────────────────────────────
    @staticmethod
    def _render_pptx(doc, abs_path: Path) -> str:
        import html as _html_mod

        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(abs_path))
        slides_html = []

        for si, slide in enumerate(prs.slides):
            shapes_html = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                safe = _html_mod.escape(text)
                try:
                    sz = shape.text_frame.paragraphs[0].runs[0].font.size
                    sz_pt = int(sz.pt) if sz else 18
                except Exception:
                    sz_pt = 18
                try:
                    bold = shape.text_frame.paragraphs[0].runs[0].font.bold
                except Exception:
                    bold = False
                try:
                    rgb = shape.text_frame.paragraphs[0].runs[0].font.color.rgb
                    color = f"#{rgb}"
                except Exception:
                    color = "#ffffff" if si == 0 else "#111827"

                if sz_pt >= 32:
                    shapes_html.append(
                        f'<div class="slide-title" style="color:{color}">{safe}</div>'
                    )
                elif sz_pt >= 20:
                    shapes_html.append(
                        f'<div class="slide-subtitle" style="color:{color}">{safe}</div>'
                    )
                else:
                    # Multi-line bullet content
                    for line in text.split("\n"):
                        if line.strip():
                            shapes_html.append(
                                f'<div class="bullet" style="color:{color}">'
                                f"▸ &nbsp;{_html_mod.escape(line.strip())}</div>"
                            )

            bg_color = "#1E1B4B" if si == 0 else ("#f8f9ff" if si % 2 == 0 else "#fff")
            try:
                bg_elem = slide.background
                fill = bg_elem.fill
                if fill.type is not None:
                    from pptx.oxml.ns import qn as _qn

                    sr = bg_elem._element.find(".//" + _qn("a:srgbClr"))
                    if sr is not None:
                        bg_color = "#" + sr.get("val", "1E1B4B")
            except Exception:
                pass

            slide_num_html = (
                f'<span class="slide-num">{si+1} / {len(prs.slides)}</span>'
            )
            slides_html.append(
                f"""
<div class="slide" style="background:{bg_color}" id="slide-{si+1}">
  <div class="slide-inner">
    {"".join(shapes_html)}
  </div>
  {slide_num_html}
</div>"""
            )

        slides_body = "\n".join(slides_html)
        total = len(prs.slides)
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<style>
  * {{ box-sizing:border-box; margin:0; padding:0; }}
  body {{
    font-family:'Segoe UI',system-ui,sans-serif;
    background:#0f0c2e; color:#fff; overflow-x:hidden;
  }}
  .topbar {{
    background:linear-gradient(90deg,#1E1B4B,#7C3AED);
    padding:12px 24px; position:sticky; top:0; z-index:10;
    display:flex; align-items:center; gap:12px;
  }}
  .topbar .badge {{ background:rgba(255,255,255,.2); color:#C7D2FE;
    font-size:11px; font-weight:700; padding:3px 10px;
    border-radius:20px; letter-spacing:1.5px; }}
  .topbar h1 {{ color:#fff; font-size:15px; font-weight:700; flex:1;
    white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }}
  .nav {{ display:flex; gap:8px; }}
  .nav button {{
    background:rgba(255,255,255,.1); border:1px solid rgba(255,255,255,.2);
    color:#fff; padding:4px 12px; border-radius:6px; cursor:pointer;
    font-size:13px; transition:background .2s;
  }}
  .nav button:hover {{ background:rgba(255,255,255,.2); }}
  .slide {{
    width:100vw; min-height:100vh;
    display:flex; align-items:center; justify-content:center;
    position:relative; padding:60px 80px;
    scroll-snap-align:start;
    border-bottom:2px solid rgba(255,255,255,.05);
  }}
  .slide-inner {{
    max-width:900px; width:100%;
  }}
  .slide-title {{
    font-size:clamp(28px,4vw,52px); font-weight:800;
    line-height:1.15; margin-bottom:20px;
  }}
  .slide-subtitle {{
    font-size:clamp(16px,2vw,22px); font-weight:500;
    margin-bottom:14px; opacity:0.85;
  }}
  .bullet {{
    font-size:clamp(14px,1.6vw,18px); padding:8px 0;
    border-bottom:1px solid rgba(255,255,255,.06);
    display:flex; align-items:flex-start; gap:8px;
  }}
  .slide-num {{
    position:absolute; bottom:20px; right:28px;
    font-size:12px; opacity:0.4; letter-spacing:1px;
  }}
  .progress-bar {{
    position:fixed; top:0; left:0; height:3px;
    background:linear-gradient(90deg,#4F46E5,#7C3AED);
    transition:width .15s; z-index:999;
  }}
  html {{ scroll-snap-type:y mandatory; scroll-behavior:smooth; }}
</style>
</head>
<body>
<div id="pb" class="progress-bar" style="width:0%"></div>
<div class="topbar">
  <span class="badge">📊 PPT</span>
  <h1>{doc.title}</h1>
  <div class="nav">
    <button onclick="prev()">◀</button>
    <button onclick="next()">▶</button>
  </div>
</div>
{slides_body}
<script>
let cur = 1;
const total = {total};
function goTo(n) {{
  cur = Math.max(1, Math.min(total, n));
  document.getElementById('slide-'+cur)?.scrollIntoView({{behavior:'smooth'}});
}}
function next() {{ goTo(cur+1); }}
function prev() {{ goTo(cur-1); }}
window.addEventListener('keydown', e => {{
  if (e.key==='ArrowRight'||e.key==='ArrowDown') next();
  if (e.key==='ArrowLeft'||e.key==='ArrowUp') prev();
}});
window.addEventListener('scroll',()=>{{
  const pct=window.scrollY/(document.body.scrollHeight-window.innerHeight)*100;
  document.getElementById('pb').style.width=Math.min(100,pct)+'%';
  // Update current slide
  for(let i=1;i<=total;i++){{
    const el=document.getElementById('slide-'+i);
    if(el){{
      const r=el.getBoundingClientRect();
      if(r.top<=window.innerHeight/2&&r.bottom>=window.innerHeight/2){{cur=i;break;}}
    }}
  }}
}});
</script>
</body>
</html>"""


def _md_inline(text: str) -> str:
    """Convert inline markdown (bold, italic, code, links) to HTML."""
    import html as _html_mod
    import re as _re

    text = _html_mod.escape(text)
    # Bold-italic
    text = _re.sub(r"\*\*\*(.+?)\*\*\*", r"<strong><em>\1</em></strong>", text)
    # Bold
    text = _re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
    # Italic
    text = _re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
    # Code
    text = _re.sub(r"`([^`]+)`", r"<code>\1</code>", text)
    # Links
    text = _re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r'<a href="\2">\1</a>', text)
    # Strikethrough
    text = _re.sub(r"~~(.+?)~~", r"<del>\1</del>", text)
    return text
